#!/usr/bin/env python3
"""gen-self-app-synthesis.py — render the c1v synthesis capstone JSON
into PPT slides + an xlsx workbook.

Reads:
  .planning/runs/self-application/synthesis/architecture_recommendation.v1.json

Writes (under --output-dir):
  synthesis.pptx                       — ~12 slides:
    cover, summary, decisions, pareto, latency, risks,
    5 mermaid PNGs (context/use_case/class/sequence/decision_network),
    residual+HoQ summary, next steps.
  synthesis.xlsx                       — 7 sheets:
    Summary, Decisions, Pareto, Risks, LatencyBudget, ResidualRisk,
    NextSteps.
  mermaid/synth-{name}.mmd             — 5 mermaid sources written from
                                         the JSON (extracted for the
                                         deck's image slides).

Designed for the master assembler's section "Synthesis —
Architecture Recommendation".

Usage:
    python3 scripts/artifact-generators/gen-self-app-synthesis.py \\
      --input .planning/runs/self-application/synthesis/architecture_recommendation.v1.json \\
      --output-dir plans/dogfooding-artifact-deck/portfolio
"""
from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path
from typing import Any, Dict, List

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _bold_header(cell, text: str) -> None:
    cell.value = text
    cell.font = Font(bold=True)
    cell.fill = PatternFill('solid', fgColor='DDDDDD')
    cell.alignment = Alignment(wrap_text=True, vertical='top')


def _normal(cell, value, *, wrap: bool = True) -> None:
    cell.value = value
    cell.alignment = Alignment(wrap_text=wrap, vertical='top')


def _set_widths(ws, widths: List[int]) -> None:
    from openpyxl.utils import get_column_letter
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w


# ---------------------------------------------------------------------------
# xlsx
# ---------------------------------------------------------------------------

def write_xlsx(d: Dict[str, Any], out_path: Path) -> None:
    wb = Workbook()
    wb.remove(wb.active)

    # Summary
    ws = wb.create_sheet('Summary')
    ws['A1'] = 'c1v — Architecture Recommendation Summary'
    ws['A1'].font = Font(bold=True, size=16)
    ws.merge_cells('A1:D1')

    rows = [
        ('Schema',          d.get('_schema')),
        ('Project',         d.get('metadata', {}).get('project_name')),
        ('Phase',           d.get('metadata', {}).get('phase_name')),
        ('Author',          d.get('metadata', {}).get('author')),
        ('Generated',       d.get('metadata', {}).get('generated_at')),
        ('Synthesized at',  d.get('synthesized_at')),
        ('Inputs hash',     d.get('inputs_hash')),
        ('Model version',   d.get('model_version')),
        ('', ''),
        ('Recommendation',  d.get('top_level_architecture', {}).get('summary')),
    ]
    for i, (k, v) in enumerate(rows, start=3):
        ws.cell(row=i, column=1, value=k).font = Font(bold=True)
        _normal(ws.cell(row=i, column=2), v)
    _set_widths(ws, [22, 110, 8, 8])

    # Decisions
    ws = wb.create_sheet('Decisions')
    headers = ['ID', 'Question', 'Winning Alternative', 'Rationale', 'Citations']
    for i, h in enumerate(headers, 1):
        _bold_header(ws.cell(row=1, column=i), h)
    for r, dec in enumerate(d.get('decisions') or [], start=2):
        _normal(ws.cell(row=r, column=1), dec.get('id'))
        _normal(ws.cell(row=r, column=2), dec.get('question'))
        _normal(ws.cell(row=r, column=3), dec.get('winning_alternative'))
        _normal(ws.cell(row=r, column=4), dec.get('rationale'))
        _normal(ws.cell(row=r, column=5), '\n'.join(dec.get('citations') or []))
    _set_widths(ws, [10, 36, 36, 60, 36])

    # Pareto
    ws = wb.create_sheet('Pareto')
    headers = ['ID', 'Name', 'Cost', 'Latency p95 (ms)', 'Availability (%)', 'Recommended', 'Dominates', 'Summary']
    for i, h in enumerate(headers, 1):
        _bold_header(ws.cell(row=1, column=i), h)
    for r, alt in enumerate(d.get('pareto_frontier') or [], start=2):
        cost = alt.get('cost', {})
        lat = alt.get('latency', {})
        av = alt.get('availability', {})
        _normal(ws.cell(row=r, column=1), alt.get('id'))
        _normal(ws.cell(row=r, column=2), alt.get('name'))
        _normal(ws.cell(row=r, column=3),
                f"{cost.get('value')} {cost.get('units', '')}".strip())
        _normal(ws.cell(row=r, column=4),
                f"{lat.get('value')} {lat.get('units', '')}".strip())
        _normal(ws.cell(row=r, column=5),
                f"{av.get('value')} {av.get('units', '')}".strip())
        _normal(ws.cell(row=r, column=6), 'YES' if alt.get('is_recommended') else '')
        _normal(ws.cell(row=r, column=7), ', '.join(alt.get('dominates') or []))
        _normal(ws.cell(row=r, column=8), alt.get('summary'))
    _set_widths(ws, [8, 44, 22, 22, 18, 14, 18, 80])

    # Risks
    ws = wb.create_sheet('Risks')
    headers = ['ID', 'Description', 'Severity', 'Mitigation']
    for i, h in enumerate(headers, 1):
        _bold_header(ws.cell(row=1, column=i), h)
    for r, risk in enumerate(d.get('risks') or [], start=2):
        _normal(ws.cell(row=r, column=1), risk.get('id'))
        _normal(ws.cell(row=r, column=2), risk.get('description'))
        _normal(ws.cell(row=r, column=3), risk.get('severity') or risk.get('rpn'))
        _normal(ws.cell(row=r, column=4), risk.get('mitigation'))
    _set_widths(ws, [10, 60, 12, 60])

    # Latency Budget
    ws = wb.create_sheet('LatencyBudget')
    chain = (d.get('tail_latency_budgets') or [{}])[0]
    ws['A1'] = f"Chain: {chain.get('chain_id', '?')}"
    ws['A1'].font = Font(bold=True, size=14)
    ws.merge_cells('A1:D1')
    headers = ['Hop', 'Interface', 'p95 (ms)', 'Notes']
    for i, h in enumerate(headers, 1):
        _bold_header(ws.cell(row=3, column=i), h)
    legs = chain.get('legs') or []
    for r, leg in enumerate(legs, start=4):
        _normal(ws.cell(row=r, column=1), leg.get('hop'))
        _normal(ws.cell(row=r, column=2), leg.get('interface_id'))
        _normal(ws.cell(row=r, column=3), leg.get('p95_ms'))
        _normal(ws.cell(row=r, column=4), leg.get('notes'))
    sum_row = 4 + len(legs)
    ws.cell(row=sum_row, column=1, value='SUM').font = Font(bold=True)
    ws.cell(row=sum_row, column=3, value=chain.get('p95_total_ms')).font = Font(bold=True)
    _set_widths(ws, [10, 30, 14, 60])

    # ResidualRisk
    ws = wb.create_sheet('ResidualRisk')
    rr = d.get('residual_risk') or {}
    ws['A1'] = (f"Threshold: RPN ≥ {rr.get('threshold', '?')}    "
                f"Flag count: {rr.get('flag_count', 0)}")
    ws['A1'].font = Font(bold=True)
    ws.merge_cells('A1:D1')
    headers = ['Failure Mode ID', 'Description', 'RPN', 'Linked Control']
    for i, h in enumerate(headers, 1):
        _bold_header(ws.cell(row=3, column=i), h)
    for r, flag in enumerate(rr.get('flags') or [], start=4):
        _normal(ws.cell(row=r, column=1), flag.get('id'))
        _normal(ws.cell(row=r, column=2), flag.get('description'))
        _normal(ws.cell(row=r, column=3), flag.get('rpn'))
        _normal(ws.cell(row=r, column=4), flag.get('control_id') or flag.get('mitigation'))
    _set_widths(ws, [16, 60, 10, 30])

    # NextSteps
    ws = wb.create_sheet('NextSteps')
    headers = ['#', 'Action', 'Owner', 'Trigger']
    for i, h in enumerate(headers, 1):
        _bold_header(ws.cell(row=1, column=i), h)
    for r, step in enumerate(d.get('next_steps') or [], start=2):
        if isinstance(step, str):
            _normal(ws.cell(row=r, column=2), step)
        else:
            _normal(ws.cell(row=r, column=1), step.get('id') or r - 1)
            _normal(ws.cell(row=r, column=2), step.get('action') or step.get('description'))
            _normal(ws.cell(row=r, column=3), step.get('owner'))
            _normal(ws.cell(row=r, column=4), step.get('trigger'))
    _set_widths(ws, [6, 80, 20, 30])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))


# ---------------------------------------------------------------------------
# mermaid extraction + rasterization
# ---------------------------------------------------------------------------

def write_mermaid_sources(d: Dict[str, Any], out_dir: Path) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    sources: List[Path] = []
    for name, src in (d.get('mermaid_diagrams') or {}).items():
        path = out_dir / f"synth-{name}.mmd"
        path.write_text(src, encoding='utf-8')
        sources.append(path)
    return sources


def rasterize(sources: List[Path], png_dir: Path) -> List[Path]:
    png_dir.mkdir(parents=True, exist_ok=True)
    pngs: List[Path] = []
    rasterizer = HERE / 'rasterize-mermaid.py'
    for src in sources:
        prefix = src.stem  # e.g. "synth-context"
        cmd = [
            sys.executable, str(rasterizer),
            '--input', str(src),
            '--output-dir', str(png_dir),
            '--prefix', prefix,
            '--theme', 'neutral',
            '--background', 'white',
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, check=False)
        if result.returncode != 0:
            sys.stderr.write(f"rasterize failed for {src}\n{result.stderr}\n")
            continue
        pngs.append(png_dir / f"{prefix}.png")
    return pngs


# ---------------------------------------------------------------------------
# pptx
# ---------------------------------------------------------------------------

def _add_title_slide(prs, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for shp in list(slide.placeholders):
        try:
            shp._element.getparent().remove(shp._element)
        except Exception:
            pass
    left = Inches(0.5)
    top = Inches(2.5)
    width = prs.slide_width - Inches(1.0)
    height = Inches(2.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(18)
        r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def _add_text_slide(prs, title: str, body: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shp in list(slide.placeholders):
        try:
            if shp.placeholder_format.idx == 0:  # keep title placeholder
                shp.text = title
                for r in shp.text_frame.paragraphs[0].runs:
                    r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
            else:
                shp._element.getparent().remove(shp._element)
        except Exception:
            pass
    left = Inches(0.5)
    top = Inches(1.4)
    width = prs.slide_width - Inches(1.0)
    height = prs.slide_height - Inches(1.8)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for line in body.split('\n'):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)


def _add_table_slide(prs, title: str, headers: List[str], rows: List[List[Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shp in list(slide.placeholders):
        try:
            if shp.placeholder_format.idx == 0:
                shp.text = title
                for r in shp.text_frame.paragraphs[0].runs:
                    r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
            else:
                shp._element.getparent().remove(shp._element)
        except Exception:
            pass
    if not rows:
        return
    left = Inches(0.3)
    top = Inches(1.4)
    width = prs.slide_width - Inches(0.6)
    height = prs.slide_height - Inches(1.8)
    table = slide.shapes.add_table(
        rows=len(rows) + 1, cols=len(headers),
        left=left, top=top, width=width, height=height,
    ).table
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = table.cell(i, j)
            c.text = '' if v is None else str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)


def _add_image_slide(prs, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shp in list(slide.placeholders):
        try:
            if shp.placeholder_format.idx == 0:
                shp.text = title
                for r in shp.text_frame.paragraphs[0].runs:
                    r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
            else:
                shp._element.getparent().remove(shp._element)
        except Exception:
            pass
    left = Inches(0.5)
    top = Inches(1.3)
    width = prs.slide_width - Inches(1.0)
    height = prs.slide_height - Inches(1.6)
    slide.shapes.add_picture(
        str(image_path), left, top, width=width, height=height,
    )


def write_pptx(d: Dict[str, Any], pngs: List[Path], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        'Synthesis — Architecture Recommendation',
        'Portfolio keystone — c1v self-application',
    )

    summary = (d.get('top_level_architecture') or {}).get('summary') or ''
    _add_text_slide(prs, 'Recommendation', summary)

    # Decisions table
    dec_rows = [
        [x.get('id'), x.get('question'), x.get('winning_alternative')]
        for x in (d.get('decisions') or [])
    ]
    _add_table_slide(prs, 'Decisions (winning alternatives)',
                     ['ID', 'Question', 'Winner'], dec_rows)

    # Pareto frontier
    par_rows = []
    for alt in d.get('pareto_frontier') or []:
        c = alt.get('cost', {})
        l = alt.get('latency', {})
        a = alt.get('availability', {})
        par_rows.append([
            ('★ ' if alt.get('is_recommended') else '') + (alt.get('id') or ''),
            alt.get('name'),
            f"{c.get('value')} {c.get('units', '')}".strip(),
            f"{l.get('value')} {l.get('units', '')}".strip(),
            f"{a.get('value')} {a.get('units', '')}".strip(),
        ])
    _add_table_slide(prs, 'Pareto Frontier',
                     ['ID', 'Architecture', 'Cost', 'Latency p95', 'Availability'],
                     par_rows)

    # Risks
    risk_rows = [
        [x.get('id'), x.get('description'), x.get('severity') or x.get('rpn')]
        for x in (d.get('risks') or [])
    ]
    _add_table_slide(prs, 'Top-Level Risks',
                     ['ID', 'Description', 'Severity / RPN'], risk_rows)

    # Latency chain
    chain = (d.get('tail_latency_budgets') or [{}])[0]
    legs = chain.get('legs') or []
    leg_rows = [
        [leg.get('hop'), leg.get('interface_id'), leg.get('p95_ms'), leg.get('notes')]
        for leg in legs
    ]
    if leg_rows:
        leg_rows.append(['SUM', '', chain.get('p95_total_ms'), ''])
    _add_table_slide(prs,
                     f"Tail Latency — chain {chain.get('chain_id', '')}",
                     ['Hop', 'Interface', 'p95 (ms)', 'Notes'], leg_rows)

    # Mermaid image slides
    title_map = {
        'synth-context': 'Context Diagram',
        'synth-use_case': 'Use Cases',
        'synth-class': 'Class Model',
        'synth-sequence': 'Sequence — Authoring Spec Emit',
        'synth-decision_network': 'Decision Network — Winners',
    }
    for png in pngs:
        title = title_map.get(png.stem, png.stem)
        _add_image_slide(prs, title, png)

    # Residual risk
    rr = d.get('residual_risk') or {}
    flag_rows = [
        [f.get('id'), f.get('description'), f.get('rpn')]
        for f in (rr.get('flags') or [])
    ]
    _add_table_slide(
        prs,
        f"Residual Risk Flags (RPN ≥ {rr.get('threshold', '?')}, n={rr.get('flag_count', 0)})",
        ['Failure Mode', 'Description', 'RPN'],
        flag_rows[:14],  # top 14
    )

    # HoQ summary
    hoq = d.get('hoq') or {}
    if hoq:
        hoq_body = (
            f"PCs:                  {hoq.get('pc_count', '?')}\n"
            f"ECs:                  {hoq.get('ec_count', '?')}\n"
            f"Matrix nonzero:       {hoq.get('matrix_nonzero', '?')} / "
            f"{hoq.get('matrix_total', '?')}\n"
            f"Sparsity:             {hoq.get('matrix_sparsity_pct', '?')}%"
        )
        _add_text_slide(prs, 'House of Quality — Summary', hoq_body)

    # Next steps
    ns_rows = []
    for step in (d.get('next_steps') or []):
        if isinstance(step, str):
            ns_rows.append(['', step, ''])
        else:
            ns_rows.append([
                step.get('id') or '',
                step.get('action') or step.get('description', ''),
                step.get('owner') or '',
            ])
    _add_table_slide(prs, 'Next Steps', ['#', 'Action', 'Owner'], ns_rows)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------

def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument('--input', required=True)
    ap.add_argument('--output-dir', required=True)
    args = ap.parse_args()

    in_path = Path(args.input).resolve()
    out_dir = Path(args.output_dir).resolve()
    d = json.loads(in_path.read_text(encoding='utf-8'))

    mmd_dir = out_dir / 'mermaid'
    sources = write_mermaid_sources(d, mmd_dir)

    png_dir = out_dir / 'mermaid-png'
    pngs = rasterize(sources, png_dir)

    pptx_path = out_dir / 'synthesis.pptx'
    write_pptx(d, pngs, pptx_path)

    xlsx_path = out_dir / 'synthesis.xlsx'
    write_xlsx(d, xlsx_path)

    result = {
        'pptx': str(pptx_path),
        'xlsx': str(xlsx_path),
        'mermaid_sources': [str(s) for s in sources],
        'mermaid_pngs': [str(p) for p in pngs],
    }
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
