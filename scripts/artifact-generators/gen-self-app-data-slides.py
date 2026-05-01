#!/usr/bin/env python3
"""gen-self-app-data-slides.py — produce per-module PPT slides for the
modules that ship JSON-only artifacts (M1, M3, M4, M5, M8).

For each module, emits a small .pptx with a section cover + 1-3 table
slides summarizing the key matrix data. The output is consumed by the
master deck assembler — see master-deck.manifest.json.

Reads (all c1v self-application JSON):
  system-design/kb-upgrade-v2/module-1-defining-scope/data_flows.v1.json
  system-design/kb-upgrade-v2/module-3-ffbd/ffbd.v1.json
  system-design/kb-upgrade-v2/module-4-decision-matrix/decision_network.v1.json
  system-design/kb-upgrade-v2/module-5-formfunction/form_function_map.v1.json
  system-design/kb-upgrade-v2/module-8-risk/fmea_early.v1.json
  system-design/kb-upgrade-v2/module-8-risk/fmea_residual.v1.json

Writes (under --output-dir):
  M1-data-flows.pptx
  M3-ffbd.pptx
  M4-decision-network.pptx
  M5-form-function.pptx
  M8-fmea.pptx

Style: B&W neutral per project convention.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# slide helpers
# ---------------------------------------------------------------------------

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _add_cover(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shp in list(slide.placeholders):
        try:
            shp._element.getparent().remove(shp._element)
        except Exception:
            pass
    tx = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        prs.slide_width - Inches(1.0), Inches(2.5),
    )
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(18)
        r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def _add_table(prs: Presentation, title: str, headers: List[str],
               rows: List[List[Any]], col_widths: List[float] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shp in list(slide.placeholders):
        try:
            if shp.placeholder_format.idx == 0:
                shp.text = title
                for r in shp.text_frame.paragraphs[0].runs:
                    r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
            else:
                shp._element.getparent().remove(shp._element)
        except Exception:
            pass
    if not rows:
        return
    left = Inches(0.3)
    top = Inches(1.4)
    width = prs.slide_width - Inches(0.6)
    height = prs.slide_height - Inches(1.8)
    table = slide.shapes.add_table(
        rows=len(rows) + 1, cols=len(headers),
        left=left, top=top, width=width, height=height,
    ).table

    if col_widths:
        total = sum(col_widths)
        unit = (prs.slide_width - Inches(0.6)) / total
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(w * unit)

    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = table.cell(i, j)
            c.text = '' if v is None else str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)


def _add_text(prs: Presentation, title: str, body: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shp in list(slide.placeholders):
        try:
            if shp.placeholder_format.idx == 0:
                shp.text = title
                for r in shp.text_frame.paragraphs[0].runs:
                    r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
            else:
                shp._element.getparent().remove(shp._element)
        except Exception:
            pass
    tx = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.4),
        prs.slide_width - Inches(1.0),
        prs.slide_height - Inches(1.8),
    )
    tf = tx.text_frame
    tf.word_wrap = True
    for line in body.split('\n'):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)


def _truncate(value: Any, n: int = 80) -> str:
    s = '' if value is None else str(value)
    return s if len(s) <= n else s[:n - 1] + '…'


# ---------------------------------------------------------------------------
# per-module renderers
# ---------------------------------------------------------------------------

def write_m1(out_dir: Path) -> Path:
    src = Path('system-design/kb-upgrade-v2/module-1-defining-scope/data_flows.v1.json')
    d = json.loads(src.read_text(encoding='utf-8'))
    prs = _new_prs()
    _add_cover(prs, 'M1 — Defining Scope', 'Data flows + system boundary')
    rows = [
        [e.get('id'), e.get('name'), _truncate(e.get('description'), 60),
         e.get('source'), e.get('sink'), e.get('criticality')]
        for e in (d.get('entries') or [])
    ]
    _add_table(prs, 'Data Flows',
               ['ID', 'Name', 'Description', 'Source', 'Sink', 'Criticality'],
               rows, col_widths=[1, 2, 5, 2, 2, 1.5])
    out_path = out_dir / 'M1-data-flows.pptx'
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def write_m3(out_dir: Path) -> Path:
    src = Path('system-design/kb-upgrade-v2/module-3-ffbd/ffbd.v1.json')
    d = json.loads(src.read_text(encoding='utf-8'))
    prs = _new_prs()
    _add_cover(prs, 'M3 — Functional Flow Block Diagram',
               '7 top-level functions / 14 arrows / 3 logic gates')
    rows = [
        [f.get('id'), _truncate(f.get('name'), 40), _truncate(f.get('description'), 60)]
        for f in (d.get('functions') or [])
    ]
    _add_table(prs, 'Functions', ['ID', 'Name', 'Description'], rows,
               col_widths=[1, 3, 6])
    arrow_rows = [
        [a.get('id'), a.get('from'), a.get('to'), _truncate(a.get('label'), 40)]
        for a in (d.get('arrows') or [])
    ]
    if arrow_rows:
        _add_table(prs, 'Arrows', ['ID', 'From', 'To', 'Label'], arrow_rows,
                   col_widths=[1, 2, 2, 5])
    gate_rows = [
        [g.get('id'), g.get('kind'), _truncate(g.get('description'), 70)]
        for g in (d.get('logic_gates') or [])
    ]
    if gate_rows:
        _add_table(prs, 'Logic Gates', ['ID', 'Kind', 'Description'], gate_rows,
                   col_widths=[1, 1.5, 7])
    out_path = out_dir / 'M3-ffbd.pptx'
    prs.save(str(out_path))
    return out_path


def write_m4(out_dir: Path) -> Path:
    src = Path('system-design/kb-upgrade-v2/module-4-decision-matrix/decision_network.v1.json')
    d = json.loads(src.read_text(encoding='utf-8'))
    prs = _new_prs()
    _add_cover(prs, 'M4 — Decision Network',
               f"Selected: {d.get('selected_architecture_id', 'AV.01')}")

    # Decision nodes summary
    node_rows = []
    for phase_key, phase in (d.get('phases') or {}).items():
        if not isinstance(phase, dict):
            continue
        for n in phase.get('decision_nodes') or []:
            node_rows.append([
                n.get('id'),
                _truncate(n.get('title'), 50),
                _truncate(n.get('question'), 70),
                len(n.get('alternatives') or []),
            ])
    _add_table(prs, 'Decision Nodes',
               ['ID', 'Title', 'Question', '# Alts'],
               node_rows, col_widths=[1, 3, 5, 0.8])

    # Score matrix (decision_id × alternative_id, raw values)
    score_rows = []
    for phase_key, phase in (d.get('phases') or {}).items():
        if not isinstance(phase, dict):
            continue
        for n in phase.get('decision_nodes') or []:
            alt_names = {a.get('id'): a.get('name') for a in n.get('alternatives') or []}
            for s in (n.get('scores') or [])[:50]:  # cap for slide readability
                score_rows.append([
                    n.get('id'),
                    s.get('alternative_id'),
                    _truncate(alt_names.get(s.get('alternative_id'), ''), 30),
                    s.get('criterion_id'),
                    s.get('raw_value'),
                    s.get('normalized_value'),
                ])
    if score_rows:
        # Cap to ~28 rows per slide for readability; split if needed.
        chunk = 28
        for i in range(0, len(score_rows), chunk):
            label = 'Scores' if i == 0 else f'Scores (cont. {i // chunk + 1})'
            _add_table(
                prs, label,
                ['Decision', 'Alt', 'Alt Name', 'Criterion', 'Raw', 'Norm'],
                score_rows[i:i + chunk],
                col_widths=[1, 0.6, 3, 1.2, 0.8, 0.8],
            )

    out_path = out_dir / 'M4-decision-network.pptx'
    prs.save(str(out_path))
    return out_path


def write_m5(out_dir: Path) -> Path:
    src = Path('system-design/kb-upgrade-v2/module-5-formfunction/form_function_map.v1.json')
    d = json.loads(src.read_text(encoding='utf-8'))
    prs = _new_prs()
    _add_cover(prs, 'M5 — Form / Function', 'Crawley concept-mapping artifacts')

    forms = (d.get('phase_1_form_inventory') or {}).get('forms') or []
    if forms:
        rows = [
            [f.get('id'), _truncate(f.get('name'), 30),
             _truncate(f.get('description', f.get('purpose', '')), 70)]
            for f in forms
        ]
        _add_table(prs, 'Phase 1 — Form Inventory',
                   ['ID', 'Name', 'Description'], rows,
                   col_widths=[1, 3, 6])

    funcs = (d.get('phase_2_function_inventory') or {}).get('functions') or []
    if funcs:
        rows = [
            [f.get('id'), _truncate(f.get('name'), 30),
             _truncate(f.get('description', f.get('purpose', '')), 70)]
            for f in funcs
        ]
        _add_table(prs, 'Phase 2 — Function Inventory',
                   ['ID', 'Name', 'Description'], rows,
                   col_widths=[1, 3, 6])

    catalog = (d.get('phase_5_operand_process_catalog') or {}).get('entries') or []
    if catalog:
        rows = [
            [e.get('operand'), _truncate(e.get('process'), 40),
             _truncate(e.get('description', ''), 60)]
            for e in catalog
        ]
        _add_table(prs, 'Phase 5 — Operand × Process Catalog',
                   ['Operand', 'Process', 'Description'], rows,
                   col_widths=[2, 4, 6])

    out_path = out_dir / 'M5-form-function.pptx'
    prs.save(str(out_path))
    return out_path


def write_m8(out_dir: Path) -> Path:
    early = json.loads(
        Path('system-design/kb-upgrade-v2/module-8-risk/fmea_early.v1.json')
        .read_text(encoding='utf-8')
    )
    residual = json.loads(
        Path('system-design/kb-upgrade-v2/module-8-risk/fmea_residual.v1.json')
        .read_text(encoding='utf-8')
    )
    prs = _new_prs()

    threshold = residual.get('high_rpn_flag_threshold', 100)
    flagged = sum(
        1 for fm in residual.get('failure_modes') or []
        if fm.get('flagged_high_rpn')
    )
    _add_cover(
        prs, 'M8 — Risk / FMEA',
        f"{len(early.get('failure_modes') or [])} early modes  /  "
        f"{len(residual.get('failure_modes') or [])} residual  /  "
        f"{flagged} flagged ≥ RPN {threshold}",
    )

    early_rows = [
        [fm.get('id'), _truncate(fm.get('failure_mode'), 40),
         _truncate(fm.get('potential_cause'), 40),
         fm.get('severity'), fm.get('likelihood'),
         fm.get('detectability'), fm.get('rpn')]
        for fm in (early.get('failure_modes') or [])
    ]
    _add_table(prs, 'FMEA Early',
               ['ID', 'Failure Mode', 'Cause', 'S', 'L', 'D', 'RPN'],
               early_rows, col_widths=[1, 4, 4, 0.5, 0.5, 0.5, 0.7])

    res_rows = [
        [
            fm.get('id'),
            _truncate(fm.get('failure_mode'), 38),
            fm.get('severity'), fm.get('likelihood'),
            fm.get('detectability'),
            fm.get('rpn'), fm.get('weighted_rpn'),
            '★' if fm.get('flagged_high_rpn') else '',
            _truncate(', '.join(
                c if isinstance(c, str) else (c.get('id') or c.get('name') or str(c))
                for c in (fm.get('landed_controls') or [])
            ), 40),
        ]
        for fm in (residual.get('failure_modes') or [])
    ]
    _add_table(prs, f'FMEA Residual (threshold ≥ {threshold})',
               ['ID', 'Failure Mode', 'S', 'L', 'D', 'RPN', 'wRPN', 'Flag', 'Controls'],
               res_rows,
               col_widths=[1, 4, 0.5, 0.5, 0.5, 0.6, 0.7, 0.5, 4])

    out_path = out_dir / 'M8-fmea.pptx'
    prs.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------

def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument('--output-dir', required=True)
    args = ap.parse_args()
    out_dir = Path(args.output_dir).resolve()

    result = {
        'm1': str(write_m1(out_dir)),
        'm3': str(write_m3(out_dir)),
        'm4': str(write_m4(out_dir)),
        'm5': str(write_m5(out_dir)),
        'm8': str(write_m8(out_dir)),
    }
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
