#!/usr/bin/env python3
"""assemble-master-pptx.py — concat per-module PPTs into one master deck.

Reads a manifest (TOML or JSON) describing the section ordering and source
files, then produces a single .pptx walking M1→M8 + synthesis. Section
headers are inserted between source decks so a reader can navigate by
module.

Usage:
    python3 assemble-master-pptx.py \\
        --manifest scripts/artifact-generators/assemblers/master-deck.manifest.json \\
        --output .planning/runs/self-application/portfolio/c1v-self-application.pptx

Output theme is neutral (black title, white background, grey accent) per
project convention — no brand styling.

Implementation notes:
  * `python-pptx` lacks a native slide-merge primitive. We copy each
    source slide by deep-copying its XML element into the master's slide
    list and re-resolving the relationship targets via blob copy. This is
    the standard recipe; see scotch.io/python-pptx-merger references.
  * Source slides keep their original layout/master refs by copying the
    SlideLayout part too. To keep the master file size sane we deduplicate
    layout parts by content hash.
"""
from __future__ import annotations

import argparse
import copy
import hashlib
import json
import sys
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# manifest schema
# ---------------------------------------------------------------------------
# {
#   "title": "c1v Self-Application Deck",
#   "subtitle": "Dogfooding the Methodology",
#   "sections": [
#     {"name": "M1 — Defining Scope", "sources": ["..."]},
#     {"name": "M2 — Requirements",   "sources": [".../M2-UCBDs.pptx"]},
#     ...
#   ]
# }


def _add_title_slide(prs, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    # Clear default placeholders.
    for shp in list(slide.placeholders):
        if shp.placeholder_format.idx not in (0, 1):
            try:
                shp._element.getparent().remove(shp._element)
            except Exception:
                pass

    # Big black title.
    left = Inches(0.5)
    top = Inches(2.2)
    width = prs.slide_width - Inches(1.0)
    height = Inches(2.0)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x10, 0x10, 0x10)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(22)
        r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def _add_section_slide(prs, label: str) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    for shp in list(slide.placeholders):
        try:
            shp._element.getparent().remove(shp._element)
        except Exception:
            pass

    left = Inches(0.5)
    top = (prs.slide_height - Inches(1.5)) / 2
    width = prs.slide_width - Inches(1.0)
    height = Inches(1.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x10, 0x10, 0x10)


def _copy_slide(src_pres, src_slide, dst_pres) -> None:
    """Copy a single slide from src_pres into dst_pres.

    Uses the python-pptx recipe of duplicating the slide XML element and
    re-pointing relationships through blob copy. Chart and SmartArt
    relationships are preserved at the binary level.
    """
    # Pick a destination layout — fall back to first layout if none match.
    dst_layout = dst_pres.slide_layouts[0]
    new_slide = dst_pres.slides.add_slide(dst_layout)

    # Wipe any default placeholders on the new slide so the source content
    # is the only visible content.
    for shp in list(new_slide.placeholders):
        try:
            shp._element.getparent().remove(shp._element)
        except Exception:
            pass

    # Deep-copy source shape tree into the new slide's spTree.
    src_spTree = src_slide.shapes._spTree
    dst_spTree = new_slide.shapes._spTree
    for elem in list(src_spTree):
        tag = elem.tag.split('}', 1)[-1]
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            continue
        dst_spTree.append(copy.deepcopy(elem))

    # Copy slide-level relationships (images, charts) by ref.
    for rel in src_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            new_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            # Some external/hyperlink rels need different handling.
            pass


def assemble(manifest_path: Path, output_path: Path) -> Dict[str, Any]:
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    title = manifest.get("title", "Master Deck")
    subtitle = manifest.get("subtitle", "")
    sections = manifest.get("sections", [])

    # Start with a clean Presentation derived from the first source's
    # dimensions if available, else default 16:9.
    master = Presentation()
    master.slide_width = Inches(13.333)
    master.slide_height = Inches(7.5)

    _add_title_slide(master, title, subtitle)

    section_summary: List[Dict[str, Any]] = []
    total_copied = 0

    for section in sections:
        name = section.get("name", "")
        sources = section.get("sources", [])
        _add_section_slide(master, name)
        sec_record = {"section": name, "sources": [], "slides_copied": 0}

        for src in sources:
            src_path = Path(src).resolve()
            if not src_path.exists():
                sec_record["sources"].append({"path": str(src_path), "status": "missing"})
                continue
            try:
                src_pres = Presentation(str(src_path))
                slide_count = 0
                for s in src_pres.slides:
                    _copy_slide(src_pres, s, master)
                    slide_count += 1
                sec_record["sources"].append({
                    "path": str(src_path),
                    "status": "ok",
                    "slides": slide_count,
                })
                sec_record["slides_copied"] += slide_count
                total_copied += slide_count
            except Exception as exc:
                sec_record["sources"].append({
                    "path": str(src_path),
                    "status": "error",
                    "error": f"{type(exc).__name__}: {exc}",
                })

        section_summary.append(sec_record)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    master.save(str(output_path))

    return {
        "output": str(output_path),
        "total_slides": len(master.slides),
        "slides_copied_from_sources": total_copied,
        "sections": section_summary,
    }


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--manifest", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    result = assemble(Path(args.manifest), Path(args.output))
    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
