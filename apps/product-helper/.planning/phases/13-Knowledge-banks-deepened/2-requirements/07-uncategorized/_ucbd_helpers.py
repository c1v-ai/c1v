"""Shared helpers for UCBD slide generation.

Layout: 13.333 x 7.5 inch slides, 0.1 inch margins.
Font: Arial. Title 14pt, metadata 10pt, table header 10pt, table body 9pt.
Tables: black 0.75pt borders, light gray header row.
Combined slides: flow + requirements on one slide when they fit.
Split slides: flow on slide A, requirements on slide B (for large UCs).
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Constants ───────────────────────────────────────────────────

FONT = "Arial"
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
HEADER_BG = RGBColor(230, 230, 230)

SW = 13.333
SH = 7.5
MARGIN = 0.1
USABLE_W = SW - 2 * MARGIN

TITLE_SIZE = Pt(14)
META_SIZE = Pt(10)
TABLE_HEADER_SIZE = Pt(10)
TABLE_BODY_SIZE = Pt(9)
STEP_NUM_SIZE = Pt(9)

BLANK_LAYOUT = 6
BORDER_WIDTH = 9525  # 0.75pt in EMU


# ── Cell helpers ────────────────────────────────────────────────

def _set_font(cell, text, size=TABLE_BODY_SIZE, bold=False, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.color.rgb = BLACK
    run.font.bold = bold
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_border(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = tcPr.makeelement(ns + edge, {})
        ln.set("w", str(BORDER_WIDTH))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        sf = ln.makeelement(ns + "solidFill", {})
        clr = sf.makeelement(ns + "srgbClr", {})
        clr.set("val", "000000")
        sf.append(clr)
        ln.append(sf)
        tcPr.append(ln)


def _fill(cell, color=WHITE):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _header_cell(cell, text, alignment=PP_ALIGN.LEFT):
    _set_font(cell, text, TABLE_HEADER_SIZE, bold=True, alignment=alignment)
    _fill(cell, HEADER_BG)
    _add_border(cell)


def _body_cell(cell, text, bold=False, alignment=PP_ALIGN.LEFT):
    _set_font(cell, text, TABLE_BODY_SIZE, bold=bold, alignment=alignment)
    _fill(cell, WHITE)
    _add_border(cell)


# ── Slide components ────────────────────────────────────────────

def add_title(slide, text, top=None):
    if top is None:
        top = MARGIN
    tb = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(top), Inches(USABLE_W), Inches(0.35)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = TITLE_SIZE
    run.font.color.rgb = BLACK
    run.font.bold = True


def add_metadata(slide, metadata, top=0.45):
    """Compact 2-line metadata strip below title."""
    tb = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(top), Inches(USABLE_W), Inches(0.70)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    # Line 1: first 3 fields (Use Case, Actor, Trigger)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(2)
    for i, (label, value) in enumerate(metadata[:3]):
        if i > 0:
            sep = p.add_run()
            sep.text = "   |   "
            sep.font.name = FONT
            sep.font.size = META_SIZE
            sep.font.color.rgb = BLACK
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.name = FONT
        rl.font.size = META_SIZE
        rl.font.color.rgb = BLACK
        rl.font.bold = True
        rv = p.add_run()
        rv.text = value
        rv.font.name = FONT
        rv.font.size = META_SIZE
        rv.font.color.rgb = BLACK

    # Line 2: remaining fields (Start/End Condition)
    if len(metadata) > 3:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        for i, (label, value) in enumerate(metadata[3:]):
            if i > 0:
                sep = p2.add_run()
                sep.text = "   |   "
                sep.font.name = FONT
                sep.font.size = META_SIZE
                sep.font.color.rgb = BLACK
            rl = p2.add_run()
            rl.text = f"{label}: "
            rl.font.name = FONT
            rl.font.size = META_SIZE
            rl.font.color.rgb = BLACK
            rl.font.bold = True
            rv = p2.add_run()
            rv.text = value
            rv.font.name = FONT
            rv.font.size = META_SIZE
            rv.font.color.rgb = BLACK


def add_step_table(slide, steps, top=1.2):
    """Full-width step table with borders."""
    rows = len(steps) + 1
    tbl_shape = slide.shapes.add_table(
        rows, 2,
        Inches(MARGIN), Inches(top),
        Inches(USABLE_W), Inches(0.01),
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(0.45)
    tbl.columns[1].width = Inches(USABLE_W - 0.45)

    _header_cell(tbl.cell(0, 0), "#", PP_ALIGN.CENTER)
    _header_cell(tbl.cell(0, 1), "Step")

    for i, (num, desc) in enumerate(steps, start=1):
        _body_cell(tbl.cell(i, 0), num, bold=True, alignment=PP_ALIGN.CENTER)
        _body_cell(tbl.cell(i, 1), desc)


def add_req_table(slide, requirements, top=1.2):
    """Full-width requirements table with borders."""
    rows = len(requirements) + 1
    tbl_shape = slide.shapes.add_table(
        rows, 3,
        Inches(MARGIN), Inches(top),
        Inches(USABLE_W), Inches(0.01),
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.0)
    tbl.columns[1].width = Inches(USABLE_W - 1.0 - 1.8)
    tbl.columns[2].width = Inches(1.8)

    _header_cell(tbl.cell(0, 0), "REQ-ID", PP_ALIGN.CENTER)
    _header_cell(tbl.cell(0, 1), "Requirement")
    _header_cell(tbl.cell(0, 2), "Constant", PP_ALIGN.CENTER)

    for i, (rid, req, const) in enumerate(requirements, start=1):
        _body_cell(tbl.cell(i, 0), rid, bold=True, alignment=PP_ALIGN.CENTER)
        _body_cell(tbl.cell(i, 1), req)
        _body_cell(tbl.cell(i, 2), const, alignment=PP_ALIGN.CENTER)


# ── Layout functions ────────────────────────────────────────────

STEP_ROW_H = 0.24
REQ_ROW_H = 0.22
HEADER_H = 0.28
META_H = 0.70
TITLE_H = 0.35
GAP = 0.15


def can_combine(n_steps, n_reqs):
    """Check if flow + requirements fit on one slide."""
    total = (MARGIN + TITLE_H + META_H
             + HEADER_H + n_steps * STEP_ROW_H
             + GAP
             + HEADER_H + n_reqs * REQ_ROW_H
             + MARGIN)
    return total <= SH


def add_combined_slide(prs, title, metadata, steps, requirements):
    """Single slide: title + metadata + step table + requirements table."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    add_title(slide, title)
    add_metadata(slide, metadata)

    step_top = MARGIN + TITLE_H + META_H
    add_step_table(slide, steps, top=step_top)

    req_top = step_top + HEADER_H + len(steps) * STEP_ROW_H + GAP
    add_req_table(slide, requirements, top=req_top)

    return slide


def add_split_slides(prs, flow_title, req_title, metadata, steps, requirements):
    """Two slides: flow (A) + requirements (B)."""
    # Slide A: flow
    slide_a = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    add_title(slide_a, flow_title)
    add_metadata(slide_a, metadata)
    add_step_table(slide_a, steps)

    # Slide B: requirements
    slide_b = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    add_title(slide_b, req_title)
    add_req_table(slide_b, requirements, top=0.50)

    return slide_a, slide_b
