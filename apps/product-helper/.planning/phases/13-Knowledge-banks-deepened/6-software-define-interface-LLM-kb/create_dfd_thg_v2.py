#!/usr/bin/env python3
"""Create THG Data Flow Diagram v2 — Option B (7 subsystems).

CESYS526 DFD with 7 subsystems + patent sub-module labels.
25 interfaces (21 internal + 4 external) on a 17 x 10" slide for readability.

Layout strategy:
  - Top row (L→R): SS4 → SS1 → SS2 → SS3  (primary operational flow)
  - Middle right: SS7 (Team Coordination, below SS3)
  - Bottom row: SS5 (center), SS6 (right of SS5)
  - External elements: edges and below
  - Labels: short (IF-ID + key words), white background, 9pt min
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# === SLIDE ===
SW, SH = 23.0, 13.0  # Extra-wide for maximum readability

# === COLORS ===
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
SS_FILL = RGBColor(218, 232, 252)
CORE_FILL = RGBColor(255, 235, 235)
PLATFORM_FILL = RGBColor(232, 245, 233)
TEAM_FILL = RGBColor(255, 243, 224)
EXT_FILL = RGBColor(240, 240, 240)
NEW_COLOR = RGBColor(139, 69, 19)
LABEL_BG = RGBColor(255, 255, 255)
GRAY = RGBColor(80, 80, 80)

# === SIZES ===
SS_W, SS_H = 2.8, 1.3
EXT_W, EXT_H = 2.0, 0.85
SS_FONT = 14
SUB_FONT = 10
EXT_FONT = 12
LABEL_FONT_SZ = 11
TITLE_FONT_SZ = 22


def _anchor_mid(shape):
    bp = shape.text_frame._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')


def add_ss(slide, cx, cy, name, fill, sub=None):
    h = SS_H + (0.3 if sub else 0)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - SS_W/2), Inches(cy - h/2), Inches(SS_W), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = BLACK; s.line.width = Pt(2.5)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(SS_FONT); p.font.bold = True; p.font.color.rgb = BLACK
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub; p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(SUB_FONT); p2.font.italic = True
        p2.font.color.rgb = RGBColor(120, 70, 20)
    _anchor_mid(s)


def add_ext(slide, cx, cy, name):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx - EXT_W/2), Inches(cy - EXT_H/2), Inches(EXT_W), Inches(EXT_H))
    s.fill.solid(); s.fill.fore_color.rgb = EXT_FILL
    s.line.color.rgb = BLACK; s.line.width = Pt(1.5)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(EXT_FONT); p.font.color.rgb = BLACK
    _anchor_mid(s)


def arrow(slide, x1, y1, x2, y2, label="", off=(0, 0),
          color=BLACK, width=Pt(1.5), lbl_color=None, bold=False):
    c = slide.shapes.add_connector(1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = width
    ln = c.line._ln
    t = ln.makeelement(qn('a:tailEnd'), {})
    t.set('type', 'triangle'); t.set('w', 'med'); t.set('len', 'med')
    ln.append(t)
    if label:
        mx = (x1 + x2)/2 + off[0]
        my = (y1 + y2)/2 + off[1]
        lw = max(1.4, len(label) * 0.07)
        lh = 0.38
        tb = slide.shapes.add_textbox(
            Inches(mx - lw/2), Inches(my - lh/2), Inches(lw), Inches(lh))
        tb.fill.background()
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(LABEL_FONT_SZ)
        p.font.color.rgb = lbl_color or GRAY
        p.font.bold = bold


def build_dfd():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── System boundary ──
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(0.8), Inches(18.5), Inches(10.0))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(250, 250, 255)
    b.line.color.rgb = RGBColor(100, 100, 140); b.line.width = Pt(2.0)
    ln = b.line._ln
    d = ln.makeelement(qn('a:prstDash'), {}); d.set('val', 'dash'); ln.append(d)
    p = b.text_frame.paragraphs[0]
    p.text = "SYSTEM BOUNDARY — Team Heat Guard (7 Subsystems)"
    p.font.size = Pt(11); p.font.color.rgb = RGBColor(100, 100, 140); p.font.italic = True

    # ── Title ──
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(20), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "THG Data Flow Diagram — Operational Interfaces (Option B: 7 Subsystems)"
    p.font.size = Pt(TITLE_FONT_SZ); p.font.bold = True; p.font.color.rgb = BLACK

    # ═══════════════════════════════════════════════════════════
    # POSITIONS — wide spacing for 23" x 13" canvas
    # ═══════════════════════════════════════════════════════════

    # Top row: primary operational flow (~5" between centers)
    S4 = (5.5, 3.0)     # Mobile App
    S1 = (10.0, 3.0)    # Prediction Engine
    S2 = (14.5, 3.0)    # Alert Engine
    S3 = (19.5, 3.0)    # Dashboard

    # Middle right: Team Coordination (below Dashboard)
    S7 = (19.5, 7.0)    # Team Coordination

    # Bottom: supporting subsystems (well below, ~4.5" gap)
    S5 = (10.0, 9.5)    # Platform Core
    S6 = (14.5, 9.5)    # Regulatory Engine

    # External elements — further out for breathing room
    WORKER = (1.5, 3.0)
    WEATHER = (5.5, 11.5)
    SSO = (7.0, 12.0)
    SAFETY = (19.5, 1.0)
    SUPER = (14.5, 12.0)
    ADMIN = (10.0, 12.0)
    AUDIT = (12.0, 11.5)
    EMERG = (18.0, 11.5)

    # ═══════════════════════════════════════════════════════════
    # DRAW ELEMENTS
    # ═══════════════════════════════════════════════════════════

    # External elements (draw first, behind arrows)
    add_ext(slide, *WORKER, "Worker")
    add_ext(slide, *WEATHER, "Weather API")
    add_ext(slide, *SSO, "SSO Provider")
    add_ext(slide, *SAFETY, "Safety Manager")
    add_ext(slide, *SUPER, "Site Supervisor")
    add_ext(slide, *ADMIN, "Employer Admin")
    add_ext(slide, *AUDIT, "Audit Systems")
    add_ext(slide, *EMERG, "Emergency Svc")

    # Subsystems
    add_ss(slide, *S4, "SS4\nMobile App", SS_FILL)
    add_ss(slide, *S1, "SS1\nPrediction Engine", CORE_FILL,
           "incl. F.4a Stratification, ML, Acclim")
    add_ss(slide, *S2, "SS2\nAlert Engine", CORE_FILL,
           "incl. F.5.4a Intervention Engine")
    add_ss(slide, *S3, "SS3\nDashboard", SS_FILL)
    add_ss(slide, *S7, "SS7\nTeam Coordination", TEAM_FILL,
           "PATENT: sessions, rotations")
    add_ss(slide, *S5, "SS5\nPlatform Core", PLATFORM_FILL)
    add_ss(slide, *S6, "SS6\nRegulatory Engine", PLATFORM_FILL)

    # ═══════════════════════════════════════════════════════════
    # INTERNAL INTERFACES — existing (black arrows)
    # ═══════════════════════════════════════════════════════════

    # IF-01: SS4 → SS1
    arrow(slide, 6.9, 2.6, 8.6, 2.6,
          "IF-01: activity, clothing, profile", off=(0, -0.4))

    # IF-02: SS1 → SS2
    arrow(slide, 11.4, 2.6, 13.1, 2.6,
          "IF-02: Trec, Dlim, risk, trigger", off=(0, -0.4))

    # IF-03: SS2 → SS4 (alert + prescription — runs below top row)
    arrow(slide, 13.1, 3.8, 6.9, 3.8,
          "IF-03: alert + intervention Rx (fluid, cooling, rest, equip)", off=(0, 0.35))

    # IF-04: SS2 → SS3
    arrow(slide, 15.9, 2.6, 18.1, 2.6,
          "IF-04: alert status, outcomes", off=(0, -0.4))

    # IF-05: SS4 → SS2 (ack — runs just below IF-01/02 line)
    arrow(slide, 6.9, 3.4, 13.1, 3.4,
          "IF-05: ack (worker, GPS, pre_trec)", off=(0, 0.32))

    # IF-06: SS1 → SS5 (prediction audit — straight down)
    arrow(slide, 9.5, 3.95, 9.5, 8.75,
          "IF-06: prediction audit\n(wbgt, ml_id)", off=(-1.3, 0))

    # IF-07: SS2 → SS5 (alert audit — diagonal)
    arrow(slide, 13.8, 3.95, 10.8, 8.75,
          "IF-07: alert lifecycle audit", off=(0.4, -0.4))

    # IF-08: SS6 → SS2 (thresholds — straight up)
    arrow(slide, 14.5, 8.75, 14.5, 3.95,
          "IF-08: thresholds (risk,\nWBGT per metabolic/regimen)", off=(1.5, 0))

    # IF-09: SS6 → SS3 (report templates — diagonal up-right)
    arrow(slide, 15.8, 9.0, 19.0, 3.95,
          "IF-09: report templates\n(jurisdiction)", off=(0.8, 0))

    # IF-10: SS5 → SS4 (auth — diagonal up-left)
    arrow(slide, 8.8, 8.8, 5.9, 3.95,
          "IF-10: JWT, RBAC, tenant", off=(-1.2, 0))

    # IF-11: SS5 → SS3 (auth — diagonal up-right)
    arrow(slide, 11.4, 8.8, 19.0, 3.95,
          "IF-11: JWT, RBAC, tenant", off=(1.5, 0.4))

    # IF-12: SS5 → SS1 (tenant config — straight up)
    arrow(slide, 10.5, 8.75, 10.5, 3.95,
          "IF-12: tenant config,\nworker profiles, acclim", off=(1.4, 0))

    # IF-13: SS3 → SS6 (report request — diagonal down-left)
    arrow(slide, 18.5, 3.95, 15.5, 8.75,
          "IF-13: report request", off=(-1.0, 0))

    # ═══════════════════════════════════════════════════════════
    # NEW INTERFACES — brown, bold, thicker
    # ═══════════════════════════════════════════════════════════
    NW = Pt(2.5)

    # IF-18: SS1 → SS7 (member state — diagonal to SS7)
    arrow(slide, 11.4, 3.5, 18.1, 6.5,
          "IF-18: member state (risk, dlim, recovery)",
          off=(0.6, -0.4), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # IF-19: SS7 → SS3 (session state — straight up)
    arrow(slide, 19.1, 6.2, 19.1, 3.95,
          "IF-19: sessions,\nrotation recs",
          off=(1.0, 0), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # IF-20: SS3 → SS7 (manager commands — straight down)
    arrow(slide, 19.9, 3.95, 19.9, 6.2,
          "IF-20: create/end,\naccept/reject",
          off=(-1.0, 0), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # IF-21: SS7 → SS4 (rotation notify — long diagonal)
    arrow(slide, 18.1, 7.4, 6.9, 3.6,
          "IF-21: rotation notification to worker",
          off=(0, 0.4), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # IF-22: SS7 → SS5 (team audit — diagonal down-left)
    arrow(slide, 18.1, 7.6, 11.4, 9.2,
          "IF-22: team coord audit",
          off=(0, -0.4), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # IF-23: SS6 → SS1 (WBGT thresholds — diagonal up-left)
    arrow(slide, 13.3, 8.8, 10.8, 3.95,
          "IF-23: WBGT thresholds (AL/TLV)",
          off=(-1.3, 0), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # IF-24: SS5 → SS7 (auth — diagonal right)
    arrow(slide, 11.4, 9.0, 18.1, 7.5,
          "IF-24: auth context (JWT, RBAC)",
          off=(0, 0.4), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # IF-25: SS2 → SS1 (intervention ML feedback — reverse of IF-02)
    arrow(slide, 13.1, 3.2, 11.4, 3.2,
          "IF-25: intervention outcomes -> ML",
          off=(0, 0.32), color=NEW_COLOR, width=NW, lbl_color=NEW_COLOR, bold=True)

    # ═══════════════════════════════════════════════════════════
    # EXTERNAL INTERFACES
    # ═══════════════════════════════════════════════════════════

    # IF-14: Weather API → SS1
    arrow(slide, 6.3, 11.0, 9.3, 3.95,
          "IF-14: temp, humidity, wind,\nsolar, dew_point, hPa", off=(-1.4, 0))

    # IF-15: SSO → SS5
    arrow(slide, 7.8, 11.6, 9.2, 10.1,
          "IF-15: SAML assertion", off=(0, 0.3))

    # IF-16: SS2 → Emergency
    arrow(slide, 15.2, 3.95, 17.5, 11.1,
          "IF-16: critical alert", off=(0.7, 0))

    # IF-17: SS5 → Audit
    arrow(slide, 11.0, 10.2, 11.5, 11.1,
          "IF-17: audit logs", off=(0.6, 0))

    # Actor interfaces
    arrow(slide, 2.5, 2.6, 4.1, 2.6,
          "biometrics, activity", off=(0, -0.35))
    arrow(slide, 4.1, 3.5, 2.5, 3.5,
          "risk alerts, predictions", off=(0, 0.3))
    arrow(slide, 19.0, 1.45, 19.0, 2.15,
          "config", off=(0.6, 0))
    arrow(slide, 20.0, 2.15, 20.0, 1.45,
          "risk grid, reports", off=(-1.0, 0))
    arrow(slide, 14.5, 3.95, 14.5, 11.55,
          "escalation alerts", off=(0.8, 0))
    arrow(slide, 10.0, 11.55, 10.0, 10.2,
          "org config, provisioning", off=(1.0, 0))

    # ── Legend ──
    lx, ly = 0.3, 10.0
    items = [
        (MSO_SHAPE.ROUNDED_RECTANGLE, CORE_FILL, "Engine (SS1, SS2)"),
        (MSO_SHAPE.ROUNDED_RECTANGLE, SS_FILL, "UI (SS3, SS4)"),
        (MSO_SHAPE.ROUNDED_RECTANGLE, PLATFORM_FILL, "Platform (SS5, SS6)"),
        (MSO_SHAPE.ROUNDED_RECTANGLE, TEAM_FILL, "Team Coord (SS7) — PATENT"),
        (MSO_SHAPE.OVAL, EXT_FILL, "External Element"),
    ]
    for i, (st, fill, lbl) in enumerate(items):
        s = slide.shapes.add_shape(st,
            Inches(lx), Inches(ly + i * 0.35), Inches(0.5), Inches(0.25))
        s.fill.solid(); s.fill.fore_color.rgb = fill
        s.line.color.rgb = BLACK; s.line.width = Pt(1)
        tb = slide.shapes.add_textbox(
            Inches(lx + 0.6), Inches(ly + i * 0.35), Inches(2.5), Inches(0.25))
        p = tb.text_frame.paragraphs[0]
        p.text = lbl; p.font.size = Pt(11); p.font.color.rgb = BLACK

    # ── New arrow legend ──
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(ly + 5 * 0.35), Inches(3.0), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Brown arrows = NEW (IF-18 to IF-25)"
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = NEW_COLOR

    # ── Stats ──
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.55), Inches(3.0), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "25 interfaces: 21 internal + 4 external\n8 new | 10 modified | 7 unchanged"
    p.font.size = Pt(10); p.font.color.rgb = RGBColor(100, 100, 100)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "THG_DFD_v2.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"  Slide: {SW}\" x {SH}\" (wide format for readability)")
    print(f"  7 subsystems + 8 external elements")
    print(f"  25 interface arrows (brown = new, black = existing)")


if __name__ == "__main__":
    build_dfd()
