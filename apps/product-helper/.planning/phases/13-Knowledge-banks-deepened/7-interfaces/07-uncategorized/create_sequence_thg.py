#!/usr/bin/env python3
"""Create THG Sequence Diagrams – Cornell CESYS526 standard (Option B: 7 subsystems).

Slide 1: UC1 — Worker Requests Heat Risk Prediction (+ WBGT screening, IF-23)
Slide 2: UC2/UC3 — Alert Delivery, Acknowledgment & Escalation (+ intervention, IF-25)
Slide 3: UC7 — Onboard New Worker (+ acclimatization init)
Slide 4: UC9 — Provision Organization (+ bulk import, jurisdiction)
Slide 5: UC6 — Generate Compliance Report (+ multi-jurisdiction, verification)
Slide 6: NEW — Team Coordination (PATENT) — IF-18 to IF-22, IF-24
Slide 7: NEW — Pre-Event Risk Stratification (PATENT)
Slide 8: NEW — Activity Switch Mid-Shift (UC13)

Uses the CESYS526 sequence diagram template (16×12" slide).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# === FILE PATHS ===
TEMPLATE = "/Users/davidancor/Documents/MBA/System Design - eCornell/6 - Defining Interfaces/Module 1/CESYS526_sequence-diagram.pptx"
OUTPUT = "/Users/davidancor/Documents/MBA/System Design - eCornell/thg-system-design-example/diagrams/THG_Sequence_v2.pptx"

# === COLORS ===
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(120, 120, 120)
SS_FILL = RGBColor(218, 232, 252)      # Light blue – internal subsystem
EXT_FILL = RGBColor(240, 240, 240)     # Light gray – external
CORE_FILL = RGBColor(255, 235, 235)    # Light red – core
PLATFORM_FILL = RGBColor(232, 245, 233)  # Light green
TEAM_FILL = RGBColor(255, 243, 224)    # Light orange – SS7 Team Coordination
FRAME_FILL = RGBColor(245, 245, 255)   # Very light blue – frames
ALT_FILL = RGBColor(255, 250, 240)     # Very light yellow – alt else
NEW_COLOR = RGBColor(139, 69, 19)      # Brown – new interfaces (IF-18+)

LINE_W = Pt(1.5)
DASH_W = Pt(1.0)


# ─────────────────────────────────────
# HELPERS
# ─────────────────────────────────────
def set_anchor_mid(shape):
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def add_participant(slide, cx, y, name, fill=SS_FILL, w=1.6, h=0.6):
    """Draw a participant box centered at cx, top at y."""
    x = cx - w / 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BLACK
    set_anchor_mid(shape)
    return cx, y + h  # return bottom center for lifeline start


def add_lifeline(slide, cx, y_start, y_end):
    """Draw a dashed vertical lifeline."""
    conn = slide.shapes.add_connector(
        1, Inches(cx), Inches(y_start), Inches(cx), Inches(y_end))
    conn.line.color.rgb = GRAY
    conn.line.width = Pt(0.75)
    ln = conn.line._ln
    prstDash = ln.makeelement(qn('a:prstDash'), {})
    prstDash.set('val', 'dash')
    ln.append(prstDash)


def add_activation(slide, cx, y_start, y_end, w=0.2):
    """Draw a thin activation rectangle on a lifeline."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - w/2), Inches(y_start), Inches(w), Inches(y_end - y_start))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(1.0)


def add_message(slide, x1, y, x2, label, dashed=False, label_offset=0):
    """Draw a message arrow from (x1,y) to (x2,y) with label above."""
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    conn.line.color.rgb = BLACK
    conn.line.width = Pt(1.2)

    if dashed:
        ln = conn.line._ln
        prstDash = ln.makeelement(qn('a:prstDash'), {})
        prstDash.set('val', 'dash')
        ln.append(prstDash)

    # Arrowhead at end
    ln = conn.line._ln
    tailEnd = ln.makeelement(qn('a:tailEnd'), {})
    if dashed:
        tailEnd.set('type', 'arrow')  # open arrowhead for returns
    else:
        tailEnd.set('type', 'triangle')  # filled arrowhead for messages
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')
    ln.append(tailEnd)

    # Label
    if label:
        mx = (x1 + x2) / 2
        lw = max(1.5, len(label) * 0.06)
        ly = y - 0.25 + label_offset
        tb = slide.shapes.add_textbox(
            Inches(mx - lw/2), Inches(ly), Inches(lw), Inches(0.25))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(8)
        p.font.color.rgb = BLACK


def add_self_call(slide, cx, y, label):
    """Draw a self-call (arrow looping back to same participant)."""
    loop_w = 0.6
    # Right segment
    conn1 = slide.shapes.add_connector(
        1, Inches(cx + 0.1), Inches(y), Inches(cx + loop_w), Inches(y))
    conn1.line.color.rgb = BLACK; conn1.line.width = Pt(1.0)
    # Down segment
    conn2 = slide.shapes.add_connector(
        1, Inches(cx + loop_w), Inches(y), Inches(cx + loop_w), Inches(y + 0.25))
    conn2.line.color.rgb = BLACK; conn2.line.width = Pt(1.0)
    # Left segment with arrow
    conn3 = slide.shapes.add_connector(
        1, Inches(cx + loop_w), Inches(y + 0.25), Inches(cx + 0.1), Inches(y + 0.25))
    conn3.line.color.rgb = BLACK; conn3.line.width = Pt(1.0)
    ln = conn3.line._ln
    tailEnd = ln.makeelement(qn('a:tailEnd'), {})
    tailEnd.set('type', 'triangle'); tailEnd.set('w', 'sm'); tailEnd.set('len', 'sm')
    ln.append(tailEnd)

    # Label
    if label:
        tb = slide.shapes.add_textbox(
            Inches(cx + loop_w + 0.05), Inches(y - 0.05), Inches(1.8), Inches(0.35))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(7); p.font.color.rgb = GRAY; p.font.italic = True


def add_frame(slide, x, y, w, h, operator, guard="", fill=FRAME_FILL):
    """Draw a frame box (alt, loop, opt) with operator label."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = RGBColor(150, 150, 180)
    rect.line.width = Pt(1.0)
    # Send to back (it'll be drawn first if called first)
    # Operator tab
    tab_w = max(0.5, len(operator) * 0.09)
    tab = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(tab_w), Inches(0.25))
    tab.fill.solid()
    tab.fill.fore_color.rgb = RGBColor(200, 200, 220)
    tab.line.color.rgb = RGBColor(150, 150, 180)
    tab.line.width = Pt(1.0)
    tf = tab.text_frame
    p = tf.paragraphs[0]
    p.text = operator
    p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = BLACK
    set_anchor_mid(tab)

    # Guard text
    if guard:
        tb = slide.shapes.add_textbox(
            Inches(x + tab_w + 0.1), Inches(y), Inches(2.0), Inches(0.25))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"[{guard}]"
        p.font.size = Pt(8); p.font.color.rgb = BLACK


def add_frame_divider(slide, x, y, w, label="else"):
    """Draw the dashed divider line in an alt frame."""
    conn = slide.shapes.add_connector(
        1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    conn.line.color.rgb = RGBColor(150, 150, 180)
    conn.line.width = Pt(0.75)
    ln = conn.line._ln
    prstDash = ln.makeelement(qn('a:prstDash'), {})
    prstDash.set('val', 'dash')
    ln.append(prstDash)

    tb = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y - 0.02), Inches(1.0), Inches(0.22))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"[{label}]"
    p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = BLACK


def add_title(slide, text, subtitle=""):
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BLACK
    if subtitle:
        tb2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.55), Inches(14.0), Inches(0.3))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(9); p2.font.color.rgb = GRAY; p2.font.italic = True


def add_note(slide, x, y, text, w=2.5, h=0.8):
    """Add a note box (folded corner style approximation)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 220)
    shape.line.color.rgb = RGBColor(200, 200, 150)
    shape.line.width = Pt(0.75)
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7); p.font.color.rgb = BLACK
    set_anchor_mid(shape)


# ─────────────────────────────────────
# SLIDE 1: UC1 — Request Prediction (OPTION B — WBGT + acclimatization + ML)
# ─────────────────────────────────────
def build_slide_prediction(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — UC1: Request Heat Risk Prediction",
              "Dual-model flow: PHS + WBGT screening  |  Interfaces: IF-01, IF-02, IF-03, IF-06, IF-10, IF-12, IF-14, IF-23")

    # Participants (7) — added SS6 for WBGT threshold lookup
    px = [1.2, 3.2, 5.2, 7.4, 9.8, 11.8, 14.2]
    py_top = 1.1
    names = ["Worker", "SS4\nMobile App", "SS5\nPlatform\nCore", "SS1\nPrediction\nEngine",
             "Weather\nAPI", "SS6\nRegulatory", "SS2\nAlert Engine"]
    fills = [EXT_FILL, SS_FILL, PLATFORM_FILL, CORE_FILL, EXT_FILL, SS_FILL, SS_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill, w=1.5, h=0.6)

    # Lifelines
    ll_top = py_top + 0.7
    ll_bot = 11.5
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activation rectangles
    add_activation(slide, px[0], 2.0, 11.2)   # Worker
    add_activation(slide, px[1], 2.2, 11.0)   # SS4 Mobile
    add_activation(slide, px[2], 2.5, 3.1)    # SS5 Platform (auth)
    add_activation(slide, px[2], 3.8, 4.4)    # SS5 Platform (profile + acclim)
    add_activation(slide, px[2], 9.2, 9.5)    # SS5 Platform (audit log)
    add_activation(slide, px[3], 3.5, 10.0)   # SS1 Prediction (long)
    add_activation(slide, px[4], 5.4, 5.9)    # Weather API
    add_activation(slide, px[5], 7.0, 7.5)    # SS6 Regulatory (WBGT thresholds)
    add_activation(slide, px[6], 10.0, 10.8)  # SS2 Alert

    y = 2.2

    # 1. Worker → SS4: select activity
    add_message(slide, px[0], y, px[1], "select activity, confirm clothing/PPE")
    y += 0.38

    # 2. SS4 → SS5: authenticate [IF-10]
    add_message(slide, px[1], y, px[2], "authenticate (token)")
    y += 0.32

    # 3. SS5 → SS4: JWT + RBAC [IF-10]
    add_message(slide, px[2], y, px[1], "JWT, RBAC, tenant ID", dashed=True)
    add_note(slide, 0.2, y - 0.12, "IF-10", w=0.45, h=0.25)
    y += 0.38

    # 4. SS4 → SS1: request prediction [IF-01]
    add_message(slide, px[1], y, px[3], "prediction request (activity, clothing, profile)")
    add_note(slide, 2.8, y + 0.12, "IF-01", w=0.45, h=0.25)
    y += 0.38

    # 5. SS1 → SS5: retrieve worker profile + acclim data [IF-12]
    add_message(slide, px[3], y, px[2], "retrieve profile + acclimatization data")
    y += 0.32

    # 6. SS5 → SS1: return profile [IF-12]
    add_message(slide, px[2], y, px[3], "profile, tenant config, acclim day, regulatory_standard", dashed=True)
    add_note(slide, 3.8, y + 0.12, "IF-12", w=0.45, h=0.25)
    y += 0.38

    # 7. SS1 self-call: metabolic rate + clothing factor
    add_self_call(slide, px[3], y, "metabolic rate (ISO 8996) + clothing factor")
    y += 0.38

    # 8. SS1 self-call: read acclimatization day
    add_self_call(slide, px[3], y, "read acclim day (1-14)")
    y += 0.38

    # 9. SS1 → Weather API: request env data [IF-14]
    add_message(slide, px[3], y, px[4], "env data request (lat, lon)")
    add_note(slide, 8.0, y + 0.12, "IF-14", w=0.45, h=0.25)
    y += 0.32

    # 10. Weather → SS1: return data [IF-14]
    add_message(slide, px[4], y, px[3], "temp, humidity, wind, solar, dew_point, pressure", dashed=True)
    y += 0.38

    # 11. SS1 self-call: execute PHS → Trec + Dlim
    add_self_call(slide, px[3], y, "execute PHS → Trec + Dlim")
    y += 0.38

    # --- WBGT SCREENING BLOCK (NEW) ---
    add_frame(slide, 6.0, y - 0.08, 7.0, 1.6, "ref", "WBGT Screening (F.4.10–F.4.13)")
    y += 0.28

    # 12. SS1 → SS6: WBGT threshold lookup [IF-23 NEW]
    add_message(slide, px[3], y, px[5], "WBGT thresholds (metabolic, regimen)")
    add_note(slide, 9.0, y + 0.12, "IF-23", w=0.5, h=0.25)
    y += 0.32

    # 13. SS6 → SS1: return AL/TLV thresholds
    add_message(slide, px[5], y, px[3], "AL/TLV per metabolic x regimen x acclim", dashed=True)
    y += 0.32

    # 14. SS1 self-call: derive WBGT, apply CAF
    add_self_call(slide, px[3], y, "derive WBGT, apply CAF")
    y += 0.32

    # 15. SS1 self-call: evaluate work-rest regimens
    add_self_call(slide, px[3], y, "evaluate 4 regimens vs AL/TLV")
    y += 0.32

    # 16. SS1 self-call: reconcile
    add_self_call(slide, px[3], y, "reconcile: risk = max(PHS, WBGT)")
    y += 0.38

    # --- ML COMPARISON (OPTIONAL) ---
    add_frame(slide, 6.0, y - 0.08, 3.0, 0.55, "opt", "ML model active")
    y += 0.25

    # 17. SS1 self-call: ML comparison
    add_self_call(slide, px[3], y, "store ml_predicted_trec")
    y += 0.5

    # 18. SS1 self-call: classify risk
    add_self_call(slide, px[3], y, "classify: GREEN/YELLOW/AMBER/RED/BLACK")
    y += 0.38

    # 19. SS1 → SS5: audit log [IF-06]
    add_message(slide, px[3], y, px[2], "log: inputs, outputs, wbgt_detail, ml_model_id")
    add_note(slide, 3.8, y + 0.12, "IF-06", w=0.45, h=0.25)
    y += 0.38

    # 20. SS1 → SS2: prediction result [IF-02]
    add_message(slide, px[3], y, px[6], "Trec, Dlim, risk, compliance_trigger, wbgt_effective")
    add_note(slide, 10.5, y + 0.12, "IF-02", w=0.45, h=0.25)
    y += 0.38

    # 21. SS2 → SS4: alert payload [IF-03]
    add_message(slide, px[6], y, px[1], "risk level, guidance, action required")
    add_note(slide, 7.5, y + 0.12, "IF-03", w=0.45, h=0.25)
    y += 0.35

    # 22. SS4 → Worker: display
    add_message(slide, px[1], y, px[0], "display prediction + risk + guidance", dashed=True)

    # Notes
    add_note(slide, 0.2, 1.85, "Starting Condition:\nWorker authenticated,\nprofile complete,\nat job site", w=1.6, h=0.55)
    add_note(slide, 13.0, 11.0, "Ending Condition:\nWorker sees Trec, Dlim,\nrisk (max of PHS/WBGT)\n+ guidance", w=2.5, h=0.6)


# ─────────────────────────────────────
# SLIDE 2: UC2/UC3 — Alert, Intervention & Escalation (OPTION B)
# ─────────────────────────────────────
def build_slide_alert(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — UC2/UC3: Alert, Intervention & Escalation",
              "Alert + intervention lifecycle  |  Interfaces: IF-02, IF-03, IF-04, IF-05, IF-07, IF-08, IF-25")

    # Participants (7) — same as before
    px = [1.3, 3.5, 5.7, 8.0, 10.3, 12.5, 14.7]
    py_top = 1.1
    names = ["SS1\nPrediction", "SS6\nRegulatory", "SS2\nAlert Engine", "SS4\nMobile App",
             "Worker", "SS3\nDashboard", "Site\nSupervisor"]
    fills = [CORE_FILL, SS_FILL, SS_FILL, SS_FILL, EXT_FILL, SS_FILL, EXT_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill)

    ll_top = py_top + 0.7
    ll_bot = 11.5
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activations
    add_activation(slide, px[0], 2.0, 2.5)    # SS1 (prediction in)
    add_activation(slide, px[0], 8.8, 9.2)    # SS1 (IF-25 ML feedback)
    add_activation(slide, px[1], 2.4, 3.0)    # SS6 (thresholds)
    add_activation(slide, px[2], 2.2, 10.8)   # SS2 main
    add_activation(slide, px[3], 4.6, 5.2)    # SS4 (alert delivery)
    add_activation(slide, px[3], 6.0, 8.5)    # SS4 (ack + intervention)
    add_activation(slide, px[4], 4.9, 5.2)    # Worker (receives)
    add_activation(slide, px[4], 6.2, 6.6)    # Worker (ack)
    add_activation(slide, px[5], 4.2, 4.5)    # SS3 (alert update)
    add_activation(slide, px[5], 9.3, 9.6)    # SS3 (outcome update)
    add_activation(slide, px[6], 10.0, 10.6)  # Supervisor

    y = 2.2

    # 1. SS1 → SS2: prediction result [IF-02]
    add_message(slide, px[0], y, px[2], "Trec, Dlim, risk, compliance_trigger")
    add_note(slide, 0.2, y + 0.1, "IF-02", w=0.45, h=0.22)
    y += 0.35

    # 2. SS2 → SS6: request thresholds [IF-08]
    add_message(slide, px[2], y, px[1], "request thresholds (org, state)")
    y += 0.3

    # 3. SS6 → SS2: return thresholds [IF-08]
    add_message(slide, px[1], y, px[2], "thresholds per state/org", dashed=True)
    add_note(slide, 2.2, y + 0.1, "IF-08", w=0.45, h=0.22)
    y += 0.3

    # 4. SS2 self-call: apply thresholds
    add_self_call(slide, px[2], y, "apply thresholds, detect risk change")
    y += 0.35

    # 5. SS2 self-call: generate intervention prescription (F.5.4a) — NEW
    add_self_call(slide, px[2], y, "gen prescription (fluid, cooling, rest, equip)")
    y += 0.3

    # 6. SS2 self-call: compute rest_duration method — NEW
    add_self_call(slide, px[2], y, "rest method: FIXED/TREC/DLIM_PROPORTIONAL")
    y += 0.38

    # === ALT FRAME: Risk changed vs no change ===
    frame_y = y - 0.1
    add_frame(slide, 3.8, frame_y, 11.3, 7.2, "alt", "risk level changed")
    y += 0.22

    # 7. SS2 → SS3: update dashboard [IF-04]
    add_message(slide, px[2], y, px[5], "alert + prescription to dashboard")
    add_note(slide, 10.8, y + 0.1, "IF-04", w=0.45, h=0.22)
    y += 0.35

    # 8. SS2 → SS4: push alert + intervention prescription [IF-03]
    add_message(slide, px[2], y, px[3], "alert + intervention prescription")
    add_note(slide, 5.5, y + 0.1, "IF-03", w=0.45, h=0.22)
    y += 0.3

    # 9. SS4 → Worker: push notification
    add_message(slide, px[3], y, px[4], "push notification + prescription")
    y += 0.35

    # === INNER ALT: Ack vs No Ack ===
    inner_frame_y = y - 0.05
    add_frame(slide, 4.2, inner_frame_y, 10.5, 4.6, "alt", "worker acknowledges within threshold")
    y += 0.22

    # --- ACK PATH ---
    # 10. Worker → SS4: acknowledge
    add_message(slide, px[4], y, px[3], "tap \"Acknowledged\"")
    y += 0.3

    # 11. SS4 → SS2: ack [IF-05]
    add_message(slide, px[3], y, px[2], "ack: worker, alert, GPS, pre_trec")
    add_note(slide, 4.5, y + 0.1, "IF-05", w=0.45, h=0.22)
    y += 0.3

    # 12. SS2 self-call: log ack [IF-07]
    add_self_call(slide, px[2], y, "log ack [IF-07]")
    y += 0.32

    # --- INTERVENTION EXECUTION (F.5.7a-g) — NEW ---
    add_frame(slide, 6.2, y - 0.06, 5.5, 1.55, "ref", "Intervention (F.5.7a-g)")
    y += 0.22

    # 13. SS4 self-call: display checklist + countdown
    add_self_call(slide, px[3], y, "display checklist + countdown")
    y += 0.3

    # 14. SS4 → SS2: intervention completed
    add_message(slide, px[3], y, px[2], "completed (GPS, completion_time)")
    y += 0.3

    # 15. SS2 self-call: record post_intervention_trec
    add_self_call(slide, px[2], y, "record post_intervention_trec")
    y += 0.28

    # 16. SS2 self-call: calculate effectiveness_score
    add_self_call(slide, px[2], y, "effectiveness_score (0-1)")
    y += 0.35

    # 17. SS2 → SS1: feed outcome to ML pipeline [IF-25 NEW]
    add_message(slide, px[2], y, px[0], "intervention outcome → ML pipeline")
    add_note(slide, 0.2, y + 0.1, "IF-25", w=0.5, h=0.22)
    y += 0.32

    # 18. SS2 → SS3: ack + effectiveness [IF-04]
    add_message(slide, px[2], y, px[5], "status: ack + effectiveness_score")
    y += 0.3

    # --- DIVIDER: no ack ---
    add_frame_divider(slide, 4.2, inner_frame_y + 3.35, 10.5, "worker does NOT acknowledge")
    y = inner_frame_y + 3.55

    # --- ESCALATION PATH ---
    # 19. SS2 self-call: escalation timer
    add_self_call(slide, px[2], y, "timer expires (AMBER=5min, RED=2min)")
    y += 0.32

    # 20. SS2 → Supervisor: escalation + prescription
    add_message(slide, px[2], y, px[6], "escalation + prescription details")
    y += 0.3

    # 21. Supervisor response (4 actions)
    add_frame(slide, 4.5, y - 0.06, 10.5, 0.55, "alt", "supervisor responds")
    y += 0.22

    add_message(slide, px[6], y, px[2], "action: Contact/Dispatch/Handled/Delegate", dashed=True)
    y += 0.35

    # 22. SS2 self-call: re-escalate if no response
    add_frame_divider(slide, 4.5, y - 0.08, 10.5, "no response in 10min")
    add_self_call(slide, px[2], y + 0.12, "re-escalate")
    y += 0.45

    # 23. SS2 self-call: log escalation [IF-07]
    add_self_call(slide, px[2], y, "log escalation [IF-07]")
    y += 0.3

    # --- DIVIDER: no risk change ---
    add_frame_divider(slide, 3.8, frame_y + 6.5, 11.3, "no risk level change")
    y = frame_y + 6.7

    add_message(slide, px[2], y, px[5], "dashboard update only (no alert)")

    # Notes
    add_note(slide, 0.2, 1.85, "Trigger:\nNew prediction\narrives from SS1", w=1.4, h=0.45)
    add_note(slide, 0.2, 10.8, "Ending Condition:\nAlert ack'd + intervention\ncomplete, OR escalated.\nIF-25 feeds ML.", w=2.0, h=0.65)
    add_note(slide, 12.0, 4.8,
             "Prescription:\n- fluid_vol\n- cooling_method\n- rest_duration\n- equip_removal_seq", w=2.2, h=0.65)


# ─────────────────────────────────────
# SLIDE 3: UC7 — Onboard New Worker
# ─────────────────────────────────────
def build_slide_onboard_worker(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — UC7: Onboard New Worker",
              "First-touch flow: invite → profile → first prediction in ≤ 5 min  |  Interfaces: IF-01, IF-02, IF-03, IF-10, IF-12, IF-14")

    # Participants (6)
    px = [1.5, 4.0, 6.5, 9.0, 11.5, 14.0]
    py_top = 1.1
    names = ["Worker", "SS4\nMobile App", "SS5\nPlatform Core", "SSO\nProvider",
             "SS1\nPrediction", "SS2\nAlert Engine"]
    fills = [EXT_FILL, SS_FILL, PLATFORM_FILL, EXT_FILL, CORE_FILL, SS_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill)

    ll_top = py_top + 0.7
    ll_bot = 11.2
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activations
    add_activation(slide, px[0], 2.0, 10.5)    # Worker
    add_activation(slide, px[1], 2.2, 10.2)    # SS4 Mobile
    add_activation(slide, px[2], 2.6, 4.0)     # SS5 (auth)
    add_activation(slide, px[2], 5.8, 6.5)     # SS5 (save profile)
    add_activation(slide, px[3], 3.0, 3.6)     # SSO Provider
    add_activation(slide, px[4], 7.2, 9.8)     # SS1 Prediction
    add_activation(slide, px[5], 9.5, 10.0)    # SS2 Alert

    y = 2.2

    # 1. Worker → SS4: Open app, enter invite code
    add_message(slide, px[0], y, px[1], "open app, enter invite code")
    y += 0.5

    # === ALT: SSO vs Email/Password ===
    add_frame(slide, 2.5, y - 0.1, 7.0, 2.2, "alt", "SSO enabled for org")
    y += 0.25

    # SSO path
    # 2. SS4 → SS5: initiate SSO
    add_message(slide, px[1], y, px[2], "initiate SSO login")
    y += 0.4

    # 3. SS5 → SSO: redirect to IdP
    add_message(slide, px[2], y, px[3], "SAML auth request")
    y += 0.35

    # 4. SSO → SS5: assertion
    add_message(slide, px[3], y, px[2], "SAML assertion (authenticated)", dashed=True)
    add_note(slide, 5.0, y + 0.12, "IF-15", w=0.5, h=0.25)
    y += 0.35

    # 5. SS5 → SS4: JWT
    add_message(slide, px[2], y, px[1], "JWT token, RBAC, tenant ID", dashed=True)
    add_note(slide, 1.5, y + 0.12, "IF-10", w=0.5, h=0.25)
    y += 0.4

    # Divider: else email/password
    add_frame_divider(slide, 2.5, y, 7.0, "email/password fallback")
    y += 0.25

    # 6. Worker → SS4: email + password
    add_message(slide, px[0], y, px[1], "enter email + create password")
    y += 0.25

    # SS4 → SS5: create account
    add_message(slide, px[1], y, px[2], "create account")
    y += 0.25

    add_message(slide, px[2], y, px[1], "JWT token", dashed=True)

    y += 0.6

    # === PROFILE COLLECTION ===
    add_note(slide, 13.0, y - 0.2, "Constraint:\nONBD-01: ≤ 5 min total\nONBD-02: ≤ 6 input fields", w=2.5, h=0.5)

    # 7. SS4 → Worker: display profile form
    add_message(slide, px[1], y, px[0], "display profile creation form", dashed=True)
    y += 0.4

    # 8. Worker → SS4: enter biometrics
    add_message(slide, px[0], y, px[1], "age, weight, height, acclimatization")
    y += 0.4

    # 9. SS4 self-call: validate
    add_self_call(slide, px[1], y, "validate inputs (range checks)")
    y += 0.4

    # 10. Worker → SS4: select activity
    add_message(slide, px[0], y, px[1], "select work activity from category list")
    y += 0.4

    # 11. Worker → SS4: select clothing
    add_message(slide, px[0], y, px[1], "select clothing / PPE")
    y += 0.4

    # 12. SS4 → SS5: save profile
    add_message(slide, px[1], y, px[2], "save worker profile")
    y += 0.35

    # 13. SS5 → SS4: profile saved
    add_message(slide, px[2], y, px[1], "profile ID confirmed", dashed=True)
    add_note(slide, 1.5, y + 0.12, "IF-12", w=0.5, h=0.25)
    y += 0.55

    # === AUTO-TRIGGER FIRST PREDICTION ===
    add_frame(slide, 2.5, y - 0.1, 12.0, 2.5, "ref", "UC1: First Prediction (auto-triggered)")
    y += 0.3

    # 14. SS4 → SS1: request first prediction [IF-01]
    add_message(slide, px[1], y, px[4], "request prediction (activity, clothing, profile ID)")
    add_note(slide, 6.0, y + 0.12, "IF-01", w=0.5, h=0.25)
    y += 0.5

    # 15. SS1 self-calls (condensed)
    add_self_call(slide, px[4], y, "metabolic rate → clothing factor → PHS")
    y += 0.45

    # 16. SS1 self-call: weather
    add_self_call(slide, px[4], y, "retrieve env data [IF-14]")
    y += 0.45

    # 16a. SS1 self-call: initialize acclimatization protocol — NEW
    add_self_call(slide, px[4], y, "init acclim (started_at, day, status)")
    y += 0.35

    # 16b. SS1 self-call: set protocol day from declared status — NEW (ONBD-08/09)
    add_self_call(slide, px[4], y, "day from status: New=0, Returning=1, Acclim=14")
    y += 0.45

    # 17. SS1 → SS2: prediction [IF-02]
    add_message(slide, px[4], y, px[5], "prediction result: Trec, Dlim, risk level")
    add_note(slide, 11.5, y + 0.12, "IF-02", w=0.5, h=0.25)
    y += 0.45

    # 18. SS2 → SS4: alert [IF-03]
    add_message(slide, px[5], y, px[1], "risk level + guidance")
    add_note(slide, 7.5, y + 0.12, "IF-03", w=0.5, h=0.25)
    y += 0.45

    # 19. SS4 → Worker: display first prediction + tutorial
    add_message(slide, px[1], y, px[0], "first prediction + tutorial overlay", dashed=True)

    # Notes
    add_note(slide, 0.2, 1.85, "Starting Condition:\nMobile device,\nvalid org invite,\norg exists in system", w=1.8, h=0.6)
    add_note(slide, 13.0, 10.3, "Ending Condition:\nProfile complete,\nfirst prediction received,\ntutorial shown.\nTotal time ≤ 5 min.", w=2.5, h=0.7)


# ─────────────────────────────────────
# SLIDE 4: UC9 — Provision Organization
# ─────────────────────────────────────
def build_slide_provision_org(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — UC9: Provision Organization",
              "Enterprise onboarding: tenant → SSO → RBAC → thresholds → sites → billing  |  Interfaces: IF-10, IF-12, IF-15")

    # Participants (5)
    px = [2.0, 5.0, 8.0, 11.0, 14.0]
    py_top = 1.1
    names = ["Employer\nAdmin", "SS5\nPlatform Core", "SSO\nProvider",
             "SS6\nRegulatory", "Billing\nSystem"]
    fills = [EXT_FILL, PLATFORM_FILL, EXT_FILL, SS_FILL, EXT_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill)

    ll_top = py_top + 0.7
    ll_bot = 11.0
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activations
    add_activation(slide, px[0], 2.0, 10.5)    # Admin
    add_activation(slide, px[1], 2.3, 10.2)    # SS5 Platform
    add_activation(slide, px[2], 3.5, 4.5)     # SSO
    add_activation(slide, px[3], 5.5, 7.2)     # SS6 Regulatory
    add_activation(slide, px[4], 9.2, 10.0)    # Billing

    y = 2.2

    # F.1.1: Create Tenant
    add_note(slide, 0.2, y - 0.2, "FFBD: F.1.1\nCreate Tenant", w=1.5, h=0.4)
    add_message(slide, px[0], y, px[1], "create new organization (name, industry, size)")
    y += 0.4

    # SS5 self-call: provision tenant
    add_self_call(slide, px[1], y, "create tenant schema, assign tenant ID")
    y += 0.5

    # SS5 → Admin: tenant created
    add_message(slide, px[1], y, px[0], "tenant ID, admin dashboard URL", dashed=True)
    y += 0.6

    # F.1.2: Configure SSO
    add_note(slide, 0.2, y - 0.2, "FFBD: F.1.2\nConfigure SSO", w=1.5, h=0.4)

    # Alt: SSO or skip
    add_frame(slide, 3.5, y - 0.1, 7.0, 1.8, "alt", "enterprise SSO required")
    y += 0.25

    add_message(slide, px[0], y, px[1], "configure SSO (IdP metadata URL)")
    y += 0.4

    add_message(slide, px[1], y, px[2], "validate SSO configuration")
    y += 0.35

    add_message(slide, px[2], y, px[1], "SSO validated, test login successful", dashed=True)
    add_note(slide, 6.0, y + 0.12, "IF-15", w=0.5, h=0.25)
    y += 0.4

    add_frame_divider(slide, 3.5, y, 7.0, "email/password only")
    y += 0.2

    add_self_call(slide, px[1], y, "configure email auth defaults")
    y += 0.55

    # F.1.3: Set Up RBAC
    add_note(slide, 0.2, y - 0.2, "FFBD: F.1.3\nSet Up RBAC", w=1.5, h=0.4)
    add_message(slide, px[0], y, px[1], "define roles: admin, manager, supervisor, worker")
    y += 0.4

    add_self_call(slide, px[1], y, "create RBAC roles + permission matrix")
    y += 0.4

    add_message(slide, px[1], y, px[0], "RBAC configured, invite links generated", dashed=True)
    add_note(slide, 1.5, y + 0.12, "IF-10", w=0.5, h=0.25)
    y += 0.6

    # F.1.4: Configure Alert Thresholds
    add_note(slide, 0.2, y - 0.2, "FFBD: F.1.4\nAlert Thresholds", w=1.5, h=0.4)
    add_message(slide, px[0], y, px[1], "select operating states (CA, OR, WA...)")
    y += 0.4

    add_message(slide, px[1], y, px[3], "request default thresholds for selected states")
    y += 0.35

    add_message(slide, px[3], y, px[1], "state thresholds, compliance requirements", dashed=True)
    add_note(slide, 7.5, y + 0.12, "IF-08", w=0.5, h=0.25)
    y += 0.4

    add_message(slide, px[1], y, px[0], "thresholds loaded, ready to customize", dashed=True)
    y += 0.35

    # Alt: customize or accept defaults
    add_frame(slide, 3.5, y - 0.1, 4.5, 0.7, "opt", "admin customizes thresholds")
    y += 0.25

    add_message(slide, px[0], y, px[1], "adjust threshold values per site")
    y += 0.65

    # F.1.5: Register Job Sites
    add_note(slide, 0.2, y - 0.2, "FFBD: F.1.5\nRegister Sites", w=1.5, h=0.4)
    add_message(slide, px[0], y, px[1], "register job sites (name, location, indoor/outdoor)")
    y += 0.4

    add_self_call(slide, px[1], y, "geocode sites, link to weather zones")
    add_note(slide, 7.0, y - 0.05, "IF-12: site locations\nstored for SS1", w=1.5, h=0.35)
    y += 0.4

    # SS5 → SS6: auto-assign regulatory standard per site — NEW
    add_message(slide, px[1], y, px[3], "auto-assign regulatory standard per jurisdiction")
    y += 0.3

    add_message(slide, px[3], y, px[1], "standard assigned (e.g. CA §3395 outdoor)", dashed=True)
    y += 0.45

    # F.1.5b: Bulk Import Workers — NEW
    add_note(slide, 0.2, y - 0.2, "FFBD: F.1.5b\nBulk Import", w=1.5, h=0.4)
    add_frame(slide, 3.5, y - 0.1, 4.5, 1.2, "opt", "admin imports workers")
    y += 0.2

    add_message(slide, px[0], y, px[1], "bulk import workers (Excel/CSV/paste)")
    y += 0.3

    add_self_call(slide, px[1], y, "validate rows, report errors per row")
    y += 0.3

    add_self_call(slide, px[1], y, "send invites (email/SMS), track status")
    y += 0.55

    # F.1.6: Activate Billing
    add_note(slide, 0.2, y - 0.2, "FFBD: F.1.6\nActivate Billing", w=1.5, h=0.4)
    add_message(slide, px[0], y, px[1], "select plan, enter payment method")
    y += 0.35

    add_message(slide, px[1], y, px[4], "create subscription (plan, payment)")
    y += 0.3

    add_message(slide, px[4], y, px[1], "subscription active, first invoice scheduled", dashed=True)
    y += 0.35

    add_message(slide, px[1], y, px[0], "organization live — ready for worker onboarding", dashed=True)

    # Notes
    add_note(slide, 0.2, 1.85, "Starting Condition:\nNew enterprise customer\nsigns contract", w=1.8, h=0.5)
    add_note(slide, 12.5, 10.3, "Ending Condition:\nTenant provisioned,\nSSO configured, RBAC set,\nsites + jurisdiction assigned,\nbulk import available,\nbilling active.\nMaps to FFBD F.1.1–F.1.6", w=2.8, h=0.9)


# ─────────────────────────────────────
# SLIDE 5: UC6 — Generate Compliance Report
# ─────────────────────────────────────
def build_slide_compliance_report(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — UC6: Generate Compliance Report",
              "Compliance value chain: request → templates → data → compile → export  |  Interfaces: IF-04, IF-06, IF-07, IF-09, IF-11, IF-13, IF-17")

    # Participants (5)
    px = [1.8, 4.8, 8.0, 11.0, 14.0]
    py_top = 1.1
    names = ["Safety\nManager", "SS3\nDashboard", "SS5\nPlatform Core",
             "SS6\nRegulatory", "Audit\nSystems"]
    fills = [EXT_FILL, SS_FILL, PLATFORM_FILL, SS_FILL, EXT_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill)

    ll_top = py_top + 0.7
    ll_bot = 11.0
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activations
    add_activation(slide, px[0], 2.0, 10.5)    # Safety Manager
    add_activation(slide, px[1], 2.3, 10.0)    # SS3 Dashboard
    add_activation(slide, px[2], 3.8, 7.0)     # SS5 Platform (data)
    add_activation(slide, px[3], 2.8, 3.8)     # SS6 (templates)
    add_activation(slide, px[3], 7.5, 8.5)     # SS6 (compliance check)
    add_activation(slide, px[4], 9.0, 9.8)     # Audit Systems

    y = 2.2

    # 1. Safety Manager → SS3: request report
    add_message(slide, px[0], y, px[1], "generate compliance report (site, date range, type)")
    y += 0.5

    # 2. SS3 → SS5: authenticate manager
    add_message(slide, px[1], y, px[2], "validate manager permissions")
    y += 0.35

    add_message(slide, px[2], y, px[1], "RBAC: manager role confirmed, org scope", dashed=True)
    add_note(slide, 4.5, y + 0.12, "IF-11", w=0.5, h=0.25)
    y += 0.55

    # 3. SS3 → SS6: request report template [IF-13]
    add_message(slide, px[1], y, px[3], "request template (report type, state)")
    add_note(slide, 8.0, y + 0.12, "IF-13", w=0.5, h=0.25)
    y += 0.4

    # 4. SS6 → SS3: return template [IF-09]
    add_message(slide, px[3], y, px[1], "template + state-specific format + required fields", dashed=True)
    add_note(slide, 6.0, y + 0.12, "IF-09", w=0.5, h=0.25)
    y += 0.6

    # 5. SS3 → SS5: query prediction data [IF-06 data]
    add_message(slide, px[1], y, px[2], "query predictions for site + date range")
    y += 0.35

    add_message(slide, px[2], y, px[1], "prediction logs (Trec, Dlim, risk, timestamp per worker)", dashed=True)
    add_note(slide, 4.5, y + 0.12, "IF-06 data", w=0.7, h=0.25)
    y += 0.5

    # 6. SS3 → SS5: query alert data [IF-07 data]
    add_message(slide, px[1], y, px[2], "query alert lifecycle for site + date range")
    y += 0.35

    add_message(slide, px[2], y, px[1], "alert logs (triggers, acks, escalations, resolutions)", dashed=True)
    add_note(slide, 4.5, y + 0.12, "IF-07 data", w=0.7, h=0.25)
    y += 0.5

    # 7. SS3 → SS5: query acknowledgment data
    add_message(slide, px[1], y, px[2], "query worker acknowledgments + supervisor interventions")
    y += 0.35

    add_message(slide, px[2], y, px[1], "ack records (worker ID, timestamp, response time)", dashed=True)
    y += 0.55

    # 8. SS6 self-call: resolve template per site jurisdiction — NEW
    add_self_call(slide, px[3], y, "resolve per jurisdiction (CA §3395/§3396)")
    y += 0.4

    # 9. SS3 self-call: compile report
    add_self_call(slide, px[1], y, "compile report: merge template + data")
    y += 0.35

    # 9a. SS6 self-call: flag incomplete data — NEW
    add_self_call(slide, px[1], y, "flag sites added mid-period + date bounds")
    y += 0.4

    # 10. SS3 → SS6: validate compliance [IF-13]
    add_message(slide, px[1], y, px[3], "validate report against compliance rules")
    y += 0.35

    # Alt: compliant vs gaps
    add_frame(slide, 6.5, y - 0.1, 6.0, 1.2, "alt", "all compliance requirements met")
    y += 0.22

    add_message(slide, px[3], y, px[1], "compliance: PASS — all fields complete", dashed=True)
    y += 0.32

    add_frame_divider(slide, 6.5, y, 6.0, "compliance gaps found")
    y += 0.22

    add_message(slide, px[3], y, px[1], "compliance: GAPS — missing fields listed", dashed=True)
    y += 0.5

    # 10a. SS6 self-call: embed verification metadata — NEW
    add_self_call(slide, px[1], y, "embed data hash, timestamp, algo version")
    y += 0.45

    # 11. SS3 → Audit Systems: export [IF-17]
    add_message(slide, px[1], y, px[4], "export audit copy (immutable)")
    add_note(slide, 10.5, y + 0.12, "IF-17", w=0.5, h=0.25)
    y += 0.35

    add_message(slide, px[4], y, px[1], "audit receipt (hash, timestamp)", dashed=True)
    y += 0.4

    # 12. SS3 → Safety Manager: deliver report
    add_message(slide, px[1], y, px[0], "download: compliance report (PDF + Excel)", dashed=True)
    add_note(slide, 0.2, y + 0.12, "DASH-05", w=0.6, h=0.25)

    # Notes
    add_note(slide, 0.2, 1.85, "Starting Condition:\nManager authenticated,\nworkers have been active,\nprediction data exists", w=1.8, h=0.6)
    add_note(slide, 12.5, 10.3, "Ending Condition:\nOSHA-aligned report\ngenerated, validated,\naudit copy stored.\nVerification metadata\nembedded. PDF + Excel.\nMaps to FFBD F.6.6–F.6.7", w=2.8, h=0.9)
    add_note(slide, 12.5, 2.0, "Report types:\n• OSHA HIIP\n• CA §3395/§3396\n• OR OAR 437-002-0156\n• WA WAC 296-62-095\n• HIPP template\n• Training records", w=2.5, h=1.0)


# ─────────────────────────────────────
# SLIDE 6: Team Coordination (PATENT) — NEW
# ─────────────────────────────────────
def build_slide_team_coordination(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — Team Coordination (PATENT)",
              "Session lifecycle: create → monitor → rotate → close  |  Interfaces: IF-18, IF-19, IF-20, IF-21, IF-22, IF-24")

    # Participants (6)
    px = [1.3, 3.5, 5.8, 8.2, 10.8, 13.8]
    py_top = 1.1
    names = ["Safety\nManager", "SS3\nDashboard", "SS7\nTeam\nCoordination",
             "SS1\nPrediction", "SS4\nMobile App", "SS5\nPlatform\nCore"]
    fills = [EXT_FILL, SS_FILL, TEAM_FILL, CORE_FILL, SS_FILL, PLATFORM_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill, w=1.5, h=0.6)

    ll_top = py_top + 0.7
    ll_bot = 11.5
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activations
    add_activation(slide, px[0], 2.0, 2.5)     # Manager (create)
    add_activation(slide, px[0], 8.5, 9.0)     # Manager (accept rotation)
    add_activation(slide, px[0], 10.2, 10.5)   # Manager (end session)
    add_activation(slide, px[1], 2.2, 2.7)     # SS3 (create cmd)
    add_activation(slide, px[1], 7.0, 7.5)     # SS3 (rotation rec)
    add_activation(slide, px[1], 8.5, 9.0)     # SS3 (accept)
    add_activation(slide, px[1], 10.2, 10.5)   # SS3 (end)
    add_activation(slide, px[2], 2.5, 10.8)    # SS7 (long running)
    add_activation(slide, px[3], 4.0, 4.5)     # SS1 (prediction push)
    add_activation(slide, px[4], 9.0, 9.5)     # SS4 (rotation notify)
    add_activation(slide, px[5], 2.8, 3.2)     # SS5 (auth)
    add_activation(slide, px[5], 9.5, 9.9)     # SS5 (audit)
    add_activation(slide, px[5], 10.5, 10.8)   # SS5 (session close audit)

    y = 2.2

    # 1. Manager → SS3: create team session
    add_message(slide, px[0], y, px[1], "create session (workers, site, duration)")
    y += 0.32

    # 2. SS3 → SS7: create session command [IF-20]
    add_message(slide, px[1], y, px[2], "create session command")
    add_note(slide, 3.8, y + 0.1, "IF-20", w=0.45, h=0.22)
    y += 0.35

    # 3. SS5 → SS7: auth context [IF-24]
    add_message(slide, px[5], y, px[2], "auth context (JWT, RBAC)")
    add_note(slide, 9.5, y + 0.1, "IF-24", w=0.45, h=0.22)
    y += 0.32

    # 4. SS7 self-call: initialize
    add_self_call(slide, px[2], y, "initialize per-member tracking")
    y += 0.45

    # === LOOP: per prediction cycle ===
    add_frame(slide, 4.2, y - 0.08, 6.0, 2.6, "loop", "per prediction cycle")
    y += 0.25

    # 5. SS1 → SS7: prediction update [IF-18]
    add_message(slide, px[3], y, px[2], "member risk_level, dlim, recovery_rate")
    add_note(slide, 5.5, y + 0.1, "IF-18", w=0.45, h=0.22)
    y += 0.35

    # 6. SS7 self-call: detect rotation need
    add_self_call(slide, px[2], y, "detect: Dlim < 30min or risk >= AMBER?")
    y += 0.35

    # 7. SS7 self-call: evaluate replacements
    add_self_call(slide, px[2], y, "eval replacements: GREEN/YELLOW, Dlim > 2hr")
    y += 0.35

    # 8. SS7 self-call: generate recommendation
    add_self_call(slide, px[2], y, "rotation rec: IMMEDIATE / WITHIN_15MIN / ADVISORY")
    y += 0.35

    # 9. SS7 → SS3: push rotation recommendation [IF-19]
    add_message(slide, px[2], y, px[1], "push rotation recommendation")
    add_note(slide, 3.8, y + 0.1, "IF-19", w=0.45, h=0.22)
    y += 0.55

    # === ALT: Manager acts vs auto-expire ===
    add_frame(slide, 1.0, y - 0.08, 14.0, 2.4, "alt", "manager accepts rotation")
    y += 0.25

    # 10. Manager → SS3 → SS7: accept [IF-20]
    add_message(slide, px[0], y, px[1], "accept rotation")
    y += 0.28

    add_message(slide, px[1], y, px[2], "accept rotation command")
    add_note(slide, 3.8, y + 0.1, "IF-20", w=0.45, h=0.22)
    y += 0.32

    # 11. SS7 → SS4: notify affected workers [IF-21]
    add_message(slide, px[2], y, px[4], "rotation notification to workers")
    add_note(slide, 8.5, y + 0.1, "IF-21", w=0.45, h=0.22)
    y += 0.32

    # 12. SS7 → SS5: log rotation decision [IF-22]
    add_message(slide, px[2], y, px[5], "log rotation decision")
    add_note(slide, 9.5, y + 0.1, "IF-22", w=0.45, h=0.22)
    y += 0.35

    # Divider: auto-expire
    add_frame_divider(slide, 1.0, y, 14.0, "no action within 30 min")
    y += 0.22

    # 13. SS7 self-call: expire recommendation
    add_self_call(slide, px[2], y, "expire unacted recommendation")
    y += 0.55

    # === END SESSION ===
    # 14. Manager → SS3 → SS7: end session [IF-20]
    add_message(slide, px[0], y, px[1], "end session")
    y += 0.25

    add_message(slide, px[1], y, px[2], "end session command")
    y += 0.32

    # 15. SS7 self-call: calculate session summary
    add_self_call(slide, px[2], y, "summary: rotations, avg risk, incidents")
    y += 0.35

    # 16. SS7 → SS5: log session close [IF-22]
    add_message(slide, px[2], y, px[5], "log session close to audit trail")
    add_note(slide, 9.5, y + 0.1, "IF-22", w=0.45, h=0.22)

    # Notes
    add_note(slide, 0.2, 1.85, "Starting Condition:\nManager authenticated,\nworkers assigned to site,\npredictions active", w=1.6, h=0.6)
    add_note(slide, 13.0, 11.0, "Ending Condition:\nSession closed,\nall rotations logged,\nsummary generated.\nPATENT feature.", w=2.5, h=0.65)


# ─────────────────────────────────────
# SLIDE 7: Pre-Event Risk Stratification (PATENT) — NEW
# ─────────────────────────────────────
def build_slide_stratification(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — Pre-Event Risk Stratification (PATENT)",
              "Proactive risk assessment: forecast → per-worker scoring → mitigation  |  Interfaces: IF-06, IF-12, IF-14")

    # Participants (5)
    px = [1.8, 4.8, 8.0, 11.5, 14.0]
    py_top = 1.1
    names = ["Safety\nManager", "SS3\nDashboard", "SS1\nPrediction\nEngine",
             "Weather\nAPI", "SS5\nPlatform\nCore"]
    fills = [EXT_FILL, SS_FILL, CORE_FILL, EXT_FILL, PLATFORM_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill)

    ll_top = py_top + 0.7
    ll_bot = 11.0
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activations
    add_activation(slide, px[0], 2.0, 10.0)    # Manager
    add_activation(slide, px[1], 2.3, 9.6)     # SS3 Dashboard
    add_activation(slide, px[2], 3.0, 8.5)     # SS1 Prediction
    add_activation(slide, px[3], 3.5, 4.2)     # Weather API
    add_activation(slide, px[4], 4.5, 5.2)     # SS5 (roster)
    add_activation(slide, px[4], 8.5, 9.0)     # SS5 (audit)

    y = 2.2

    # 1. Manager → SS3: create assessment
    add_message(slide, px[0], y, px[1], "create stratification (type, target date, site)")
    y += 0.45

    # 2. SS3 → SS1: request stratification
    add_message(slide, px[1], y, px[2], "request stratification assessment")
    y += 0.45

    # 3. SS1 → Weather API: retrieve forecast [IF-14]
    add_message(slide, px[2], y, px[3], "retrieve forecast for target date/location")
    add_note(slide, 8.5, y + 0.12, "IF-14", w=0.45, h=0.25)
    y += 0.35

    add_message(slide, px[3], y, px[2], "forecast: temp, humidity, wind, solar, WBGT", dashed=True)
    y += 0.5

    # 4. SS1 → SS5: retrieve worker roster [IF-12]
    add_message(slide, px[2], y, px[4], "retrieve worker roster with profiles")
    add_note(slide, 9.0, y + 0.12, "IF-12", w=0.45, h=0.25)
    y += 0.35

    add_message(slide, px[4], y, px[2], "roster: profiles, acclim status, history", dashed=True)
    y += 0.5

    # 5. SS1 self-call: map activities → metabolic categories
    add_self_call(slide, px[2], y, "map activities → metabolic categories")
    y += 0.45

    # 6. SS1 self-call: per-worker risk score
    add_self_call(slide, px[2], y, "per-worker risk score (0-100)")
    y += 0.4

    add_note(slide, 10.5, y - 0.3, "Score inputs:\n• forecast severity\n• acclim day (1-14)\n• BMI / age\n• recent incidents\n• clothing factor", w=2.2, h=0.8)

    # 7. SS1 self-call: rank by vulnerability
    add_self_call(slide, px[2], y, "rank by vulnerability (non-acclim, high BMI, incidents)")
    y += 0.45

    # 8. SS1 self-call: generate mitigations
    add_self_call(slide, px[2], y, "generate mitigation recommendations")
    y += 0.5

    # 9. SS1 → SS3: return assessment
    add_message(slide, px[2], y, px[1], "assessment: risk_score, worker_scores, recommendations")
    y += 0.45

    # 10. Manager → SS3: review and publish
    add_frame(slide, 1.5, y - 0.08, 5.0, 0.55, "opt", "manager publishes")
    y += 0.25

    add_message(slide, px[0], y, px[1], "publish assessment (DRAFT → PUBLISHED)")
    y += 0.55

    # 11. SS1 → SS5: log to audit trail [IF-06]
    add_message(slide, px[2], y, px[4], "log assessment to audit trail")
    add_note(slide, 9.0, y + 0.12, "IF-06", w=0.45, h=0.25)

    # Notes
    add_note(slide, 0.2, 1.85, "Starting Condition:\nManager authenticated,\nevent/shift planned,\nworker roster exists", w=1.8, h=0.55)
    add_note(slide, 12.5, 10.0, "Ending Condition:\nAssessment created,\nworkers ranked by risk,\nmitigations recommended.\nPATENT feature.", w=2.5, h=0.7)


# ─────────────────────────────────────
# SLIDE 8: Activity Switch Mid-Shift (UC13) — NEW
# ─────────────────────────────────────
def build_slide_activity_switch(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    add_title(slide,
              "Sequence Diagram — UC13: Activity Switch Mid-Shift",
              "Worker changes activity → recalculate prediction  |  Interfaces: IF-01, IF-02, IF-03, IF-06")

    # Participants (5)
    px = [1.8, 4.8, 8.0, 11.5, 14.0]
    py_top = 1.1
    names = ["Worker", "SS4\nMobile App", "SS1\nPrediction\nEngine",
             "SS2\nAlert Engine", "SS5\nPlatform\nCore"]
    fills = [EXT_FILL, SS_FILL, CORE_FILL, SS_FILL, PLATFORM_FILL]

    for x, name, fill in zip(px, names, fills):
        add_participant(slide, x, py_top, name, fill)

    ll_top = py_top + 0.7
    ll_bot = 11.0
    for x in px:
        add_lifeline(slide, x, ll_top, ll_bot)

    # Activations
    add_activation(slide, px[0], 2.0, 4.5)     # Worker (input)
    add_activation(slide, px[1], 2.2, 10.0)    # SS4 Mobile
    add_activation(slide, px[2], 5.5, 9.0)     # SS1 Prediction
    add_activation(slide, px[3], 8.5, 9.5)     # SS2 Alert
    add_activation(slide, px[4], 9.0, 9.5)     # SS5 (audit)

    y = 2.2

    # 1. Worker → SS4: select "Change Activity"
    add_message(slide, px[0], y, px[1], "select \"Change Activity\"")
    y += 0.45

    # 2. SS4 self-call: display current activity + clothing
    add_self_call(slide, px[1], y, "display current activity + clothing")
    y += 0.45

    # 3. Worker → SS4: select new activity
    add_message(slide, px[0], y, px[1], "select new activity")
    y += 0.4

    # 4. SS4 self-call: reclassify metabolic category
    add_self_call(slide, px[1], y, "reclassify metabolic (NIOSH Table 8-1)")
    y += 0.45

    # 5. [OPT] Worker → SS4: update clothing/PPE
    add_frame(slide, 1.0, y - 0.08, 6.0, 0.55, "opt", "clothing changed")
    y += 0.25

    add_message(slide, px[0], y, px[1], "update clothing/PPE → recalculate CAF")
    y += 0.55

    # 6. [OPT] Worker → SS4: update surface type
    add_frame(slide, 1.0, y - 0.08, 6.0, 0.55, "opt", "location changed")
    y += 0.25

    add_message(slide, px[0], y, px[1], "update surface_type (outdoor ↔ indoor)")
    y += 0.55

    # 7. SS4 → SS1: submit new prediction [IF-01]
    add_message(slide, px[1], y, px[2], "new prediction with updated inputs")
    add_note(slide, 5.5, y + 0.12, "IF-01", w=0.45, h=0.25)
    y += 0.45

    # 8. SS1 self-call: close previous window
    add_self_call(slide, px[2], y, "close previous prediction window (end ts)")
    y += 0.45

    # 9. SS1 self-call: execute full UC1 dual-model flow
    add_self_call(slide, px[2], y, "execute full UC1 dual-model flow")
    y += 0.5

    # 10. SS1 → SS2: new prediction result [IF-02]
    add_message(slide, px[2], y, px[3], "new prediction: Trec, Dlim, risk")
    add_note(slide, 9.0, y + 0.12, "IF-02", w=0.45, h=0.25)
    y += 0.45

    # 11. [ALT: risk changed] SS2 → SS4: alert [IF-03]
    add_frame(slide, 3.5, y - 0.08, 9.0, 0.55, "alt", "risk level changed")
    y += 0.25

    add_message(slide, px[3], y, px[1], "alert + updated intervention")
    add_note(slide, 7.0, y + 0.12, "IF-03", w=0.45, h=0.25)
    y += 0.55

    # 12. SS1 → SS5: log activity switch [IF-06]
    add_message(slide, px[2], y, px[4], "log switch: old/new activity, GPS")
    add_note(slide, 9.5, y + 0.12, "IF-06", w=0.45, h=0.25)

    # Notes
    add_note(slide, 0.2, 1.85, "Starting Condition:\nWorker has active\nprediction, changes\nwork activity mid-shift", w=1.8, h=0.6)
    add_note(slide, 12.5, 10.0, "Ending Condition:\nPrevious window closed,\nnew prediction active,\nalerts updated if risk\nchanged. Audit logged.", w=2.5, h=0.7)
    add_note(slide, 12.5, 5.5, "Activity switch\ntriggers full UC1\ndual-model flow\n(PHS + WBGT)", w=2.0, h=0.6)


# ─────────────────────────────────────
# MAIN
# ─────────────────────────────────────
def build():
    prs = Presentation(TEMPLATE)

    # Delete template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    build_slide_prediction(prs)          # Slide 1: UC1
    build_slide_alert(prs)               # Slide 2: UC2/UC3
    build_slide_onboard_worker(prs)      # Slide 3: UC7
    build_slide_provision_org(prs)       # Slide 4: UC9
    build_slide_compliance_report(prs)   # Slide 5: UC6
    build_slide_team_coordination(prs)   # Slide 6: Team Coord (PATENT)
    build_slide_stratification(prs)      # Slide 7: Stratification (PATENT)
    build_slide_activity_switch(prs)     # Slide 8: UC13 Activity Switch

    prs.save(OUTPUT)
    print(f"✓ Saved: {OUTPUT}")
    print(f"  8 slides: UC1 + UC2/3 + UC7 + UC9 + UC6 + Team Coord + Stratification + UC13")


if __name__ == "__main__":
    build()
