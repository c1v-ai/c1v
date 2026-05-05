#!/usr/bin/env python3
"""
Generate Cornell-standard Use Case Diagram for Team Heat Guard (THG)
Following CESYS521/522 formatting rules:
  - Black and white only
  - Actors as square-cornered boxes OUTSIDE boundary
  - Use cases as rounded rectangles INSIDE boundary
  - Dashed rectangle = system boundary
  - Solid lines: actor-to-use-case
  - Dashed lines: include/extend relationships (detailed slide)
  - 6 actors (4 left, 2 right), 14 use cases

Two slides:
  Slide 1: Basic  — actors, boundary, use cases, actor connections
  Slide 2: Detailed — adds include/extend relationships with labels

Usage:
  pip install python-pptx
  python generate_use_case_diagram_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ─── CONSTANTS ───────────────────────────────────────────────────

FONT_NAME = "Arial"
ACTOR_FONT = Pt(9)
UC_FONT = Pt(8)
REL_FONT = Pt(6.5)
TITLE_FONT = Pt(18)
BND_LABEL_FONT = Pt(10)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LINE_W = Pt(1.0)
DASH_W = Pt(2.0)
REL_W = Pt(0.8)

SW = 13.333
SH = 7.5

# Actor box
AW = 1.3
AH = 0.5

# Use case box (rounded rect)
UW = 2.0
UH = 0.5


# ─── DRAWING ─────────────────────────────────────────────────────

def add_actor(slide, left, top, text):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(AW), Inches(AH)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BLACK
    s.line.width = LINE_W
    tf = s.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text.upper()
    r.font.size = ACTOR_FONT
    r.font.name = FONT_NAME
    r.font.color.rgb = BLACK
    r.font.bold = True


def add_usecase(slide, left, top, text):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(UW), Inches(UH)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BLACK
    s.line.width = LINE_W
    # Reduce corner rounding
    s.adjustments[0] = 0.16667
    tf = s.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = UC_FONT
    r.font.name = FONT_NAME
    r.font.color.rgb = BLACK


def add_boundary(slide, left, top, w, h):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h)
    )
    s.fill.background()
    s.line.color.rgb = BLACK
    s.line.width = DASH_W
    s.line.dash_style = 4


def edge_pt(cx, cy, hw, hh, tx, ty):
    dx, dy = tx - cx, ty - cy
    if dx == 0 and dy == 0:
        return cx, cy
    hits = []
    if dx != 0:
        t = hw / abs(dx)
        py = cy + t * dy
        if abs(py - cy) <= hh + 0.01:
            hits.append((cx + hw * (1 if dx > 0 else -1), py))
    if dy != 0:
        t = hh / abs(dy)
        px = cx + t * dx
        if abs(px - cx) <= hw + 0.01:
            hits.append((px, cy + hh * (1 if dy > 0 else -1)))
    if not hits:
        return cx, cy
    return min(hits, key=lambda p: (p[0] - tx) ** 2 + (p[1] - ty) ** 2)


def draw_line(slide, x1, y1, x2, y2, width=LINE_W, dashed=False):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK
    c.line.width = width
    if dashed:
        c.line.dash_style = 4


def connect_actor_uc(slide, al, at, aw, ah, ul, ut, uw, uh):
    acx, acy = al + aw / 2, at + ah / 2
    ucx, ucy = ul + uw / 2, ut + uh / 2
    ax, ay = edge_pt(acx, acy, aw / 2, ah / 2, ucx, ucy)
    ux, uy = edge_pt(ucx, ucy, uw / 2, uh / 2, acx, acy)
    draw_line(slide, ax, ay, ux, uy)


def connect_uc_uc(slide, fl, ft, tl, tt, label):
    fcx, fcy = fl + UW / 2, ft + UH / 2
    tcx, tcy = tl + UW / 2, tt + UH / 2
    fx, fy = edge_pt(fcx, fcy, UW / 2, UH / 2, tcx, tcy)
    tx, ty = edge_pt(tcx, tcy, UW / 2, UH / 2, fcx, fcy)
    draw_line(slide, fx, fy, tx, ty, width=REL_W, dashed=True)
    mx, my = (fx + tx) / 2, (fy + ty) / 2
    w, h = 1.0, 0.22
    tb = slide.shapes.add_textbox(
        Inches(mx - w / 2), Inches(my - h / 2), Inches(w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = label
    r.font.size = REL_FONT
    r.font.name = FONT_NAME
    r.font.color.rgb = BLACK


# ─── LAYOUT ──────────────────────────────────────────────────────

# Left actors (primary stakeholders)
LEFT_ACTORS = [
    ("WORKER",           0.15, 1.2),
    ("SAFETY\nMANAGER",  0.15, 2.8),
    ("SITE\nSUPERVISOR", 0.15, 4.2),
    ("EMPLOYER\nADMIN",  0.15, 5.6),
]

# Right actors (external systems)
RIGHT_ACTORS = [
    ("WEATHER\nAPI",   SW - AW - 0.15, 0.7),
    ("SSO\nPROVIDER",  SW - AW - 0.15, 2.5),
]

ALL_ACTORS = {name: (x, y) for name, x, y in LEFT_ACTORS + RIGHT_ACTORS}

# System boundary
BND = (1.8, 0.35, 9.6, 6.65)  # left, top, width, height

# Use cases: id -> (label, x, y)
# Column 1 (x=2.2): Worker use cases — near left actors
# Column 2 (x=5.0): Manager/Supervisor use cases
# Column 3 (x=7.8): Admin + System use cases
C1, C2, C3 = 2.2, 5.0, 7.8

UCS = {
    # Worker column
    "UC1":  ("UC1: Request Heat\nRisk Prediction",  C1, 0.55),
    "UC2":  ("UC2: Receive\nSafety Alert",          C1, 1.4),
    "UC3":  ("UC3: Acknowledge\nAlert",             C1, 2.25),
    "UC7":  ("UC7: Onboard\nNew Worker",            C1, 3.1),
    "UC13": ("UC13: Switch Activity\nMid-Shift",    C1, 3.95),
    # Manager / Supervisor column
    "UC4":  ("UC4: Monitor Team\nRisk Status",      C2, 2.0),
    "UC5":  ("UC5: Configure Alert\nThresholds",    C2, 2.85),
    "UC6":  ("UC6: Generate\nCompliance Report",    C2, 3.7),
    "UC14": ("UC14: View Historical\nTrends",       C2, 4.55),
    "UC8":  ("UC8: Receive\nEscalation Alert",      C2, 5.4),
    # Admin + System column
    "UC9":  ("UC9: Provision\nOrganization",        C3, 4.8),
    "UC10": ("UC10: Export\nAudit Trail",            C3, 5.65),
    "UC11": ("UC11: Update\nEnvironmental Data",    C3, 0.7),
    "UC12": ("UC12: Authenticate\nvia SSO",         C3, 2.5),
}

# Actor → use-case connections
CONNECTIONS = [
    ("WORKER",           ["UC1", "UC2", "UC3", "UC7", "UC13"]),
    ("SAFETY\nMANAGER",  ["UC4", "UC5", "UC6", "UC14"]),
    ("SITE\nSUPERVISOR", ["UC8"]),
    ("EMPLOYER\nADMIN",  ["UC9", "UC10"]),
    ("WEATHER\nAPI",     ["UC11"]),
    ("SSO\nPROVIDER",    ["UC12"]),
]

# Include / Extend relationships (from → to)
RELATIONSHIPS = [
    ("UC1",  "UC11", "<<include>>"),
    ("UC1",  "UC12", "<<include>>"),
    ("UC7",  "UC1",  "<<include>>"),
    ("UC13", "UC1",  "<<extend>>"),
    ("UC3",  "UC2",  "<<extend>>"),
    ("UC8",  "UC3",  "<<extend>>"),
]


# ─── SLIDE BUILDER ───────────────────────────────────────────────

def build(prs, title, show_rels=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.01), Inches(12.8), Inches(0.3))
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.size = TITLE_FONT
    r.font.name = FONT_NAME
    r.font.color.rgb = BLACK
    r.font.bold = True

    # System boundary
    bl, bt, bw, bh = BND
    add_boundary(slide, bl, bt, bw, bh)

    # Boundary label
    lb = slide.shapes.add_textbox(
        Inches(bl + 0.15), Inches(bt + 0.05), Inches(3), Inches(0.25)
    )
    tf = lb.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "THE SYSTEM"
    r.font.size = BND_LABEL_FONT
    r.font.name = FONT_NAME
    r.font.color.rgb = BLACK
    r.font.bold = True

    # Actors
    for name, (ax, ay) in ALL_ACTORS.items():
        add_actor(slide, ax, ay, name)

    # Use cases
    for uid, (label, ux, uy) in UCS.items():
        add_usecase(slide, ux, uy, label)

    # Actor → UC lines
    for actor_name, uc_ids in CONNECTIONS:
        ax, ay = ALL_ACTORS[actor_name]
        for uid in uc_ids:
            _, ux, uy = UCS[uid]
            connect_actor_uc(slide, ax, ay, AW, AH, ux, uy, UW, UH)

    # Include/Extend (detailed only)
    if show_rels:
        for from_id, to_id, rel in RELATIONSHIPS:
            _, fx, fy = UCS[from_id]
            _, tx, ty = UCS[to_id]
            connect_uc_uc(slide, fx, fy, tx, ty, rel)


# ─── MAIN ────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    build(prs, "Use Case Diagram", show_rels=False)
    build(prs, "Use Case Diagram - Detailed", show_rels=True)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(script_dir, "THG_Use_Case_Diagram.pptx")
    prs.save(out)

    print(f"Saved: {out}")
    print(f"Actors: {len(ALL_ACTORS)}")
    print(f"Use Cases: {len(UCS)}")
    print(f"Relationships: {len(RELATIONSHIPS)}")
