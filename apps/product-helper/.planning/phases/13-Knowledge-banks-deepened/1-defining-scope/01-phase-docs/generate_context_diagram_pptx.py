#!/usr/bin/env python3
"""
Generate Cornell-standard Context Diagram for Team Heat Guard (THG)
Following CESYS521 formatting rules exactly:
  - Black and white only (no color)
  - Square corners (never rounded)
  - All boxes the SAME size (including system box)
  - Names CAPITALIZED
  - Same font and font size
  - Dashed border = system boundary
  - 8-20 external elements
  - Labels placed near each box (not floating midpoint)
  - Two labels per line: "from" near element, "to" near system
  - Balanced layout: ~4 boxes per side
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math

# --- Constants ---
FONT_SIZE = Pt(10)
LABEL_FONT = Pt(6.5)
FONT_NAME = "Arial"
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LINE_WIDTH = Pt(1.2)
DASH_WIDTH = Pt(2.0)

# Slide dims in inches
SW = 13.333
SH = 7.5
CX = SW / 2
CY = SH / 2

# ALL boxes same size (Cornell standard)
BW = 1.45
BH = 0.6


def add_box(slide, left, top, text, bold=False, font_size=FONT_SIZE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(BW), Inches(BH)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = LINE_WIDTH

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text.upper()
    run.font.size = font_size
    run.font.name = FONT_NAME
    run.font.color.rgb = BLACK
    run.font.bold = bold
    return shape


def add_dashed_rect(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.background()
    shape.line.color.rgb = BLACK
    shape.line.width = DASH_WIDTH
    shape.line.dash_style = 4
    return shape


def get_edge_point(cx, cy, hw, hh, tx, ty):
    """Get point on box edge closest to target (tx, ty)."""
    dx = tx - cx
    dy = ty - cy
    if dx == 0 and dy == 0:
        return cx, cy

    candidates = []
    if dx != 0:
        t = hw / abs(dx)
        py = cy + t * dy
        if abs(py - cy) <= hh + 0.01:
            candidates.append((cx + hw * (1 if dx > 0 else -1), py))
    if dy != 0:
        t = hh / abs(dy)
        px = cx + t * dx
        if abs(px - cx) <= hw + 0.01:
            candidates.append((px, cy + hh * (1 if dy > 0 else -1)))

    if not candidates:
        return cx, cy
    return min(candidates, key=lambda p: (p[0] - tx)**2 + (p[1] - ty)**2)


def add_line(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = BLACK
    conn.line.width = LINE_WIDTH
    return conn


def connect_orthogonal(slide, elem_left, elem_top, sys_left, sys_top, side):
    """Draw right-angle (L-shaped) connector from element to system.

    side: 'top', 'bottom', 'left', 'right' — which side of the diagram the element is on.

    Routing:
      - top elements:    exit bottom of element → go down vertically → turn right/left → enter top of system
      - bottom elements: exit top of element → go up vertically → turn right/left → enter bottom of system
      - left elements:   exit right of element → go right horizontally → turn up/down → enter left of system
      - right elements:  exit left of element → go left horizontally → turn up/down → enter right of system
    """
    ecx = elem_left + BW / 2
    ecy = elem_top + BH / 2
    scx = sys_left + BW / 2
    scy = sys_top + BH / 2

    if side == 'top':
        # Exit bottom of element, enter top of system
        x1 = ecx
        y1 = elem_top + BH  # bottom edge of element
        x3 = scx
        y3 = sys_top  # top edge of system (really top of dashed boundary area)
        # Corner point: go down to the system's y-level approach, then across
        corner_y = y3  # turn at the system's top edge level
        # Segment 1: vertical down from element
        add_line(slide, x1, y1, x1, corner_y)
        # Segment 2: horizontal to system
        add_line(slide, x1, corner_y, x3, corner_y)

    elif side == 'bottom':
        # Exit top of element, enter bottom of system
        x1 = ecx
        y1 = elem_top  # top edge of element
        x3 = scx
        y3 = sys_top + BH  # bottom edge of system
        corner_y = y3
        add_line(slide, x1, y1, x1, corner_y)
        add_line(slide, x1, corner_y, x3, corner_y)

    elif side == 'left':
        # Exit right of element, enter left of system
        x1 = elem_left + BW  # right edge of element
        y1 = ecy
        x3 = sys_left  # left edge of system
        y3 = scy
        corner_x = x3
        add_line(slide, x1, y1, corner_x, y1)
        add_line(slide, corner_x, y1, corner_x, y3)

    elif side == 'right':
        # Exit left of element, enter right of system
        x1 = elem_left  # left edge of element
        y1 = ecy
        x3 = sys_left + BW  # right edge of system
        y3 = scy
        corner_x = x3
        add_line(slide, x1, y1, corner_x, y1)
        add_line(slide, corner_x, y1, corner_x, y3)


def add_label(slide, x, y, text, font_size=LABEL_FONT, max_w=1.8):
    """Add interaction label near a box. x,y = center of label."""
    w = max_w
    h = 0.55
    left = x - w / 2
    top = y - h / 2
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.name = FONT_NAME
    run.font.color.rgb = BLACK
    return txBox


# ============================================================
# LAYOUT: 4 sides, 4 boxes each = 16 elements
# ============================================================

# System box at center (same size as all other boxes)
sys_left = CX - BW / 2
sys_top = CY - BH / 2

# Dashed boundary
BPAD = 0.35
bound_left = sys_left - BPAD
bound_top = sys_top - BPAD
bound_w = BW + 2 * BPAD
bound_h = BH + 2 * BPAD

# TOP ROW — 4 boxes evenly spaced
top_y = 0.5
top_xs = [1.5, 4.2, 7.0, 10.0]
top_elements = [
    ("Workers",              top_xs[0], top_y),
    ("Weather\nAPI",         top_xs[1], top_y),
    ("Outdoor\nEnvironment", top_xs[2], top_y),
    ("Indoor\nEnvironment",  top_xs[3], top_y),
]

# LEFT COLUMN — 4 boxes evenly spaced vertically
left_x = 0.3
left_ys = [1.8, 3.0, 4.2, 5.4]
left_elements = [
    ("Work\nActivities",   left_x, left_ys[0]),
    ("Clothing /\nPPE",    left_x, left_ys[1]),
    ("Employer\nUsers",    left_x, left_ys[2]),
    ("Mobile\nDevices",    left_x, left_ys[3]),
]

# RIGHT COLUMN — 4 boxes evenly spaced vertically
right_x = SW - BW - 0.3
right_ys = [1.8, 3.0, 4.2, 5.4]
right_elements = [
    ("OSHA",              right_x, right_ys[0]),
    ("State\nStandards",  right_x, right_ys[1]),
    ("NIOSH /\nACGIH",    right_x, right_ys[2]),
    ("Workers\nComp",     right_x, right_ys[3]),
]

# BOTTOM ROW — 4 boxes evenly spaced
bot_y = 6.4
bot_xs = [1.5, 4.2, 7.0, 10.0]
bot_elements = [
    ("SSO /\nIdentity",      bot_xs[0], bot_y),
    ("Web\nBrowser",         bot_xs[1], bot_y),
    ("Audit\nSystems",       bot_xs[2], bot_y),
    ("Emergency\nServices",  bot_xs[3], bot_y),
]

# Tag each element with its side for orthogonal routing
all_elements = (
    [(name, x, y, 'top') for name, x, y in top_elements] +
    [(name, x, y, 'left') for name, x, y in left_elements] +
    [(name, x, y, 'right') for name, x, y in right_elements] +
    [(name, x, y, 'bottom') for name, x, y in bot_elements]
)

# ============================================================
# INTERACTION LABELS
# Format: element_name -> (from_element_label, from_system_label, label_position)
# label_position: "above", "below", "left", "right" — relative to the element
# ============================================================

interactions = {
    # TOP elements — labels below element box
    "Workers":              ("biometrics, activity data",               "risk alerts, predictions"),
    "Weather\nAPI":         ("temp, humidity, wind, radiation",         ""),
    "Outdoor\nEnvironment": ("solar load, ambient conditions",          ""),
    "Indoor\nEnvironment":  ("enclosure temp, ventilation data",        ""),

    # LEFT elements — labels right of element box
    "Work\nActivities":     ("metabolic rate, task intensity",          ""),
    "Clothing /\nPPE":      ("insulation factor, vapor resistance",     ""),
    "Employer\nUsers":      ("configure thresholds,\nprovision orgs, acknowledge alerts",
                             "team dashboards,\nescalation alerts, compliance reports"),
    "Mobile\nDevices":      ("",                                        "mobile app delivery"),

    # RIGHT elements — labels left of element box
    "OSHA":                 ("compliance requirements",                 "compliance evidence"),
    "State\nStandards":     ("state-specific thresholds",               ""),
    "NIOSH /\nACGIH":       ("accuracy benchmarks, TLV guidelines",     ""),
    "Workers\nComp":        ("",                                        "risk evidence"),

    # BOTTOM elements — labels above element box
    "SSO /\nIdentity":      ("",                                        "authentication"),
    "Web\nBrowser":         ("",                                        "web dashboard"),
    "Audit\nSystems":       ("",                                        "audit logs"),
    "Emergency\nServices":  ("",                                        "critical alerts"),
}


def place_label_near_element(slide, elem_left, elem_top, sys_left, sys_top, text):
    """Place label between element and system, near the element side."""
    ecx = elem_left + BW / 2
    ecy = elem_top + BH / 2
    scx = sys_left + BW / 2
    scy = sys_top + BH / 2

    # Place at 25% of the way from element to system
    lx = ecx + (scx - ecx) * 0.25
    ly = ecy + (scy - ecy) * 0.25
    add_label(slide, lx, ly, text)


def place_label_near_system(slide, elem_left, elem_top, sys_left, sys_top, text):
    """Place label between element and system, near the system side."""
    ecx = elem_left + BW / 2
    ecy = elem_top + BH / 2
    scx = sys_left + BW / 2
    scy = sys_top + BH / 2

    # Place at 75% of the way from element to system
    lx = ecx + (scx - ecx) * 0.72
    ly = ecy + (scy - ecy) * 0.72
    add_label(slide, lx, ly, text)


def build_slide(prs, title, detailed=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.02), Inches(12.5), Inches(0.42))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.name = FONT_NAME
    run.font.color.rgb = BLACK
    run.font.bold = True

    # Dashed system boundary
    add_dashed_rect(slide, bound_left, bound_top, bound_w, bound_h)

    # System box (same size as all others)
    add_box(slide, sys_left, sys_top, "The System", bold=True)

    # Element boxes + connections + labels
    for name, left, top, side in all_elements:
        add_box(slide, left, top, name)
        connect_orthogonal(slide, left, top, sys_left, sys_top, side)

        if detailed and name in interactions:
            from_lbl, to_lbl = interactions[name]
            if from_lbl:
                place_label_near_element(slide, left, top, sys_left, sys_top, text=from_lbl)
            if to_lbl:
                place_label_near_system(slide, left, top, sys_left, sys_top, text=to_lbl)

    return slide


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Basic (boxes + lines only)
build_slide(prs, "Context Diagram", detailed=False)

# Slide 2: Detailed (boxes + lines + interaction labels)
build_slide(prs, "Context Diagram - Detailed", detailed=True)

out_path = "/Users/davidancor/Projects/peak-genics-heatgaurd/system-design/diagrams/THG_Context_Diagram.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Elements: {len(all_elements)}")
print("Layout: 4 top, 4 left, 4 right, 4 bottom")
