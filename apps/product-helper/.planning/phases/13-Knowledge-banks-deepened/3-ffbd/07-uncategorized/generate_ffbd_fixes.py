#!/usr/bin/env python3
"""
Generate corrected FFBD slides for F.4 and F.5.

Fixes:
  F.4 — Remove duplicate AND-gate-to-F.4.14 connector (Path A at y=15.5)
  F.5 — Bridge F.5.9→F.5.10 gap; route F.5.7g→F.5.10

Creates THG_FFBD_fixes.pptx with 2 replacement slides.
User plugs them into the main THG_FFBD.pptx manually.

Usage:
  pip install python-pptx lxml
  python generate_ffbd_fixes.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn

# ── Constants ──────────────────────────────────────────────
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
PINK  = RGBColor(0xFF, 0xEB, 0xEB)   # F.4.5 highlight
CREAM = RGBColor(0xFF, 0xF9, 0xE0)   # decision diamonds
LINE_W = Pt(2)                        # 25400 EMU


# ── Helpers ────────────────────────────────────────────────

def add_func_box(slide, x, y, func_id, name_lines, fill=WHITE,
                 w=2.8, id_h=0.4, body_h=1.1):
    """Two-part function box: ID header + name body."""
    # Header
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(id_h))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = fill
    hdr.line.color.rgb = BLACK; hdr.line.width = LINE_W
    tf = hdr.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = func_id
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = BLACK

    # Body
    body = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y + id_h), Inches(w), Inches(body_h))
    body.fill.solid(); body.fill.fore_color.rgb = fill
    body.line.color.rgb = BLACK; body.line.width = LINE_W
    tf2 = body.text_frame; tf2.word_wrap = True
    for i, line in enumerate(name_lines):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = line
        r2.font.size = Pt(12); r2.font.color.rgb = BLACK


def add_gate(slide, x, y, label="AND", size=0.9):
    """AND/OR gate (ellipse)."""
    g = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    g.fill.solid(); g.fill.fore_color.rgb = WHITE
    g.line.color.rgb = BLACK; g.line.width = LINE_W
    tf = g.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = BLACK


def add_diamond(slide, x, y, text_lines, w=2.2, h=1.5):
    """Decision diamond."""
    d = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
    d.fill.solid(); d.fill.fore_color.rgb = CREAM
    d.line.color.rgb = BLACK; d.line.width = LINE_W
    tf = d.text_frame; tf.word_wrap = True
    for i, line in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(10); r.font.color.rgb = BLACK


def conn(slide, x1, y1, x2, y2, arrow=True):
    """Straight connector line built from raw XML (matches original PPTX).

    Uses p:cxnSp with prstGeom="line", black 2pt stroke, and optional
    triangle arrowhead at the tail end.  Handles flipH / flipV so the
    line always draws from (x1,y1) → (x2,y2).
    """
    # Convert to EMU
    ex1, ey1 = int(Inches(x1)), int(Inches(y1))
    ex2, ey2 = int(Inches(x2)), int(Inches(y2))

    # Bounding box — offset is always top-left, extent is always positive
    off_x = min(ex1, ex2)
    off_y = min(ey1, ey2)
    cx = abs(ex2 - ex1)
    cy = abs(ey2 - ey1)

    # Determine flip flags so the line goes from (x1,y1) → (x2,y2)
    flip_h = ex2 < ex1
    flip_v = ey2 < ey1

    # Build XML matching the original THG_FFBD.pptx connector structure
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    cxnSp = etree.SubElement(
        slide.shapes._spTree, qn('p:cxnSp'))

    # nvCxnSpPr
    nvCxnSpPr = etree.SubElement(cxnSp, qn('p:nvCxnSpPr'))
    cNvPr = etree.SubElement(nvCxnSpPr, qn('p:cNvPr'))
    cNvPr.set('id', '0')
    cNvPr.set('name', '')
    etree.SubElement(nvCxnSpPr, qn('p:cNvCxnSpPr'))
    etree.SubElement(nvCxnSpPr, qn('p:nvPr'))

    # spPr
    spPr = etree.SubElement(cxnSp, qn('p:spPr'))

    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    if flip_h:
        xfrm.set('flipH', '1')
    if flip_v:
        xfrm.set('flipV', '1')
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(off_x))
    off.set('y', str(off_y))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(cx))
    ext.set('cy', str(cy))

    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', 'line')
    etree.SubElement(prstGeom, qn('a:avLst'))

    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', '25400')  # 2 pt
    sf = etree.SubElement(ln, qn('a:solidFill'))
    clr = etree.SubElement(sf, qn('a:srgbClr'))
    clr.set('val', '000000')

    if arrow:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')

    # style block (matches original)
    style = etree.SubElement(cxnSp, qn('p:style'))
    for tag, idx in [('a:lnRef', '2'), ('a:fillRef', '0'),
                     ('a:effectRef', '1'), ('a:fontRef', None)]:
        ref = etree.SubElement(style, qn(tag))
        if idx is not None:
            ref.set('idx', idx)
        else:
            ref.set('idx', 'minor')
        sc = etree.SubElement(ref, qn('a:schemeClr'))
        sc.set('val', 'accent1' if 'font' not in tag else 'tx1')


def label(slide, x, y, text, w=5.0, h=0.5, size=12, italic=True,
          bold=False, align=PP_ALIGN.LEFT):
    """Text annotation."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.italic = italic
    r.font.bold = bold; r.font.color.rgb = BLACK


def title_block(slide, title, subtitle):
    """Slide title + subtitle."""
    label(slide, 1, 0.5, title, w=28, size=34, italic=False, bold=True,
          align=PP_ALIGN.CENTER)
    label(slide, 1, 1.4, subtitle, w=28, size=16, italic=False,
          align=PP_ALIGN.CENTER)


# ── F.4  Predict Heat Strain ──────────────────────────────

def build_f4(slide):
    title_block(slide,
        "F.4 Predict Heat Strain \u2605 CORE VALUE",
        "16 functions \u2013 dual-model (PHS+WBGT), AND gate, "
        "ML decision \u2013 MAE \u2264 0.3\u00b0C target")

    # ── Top row: preparation (F.4.1 → F.4.4) ──
    boxes_top = [
        (0.6,  "F.4.1",  ["Retrieve", "Worker Profile"]),
        (4.1,  "F.4.2",  ["Look Up", "Metabolic Rate", "(ISO 8996)"]),
        (7.6,  "F.4.2a", ["Apply", "Acclimatization", "Factor"]),
        (11.1, "F.4.3",  ["Apply Clothing", "Factor", "(ISO 7933)"]),
        (14.6, "F.4.3a", ["Retrieve", "Environmental", "Snapshot"]),
        (18.1, "F.4.3b", ["Compute", "WBGT Index"]),
        (21.6, "F.4.3c", ["Select Model", "Path"]),
        (25.1, "F.4.4",  ["Assemble", "Algorithm Inputs"]),
    ]
    for bx, fid, lines in boxes_top:
        add_func_box(slide, bx, 4.25, fid, lines)

    # Top row arrows
    xs_top = [0.6, 4.1, 7.6, 11.1, 14.6, 18.1, 21.6, 25.1]
    W = 2.8
    for i in range(len(xs_top) - 1):
        conn(slide, xs_top[i] + W, 5.0, xs_top[i+1], 5.0)

    # ── F.4.4 drops to horizontal backbone ──
    conn(slide, 26.5, 5.75, 26.5, 9.0, arrow=False)  # F.4.4 center-bottom → y=9
    conn(slide, 4.5, 9.0, 26.5, 9.0, arrow=False)    # backbone horizontal

    # ── Drop to split AND gate ──
    conn(slide, 4.5, 9.0, 4.5, 10.0, arrow=False)    # backbone → AND entry level
    conn(slide, 3.45, 10.0, 4.5, 10.0, arrow=False)  # → AND right edge

    add_gate(slide, 2.55, 9.55, "AND")            # split AND

    # ── AND splits into PHS (up) and WBGT (down) ──
    conn(slide, 3.45, 10.0, 3.75, 10.0, arrow=False)  # AND right → branch point
    conn(slide, 3.75, 8.5, 3.75, 10.0, arrow=False)   # branch ↑ to PHS level
    conn(slide, 3.75, 8.5, 5.6, 8.5)              # → F.4.5
    conn(slide, 3.75, 10.0, 3.75, 12.5, arrow=False)  # branch ↓ to WBGT level
    conn(slide, 3.75, 12.5, 5.6, 12.5)            # → F.4.10

    # ── PHS path ──
    add_func_box(slide, 5.6, 7.75, "F.4.5",
                 ["Execute PHS", "Algorithm", "(ISO 7933)"], fill=PINK)
    conn(slide, 8.4, 8.5, 10.55, 8.5)            # F.4.5 → AND

    add_gate(slide, 10.55, 8.05, "AND")           # Trec/Dlim split AND
    conn(slide, 11.45, 8.5, 11.7, 8.5, arrow=False)  # AND right edge → branch
    conn(slide, 11.7, 7.5, 11.7, 8.5, arrow=False)   # branch ↑ Trec
    conn(slide, 11.7, 7.5, 13.1, 7.5)            # → F.4.6
    conn(slide, 11.7, 8.5, 11.7, 9.5, arrow=False)   # branch ↓ Dlim
    conn(slide, 11.7, 9.5, 13.1, 9.5)            # → F.4.7

    add_func_box(slide, 13.1, 6.75, "F.4.6", ["Calculate", "Trec"])
    add_func_box(slide, 13.1, 8.75, "F.4.7", ["Calculate", "Dlim"])

    # Trec/Dlim → join AND
    conn(slide, 15.9, 7.5, 17.0, 7.5, arrow=False)   # F.4.6 → join
    conn(slide, 17.0, 7.5, 17.0, 8.5, arrow=False)   # ↓ to AND center
    conn(slide, 15.9, 9.5, 17.0, 9.5, arrow=False)   # F.4.7 → join
    conn(slide, 17.0, 8.5, 17.0, 9.5, arrow=False)   # ↑ to AND center
    conn(slide, 17.0, 8.5, 17.25, 8.5, arrow=False)  # → AND left edge

    add_gate(slide, 17.25, 8.05, "AND")           # join AND

    # PHS join AND → convergence AND
    conn(slide, 18.15, 8.5, 21.3, 8.5, arrow=False)  # AND right → x=21.3
    conn(slide, 21.3, 8.5, 21.3, 10.0, arrow=False)  # ↓ to convergence level

    # ── WBGT path ──
    add_func_box(slide, 5.6, 11.75, "F.4.10", ["Compute", "WBGT Screening"])
    add_func_box(slide, 9.6, 11.75, "F.4.11", ["Apply OSHA", "Action Limits"])
    add_func_box(slide, 13.6, 11.75, "F.4.12", ["Classify", "WBGT Risk", "Level"])
    add_func_box(slide, 17.6, 11.75, "F.4.13", ["Generate", "WBGT", "Recommendation"])

    conn(slide, 8.4, 12.5, 9.6, 12.5)            # F.4.10 → F.4.11
    conn(slide, 12.4, 12.5, 13.6, 12.5)          # F.4.11 → F.4.12
    conn(slide, 16.4, 12.5, 17.6, 12.5)          # F.4.12 → F.4.13
    conn(slide, 20.4, 12.5, 21.3, 12.5)          # F.4.13 → convergence

    # WBGT → convergence AND (vertical from y=10 to y=12.5)
    conn(slide, 21.3, 10.0, 21.3, 12.5, arrow=False)

    # ── Convergence AND ──
    conn(slide, 21.3, 10.0, 21.55, 10.0, arrow=False)  # → AND left edge
    add_gate(slide, 21.55, 9.55, "AND")

    # ── FIX: single exit path from convergence AND to F.4.14 ──
    # Path B only: AND → right → down to y=17.5 → left → up → into F.4.14
    conn(slide, 22.45, 10.0, 22.8, 10.0, arrow=False)   # AND right edge → x=22.8
    conn(slide, 22.8, 10.0, 22.8, 17.5, arrow=False)    # ↓ to y=17.5
    conn(slide, 3.0, 17.5, 22.8, 17.5, arrow=False)     # ← to x=3.0
    conn(slide, 3.0, 15.5, 3.0, 17.5, arrow=False)      # ↑ to y=15.5
    conn(slide, 3.0, 15.5, 3.6, 15.5)             # → into F.4.14

    # (Path A removed — no (22.5,10→15.5) drop, no (6.9,15.5) horizontal)

    # ── Bottom row ──
    add_func_box(slide, 3.6, 14.75, "F.4.14",
                 ["Reconcile", "PHS + WBGT", "Results"])
    add_func_box(slide, 7.6, 14.75, "F.4.15",
                 ["Check", "Compliance", "Trigger"])
    add_diamond(slide, 11.9, 14.75, ["ML Model", "Available?"])
    add_func_box(slide, 11.6, 16.75, "F.4.16", ["Apply ML", "Correction"])
    add_func_box(slide, 16.6, 14.75, "F.4.8",
                 ["Store", "Prediction", "Result"])
    add_func_box(slide, 20.6, 14.75, "F.4.9", ["Log to", "Audit Trail"])

    conn(slide, 6.4, 15.5, 7.6, 15.5)            # F.4.14 → F.4.15
    conn(slide, 10.4, 15.5, 11.9, 15.5)          # F.4.15 → ML diamond
    conn(slide, 14.1, 15.5, 16.6, 15.5)          # ML "No" → F.4.8
    conn(slide, 13.0, 16.25, 13.0, 17.5, arrow=False)  # ML "Yes" ↓ F.4.16
    conn(slide, 14.4, 17.5, 18.0, 17.5, arrow=False)  # F.4.16 → (below F.4.8)
    conn(slide, 18.0, 16.25, 18.0, 17.5, arrow=False)  # ↑ into F.4.8 bottom
    conn(slide, 19.4, 15.5, 20.6, 15.5)          # F.4.8 → F.4.9

    # ── Labels ──
    label(slide, 3.95, 7.2, "PHS deep analysis")
    label(slide, 3.95, 13.3, "WBGT screening")
    label(slide, 11.9, 6.9, "Core body temperature")
    label(slide, 11.9, 10.3, "Safe work duration limit")
    label(slide, 14.2, 15.0, "No \u2013 skip", w=3.0)
    label(slide, 13.5, 16.5, "Yes \u2013 apply\nML correction", w=4.0)


# ── F.5  Classify & Alert ─────────────────────────────────

def build_f5(slide):
    title_block(slide,
        "F.5 Classify & Alert",
        "17 functions \u2013 Decision gates + intervention effectiveness chain")

    # ── Main flow ──
    add_func_box(slide, 0.6, 6.25, "F.5.1",
                 ["Apply Risk", "Thresholds"])
    add_diamond(slide, 3.4, 6.25, ["Risk Level", "Changed?"])

    # "No change" branch
    add_func_box(slide, 6.1, 3.25, "F.5.3",
                 ["Update", "Dashboard Only"])

    # "Changed" branch
    add_func_box(slide, 6.1, 6.25, "F.5.4",
                 ["Generate Alert", "+ Guidance"])
    add_func_box(slide, 9.1, 6.25, "F.5.4a",
                 ["Generate", "Intervention", "Prescription"])
    add_func_box(slide, 12.1, 6.25, "F.5.5",
                 ["Push to", "Worker"])
    add_diamond(slide, 14.9, 6.25, ["Acknowledged", "in time?"])

    # "Yes" branch
    add_func_box(slide, 17.6, 3.25, "F.5.7",
                 ["Log", "Acknowledgment"])

    # "No" branch — escalation
    add_func_box(slide, 17.6, 9.75, "F.5.8",
                 ["Escalate to", "Supervisor"])
    add_func_box(slide, 20.6, 9.75, "F.5.9",
                 ["Log", "Escalation"])

    # Terminal function
    add_func_box(slide, 25.1, 6.25, "F.5.10",
                 ["Update", "Manager", "Dashboard"])

    # ── Intervention Effectiveness Chain ──
    chain = [
        (1.6,  "F.5.7a", ["Record", "Pre-Intervention", "Trec"]),
        (5.1,  "F.5.7b", ["Worker Completes", "Intervention"]),
        (8.6,  "F.5.7c", ["Record", "Completion GPS"]),
        (12.1, "F.5.7d", ["Trigger", "Re-Prediction"]),
        (15.6, "F.5.7e", ["Record", "Post-Intervention", "Trec"]),
        (19.1, "F.5.7f", ["Compute", "Effectiveness", "Score"]),
        (22.6, "F.5.7g", ["Feed ML", "Training Pipeline"]),
    ]
    for bx, fid, lines in chain:
        add_func_box(slide, bx, 13.25, fid, lines)

    # Chain arrows (F.5.7a → F.5.7g)
    chain_xs = [c[0] for c in chain]
    for i in range(len(chain_xs) - 1):
        conn(slide, chain_xs[i] + 2.8, 14.0, chain_xs[i+1], 14.0)

    # ── Main flow connectors ──

    # F.5.1 → diamond (right edge of F.5.1 to left edge of diamond)
    conn(slide, 3.4, 7.0, 3.4, 7.0, arrow=False)  # structural anchor
    # Diamond "No change" → F.5.3
    conn(slide, 4.5, 6.25, 4.5, 4.0, arrow=False) # diamond top → y=4.0 (vertical routing)
    conn(slide, 4.5, 4.0, 6.1, 4.0)               # → F.5.3

    # Diamond "Changed" → F.5.4
    conn(slide, 5.6, 7.0, 6.1, 7.0)               # diamond right → F.5.4

    # F.5.4 → F.5.4a → F.5.5
    conn(slide, 8.9, 7.0, 9.1, 7.0)
    conn(slide, 11.9, 7.0, 12.1, 7.0)

    # F.5.5 → ack diamond
    conn(slide, 14.9, 7.0, 14.9, 7.0, arrow=False)  # (zero-size anchor)

    # Diamond "Yes" → F.5.7
    conn(slide, 16.0, 6.25, 16.0, 4.0, arrow=False)  # diamond top → y=4.0
    conn(slide, 16.0, 4.0, 17.6, 4.0)             # → F.5.7

    # Diamond "No" → F.5.8
    conn(slide, 16.0, 7.75, 16.0, 10.5, arrow=False)  # diamond bottom → y=10.5
    conn(slide, 16.0, 10.5, 17.6, 10.5)           # → F.5.8

    # F.5.8 → F.5.9
    conn(slide, 20.4, 10.5, 20.6, 10.5)

    # ── FIX 1: F.5.9 → F.5.10 (bridge the gap) ──
    conn(slide, 23.4, 10.5, 24.7, 10.5, arrow=False)  # F.5.9 right → x=24.7
    conn(slide, 24.7, 7.0, 24.7, 10.5, arrow=False)   # ↑ to F.5.10 level
    conn(slide, 24.7, 7.0, 25.1, 7.0)             # → F.5.10 left edge (BRIDGED)

    # F.5.3 "No change" → F.5.10 (across top)
    conn(slide, 8.9, 4.0, 25.0, 4.0, arrow=False)    # long horizontal at y=4
    conn(slide, 25.0, 4.0, 25.0, 6.25, arrow=False)   # ↓ into F.5.10 top

    # F.5.7 "Acknowledged" → F.5.10 (via same top horizontal)
    # (already covered by the 8.9→25.0 line above, since F.5.7 at y=4.0
    #  feeds into the same horizontal)

    # F.5.7 ack side-path → effectiveness chain
    conn(slide, 19.0, 4.75, 19.0, 12.5, arrow=False)  # F.5.7 bottom → chain level
    conn(slide, 1.5, 12.5, 19.0, 12.5, arrow=False)   # horizontal at chain level
    conn(slide, 1.5, 12.5, 1.5, 14.0, arrow=False)    # ↓ into chain
    conn(slide, 1.5, 14.0, 1.6, 14.0)             # → F.5.7a

    # ── FIX 2: F.5.7g → F.5.10 ──
    conn(slide, 25.4, 14.0, 26.5, 14.0, arrow=False)  # F.5.7g right → x=26.5
    conn(slide, 26.5, 7.75, 26.5, 14.0, arrow=False)  # ↑ to F.5.10 bottom

    # F.5.10 exit arrow →
    conn(slide, 27.9, 7.0, 28.2, 7.0)             # F.5.10 right → exit

    # ── Labels ──
    label(slide, 5.0, 3.4, "No change", w=3.0)
    label(slide, 5.7, 6.5, "Changed", w=3.0)
    label(slide, 16.5, 3.4, "Yes \u2013 acknowledged", w=5.0)
    label(slide, 16.5, 11.0,
          "No \u2013 escalate\n(AMBER: 5 min, RED: 2 min)", w=6.0)
    label(slide, 1.5, 12.0, "Intervention Effectiveness Chain",
          w=8.0, bold=True, italic=False)


# ── Main ───────────────────────────────────────────────────

def main():
    # Read original for slide dimensions
    orig = Presentation("THG_FFBD.pptx")
    sw = orig.slide_width
    sh = orig.slide_height
    print(f"Slide size: {sw/914400:.0f} x {sh/914400:.0f} inches")

    prs = Presentation()
    prs.slide_width = sw
    prs.slide_height = sh

    blank = prs.slide_layouts[6]  # blank layout

    # Slide 1: corrected F.4
    s4 = prs.slides.add_slide(blank)
    build_f4(s4)
    print("Built F.4 slide (duplicate AND→F.4.14 path removed)")

    # Slide 2: corrected F.5
    s5 = prs.slides.add_slide(blank)
    build_f5(s5)
    print("Built F.5 slide (F.5.9→F.5.10 bridged, F.5.7g→F.5.10 routed)")

    out = "THG_FFBD_fixes.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")
    print("Plug slide 1 into main PPTX as replacement for F.4 (slide 5)")
    print("Plug slide 2 into main PPTX as replacement for F.5 (slide 6)")


if __name__ == "__main__":
    main()
