#!/usr/bin/env python3
"""Create all THG FFBD slides – Cornell CESYS523 standard.  (V3)

Changes from V2:
- Font: body 24pt, header 20pt, gates 22pt, labels 20pt italic, title 28pt
- Arrows: 2.5pt weight, type-2 elbow connectors for simple paths, chained
  straight segments for complex routing, filled triangle arrowheads
- Block borders: 2.0pt (arrows one setting thicker per CESYS523)
- Titles: left-aligned (CESYS523 guideline)
- Decision diamonds replaced with OR gate pairs (FFBD convention)
- Reference blocks: bracket style (CESYS523 guideline)
- F.4: Sport-specific OR branch + F.4a reference block
- F.4a: New slide - Pre-Event Risk Stratification (6 blocks)
- F.6: Blocks renamed to functional names (not structural/display)
- F.7: AND gate + F.7.2 Advance Acclimatization Protocol

Slides (9 total):
  1: Top-Level (F.1-F.7) with IT loops
  2: F.1 Onboard Organization (6 blocks)
  3: F.2 Onboard Workers (7 blocks)
  4: F.3 Collect Environmental Data (OR gates)
  5: F.4 Predict Heat Strain (AND + OR gates, sport path, F.4a ref)
  6: F.5 Classify & Alert (OR gates, intervention chain)
  7: F.6 Monitor & Report (AND + OR gates)
  8: F.7 Maintain Compliance (AND + OR gates, acclimatization cron)
  9: F.4a Pre-Event Risk Stratification (6 blocks)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# === FILE PATHS ===
TEMPLATE = "/Users/davidancor/Documents/MBA/System Design - eCornell/3 - Exploring Your System's Architecture/CESYS523_FFBD_Template - MASTER.pptx"
OUTPUT = "/Users/davidancor/Documents/MBA/System Design - eCornell/thg-system-design-example/diagrams/THG_FFBD_v3.pptx"

# === CONSTANTS (30x20 inch slide) ===
LINE_W = Pt(2.5)       # Arrow weight (one setting thicker than borders)
BORDER_W = Pt(2.0)     # Block border weight
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
CORE_FILL = RGBColor(255, 235, 235)    # Light red - core algorithm highlight

# Block dimensions
BW = 2.8;   BBH = 1.4;  BHH = 0.50    # Functional block (width, body height, header height)
GD = 1.1                                # Gate diameter
RW = 2.8;   RH = 1.9                   # Reference block (same size as func blocks per guideline)

# Font sizes (CESYS523: body=24, title=one size larger, gates/labels=within one size)
FB = 24     # Body text
FH = 20     # Header (block ID)
FG = 22     # Gate text (IT, OR, AND)
FL = 20     # Label text (italic)
FS_TITLE = 28


# ─────────────────────────────────────
# HELPERS
# ─────────────────────────────────────
def set_anchor_mid(shape):
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def add_func_block(slide, cx, cy, fid, fname, fill=WHITE):
    x = cx - BW / 2
    y = cy - (BBH + BHH) / 2
    # Header
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(BW), Inches(BHH))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = WHITE
    hdr.line.color.rgb = BLACK; hdr.line.width = BORDER_W
    tf = hdr.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = fid
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(FH); p.font.color.rgb = BLACK
    set_anchor_mid(hdr)
    # Body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y + BHH), Inches(BW), Inches(BBH))
    body.fill.solid(); body.fill.fore_color.rgb = fill
    body.line.color.rgb = BLACK; body.line.width = BORDER_W
    tf = body.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = fname
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(FB); p.font.color.rgb = BLACK
    set_anchor_mid(body)


def add_ref_block(slide, cx, cy, ref_id, ref_name):
    """Reference block: bracket style per CESYS523.
    Drawn as two thick vertical bracket lines with text between them."""
    w, h = RW, RH
    x, y = cx - w / 2, cy - h / 2
    bk_w = 0.25  # bracket tick width
    # Left bracket: vertical line + top/bottom ticks
    add_line_seg(slide, x, y, x, y + h, border=True)
    add_line_seg(slide, x, y, x + bk_w, y, border=True)
    add_line_seg(slide, x, y + h, x + bk_w, y + h, border=True)
    # Right bracket
    add_line_seg(slide, x + w, y, x + w, y + h, border=True)
    add_line_seg(slide, x + w, y, x + w - bk_w, y, border=True)
    add_line_seg(slide, x + w, y + h, x + w - bk_w, y + h, border=True)
    # Text
    tb = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(h - 0.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = ref_id
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(FH); p.font.bold = True; p.font.color.rgb = BLACK
    p2 = tf.add_paragraph(); p2.text = ref_name
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(FB); p2.font.color.rgb = BLACK
    set_anchor_mid(tb)


def add_gate(slide, cx, cy, text):
    gr = GD / 2
    g = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx - gr), Inches(cy - gr), Inches(GD), Inches(GD))
    g.fill.solid(); g.fill.fore_color.rgb = WHITE
    g.line.color.rgb = BLACK; g.line.width = BORDER_W
    tf = g.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(FG); p.font.bold = True; p.font.color.rgb = BLACK
    set_anchor_mid(g)


def _set_arrowhead(connector):
    spPr = connector._element.find(qn('p:spPr'))
    if spPr is None:
        return
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')


def add_line_seg(slide, x1, y1, x2, y2, arrowhead=False, dashed=False, border=False):
    """Single rectilinear line segment. border=True uses BORDER_W instead of LINE_W."""
    cxn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = BLACK
    cxn.line.width = BORDER_W if border else LINE_W
    if dashed:
        cxn.line.dash_style = 2
    if arrowhead:
        _set_arrowhead(cxn)
    return cxn


def add_arrow(slide, points, dashed=False):
    """Arrow with filled arrowhead.
    2-point paths use a native elbow connector (type 2) for PowerPoint editability.
    3+ point paths use chained straight segments for precise rectilinear routing."""
    if len(points) == 2:
        x1, y1 = points[0]
        x2, y2 = points[1]
        cxn = slide.shapes.add_connector(
            2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        cxn.line.color.rgb = BLACK
        cxn.line.width = LINE_W
        if dashed:
            cxn.line.dash_style = 2
        _set_arrowhead(cxn)
        return cxn
    for i in range(len(points) - 1):
        is_last = (i == len(points) - 2)
        add_line_seg(slide,
                     points[i][0], points[i][1],
                     points[i+1][0], points[i+1][1],
                     arrowhead=is_last, dashed=dashed)


def add_line(slide, points):
    """Multi-segment rectilinear line (no arrowhead)."""
    for i in range(len(points) - 1):
        add_line_seg(slide,
                     points[i][0], points[i][1],
                     points[i+1][0], points[i+1][1])


def add_label(slide, x, y, text, width=4.0, italic=True, fs=FL):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(fs); p.font.italic = italic; p.font.color.rgb = BLACK


def add_title(slide, text, subtitle=None):
    """Title in upper left corner, left-aligned (CESYS523 guideline)."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(28), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(FS_TITLE); p.font.bold = True; p.font.color.rgb = BLACK
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(28), Inches(0.6))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.LEFT
        p2.font.size = Pt(FL); p2.font.italic = True; p2.font.color.rgb = RGBColor(80, 80, 80)


def delete_all_slides(prs):
    while len(prs.slides._sldIdLst) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


# Edge helpers
def blk_l(cx):  return cx - BW / 2
def blk_r(cx):  return cx + BW / 2
def blk_t(cy):  return cy - (BBH + BHH) / 2
def blk_b(cy):  return cy + (BBH + BHH) / 2
def gate_l(cx): return cx - GD / 2
def gate_r(cx): return cx + GD / 2
def gate_t(cy): return cy - GD / 2
def gate_b(cy): return cy + GD / 2
def ref_l(cx):  return cx - RW / 2
def ref_r(cx):  return cx + RW / 2


# ════════════════════════════════════════
#  BUILD
# ════════════════════════════════════════
prs = Presentation(TEMPLATE)
delete_all_slides(prs)
blank = prs.slide_layouts[6]


# ════════════════════════════════════════
#  SLIDE 1 – TOP-LEVEL FFBD
#  F.1 → F.2 → IT(F.3) → IT(F.4) → IT(F.5) → F.6 → F.7
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "Team Heat Guard – Top-Level FFBD",
          "F.1 Onboard Org → F.7 Maintain Compliance")

YM = 9.5      # Main flow Y
YR1 = 13.5    # IT return row 1
YR2 = 15.0    # IT return row 2
YR3 = 16.5    # IT return row 3

xs = [2.5, 6.0, 9.5, 13.0, 16.5, 20.0, 23.5]
IT3_O = 8.0;  IT3_C = 11.0
IT4_O = 11.5; IT4_C = 14.5
IT5_O = 15.0; IT5_C = 18.0

add_func_block(s, xs[0], YM, "F.1", "Onboard\nOrganization")
add_func_block(s, xs[1], YM, "F.2", "Onboard\nWorkers")
add_func_block(s, xs[2], YM, "F.3", "Collect\nEnvironmental\nData")
add_func_block(s, xs[3], YM, "F.4", "Predict\nHeat Strain", fill=CORE_FILL)
add_func_block(s, xs[4], YM, "F.5", "Classify\n& Alert")
add_func_block(s, xs[5], YM, "F.6", "Monitor\n& Report")
add_func_block(s, xs[6], YM, "F.7", "Maintain\nCompliance")

add_gate(s, IT3_O, YM, "IT"); add_gate(s, IT3_C, YM, "IT")
add_gate(s, IT4_O, YM, "IT"); add_gate(s, IT4_C, YM, "IT")
add_gate(s, IT5_O, YM, "IT"); add_gate(s, IT5_C, YM, "IT")

add_arrow(s, [(blk_r(xs[0]), YM), (blk_l(xs[1]), YM)])
add_arrow(s, [(blk_r(xs[1]), YM), (gate_l(IT3_O), YM)])
add_arrow(s, [(gate_r(IT3_O), YM), (blk_l(xs[2]), YM)])
add_arrow(s, [(blk_r(xs[2]), YM), (gate_l(IT3_C), YM)])
add_arrow(s, [(gate_r(IT3_C), YM), (gate_l(IT4_O), YM)])
add_arrow(s, [(gate_r(IT4_O), YM), (blk_l(xs[3]), YM)])
add_arrow(s, [(blk_r(xs[3]), YM), (gate_l(IT4_C), YM)])
add_arrow(s, [(gate_r(IT4_C), YM), (gate_l(IT5_O), YM)])
add_arrow(s, [(gate_r(IT5_O), YM), (blk_l(xs[4]), YM)])
add_arrow(s, [(blk_r(xs[4]), YM), (gate_l(IT5_C), YM)])
add_arrow(s, [(gate_r(IT5_C), YM), (blk_l(xs[5]), YM)])
add_arrow(s, [(blk_r(xs[5]), YM), (blk_l(xs[6]), YM)])

# IT return loops
add_arrow(s, [(IT3_C, gate_b(YM)), (IT3_C, YR1), (IT3_O, YR1), (IT3_O, gate_b(YM))])
add_arrow(s, [(IT4_C, gate_b(YM)), (IT4_C, YR2), (IT4_O, YR2), (IT4_O, gate_b(YM))])
add_arrow(s, [(IT5_C, gate_b(YM)), (IT5_C, YR3), (IT5_O, YR3), (IT5_O, gate_b(YM))])

add_label(s, IT3_O - 0.5, YR1 - 0.7, "Continuous polling loop", width=6.0)
add_label(s, IT4_O - 0.5, YR2 - 0.7, "Per-worker per-refresh cycle", width=6.0)
add_label(s, IT5_O - 0.5, YR3 - 0.7, "On every new prediction", width=6.0)


# ════════════════════════════════════════
#  SLIDE 2 – F.1 Onboard Organization
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "F.1 Onboard Organization", "6 functions – sequential tenant setup flow")

YM = 9.5
xs2 = [3.0, 7.5, 12.0, 16.5, 21.0, 25.5]

add_func_block(s, xs2[0], YM, "F.1.1", "Create\nTenant")
add_func_block(s, xs2[1], YM, "F.1.2", "Configure\nSSO")
add_func_block(s, xs2[2], YM, "F.1.3", "Set Up\nRBAC Roles")
add_func_block(s, xs2[3], YM, "F.1.4", "Configure Alert\nThresholds")
add_func_block(s, xs2[4], YM, "F.1.5", "Register\nJob Sites")
add_func_block(s, xs2[5], YM, "F.1.6", "Activate\nBilling")

for i in range(5):
    add_arrow(s, [(blk_r(xs2[i]), YM), (blk_l(xs2[i+1]), YM)])


# ════════════════════════════════════════
#  SLIDE 3 – F.2 Onboard Workers
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "F.2 Onboard Workers",
          "7 functions – ≤5 minutes, ≤6 input fields (from QFD target)")

YM = 9.5
xs3 = [2.5, 6.5, 10.5, 14.5, 18.5, 22.5, 26.5]

add_func_block(s, xs3[0], YM, "F.2.1", "Authenticate\nWorker")
add_func_block(s, xs3[1], YM, "F.2.2", "Collect\nBiometrics")
add_func_block(s, xs3[2], YM, "F.2.3", "Select\nActivity")
add_func_block(s, xs3[3], YM, "F.2.4", "Select\nClothing/PPE")
add_func_block(s, xs3[4], YM, "F.2.5", "Create Worker\nProfile")
add_func_block(s, xs3[5], YM, "F.2.5a", "Initialize\nAcclimatization\nProtocol")
add_func_block(s, xs3[6], YM, "F.2.6", "Trigger First\nPrediction")

for i in range(6):
    add_arrow(s, [(blk_r(xs3[i]), YM), (blk_l(xs3[i+1]), YM)])


# ════════════════════════════════════════
#  SLIDE 4 – F.3 Collect Environmental Data
#  OR gate (indoor/outdoor) + OR gate (data freshness, was decision diamond)
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "F.3 Collect Environmental Data",
          "8 functions – OR gates for indoor/outdoor routing and data freshness")

YM = 9.5
YU = 6.5     # Upper branch
YL = 12.5    # Lower branch

X_31 = 2.5
X_OR1 = 5.5;  JX_OR1 = 6.3
X_33 = 8.5    # Outdoor (upper)
X_34 = 8.5    # Indoor (lower)
MX_OR1 = 11.0; X_OR2 = 11.8
X_35 = 14.5
X_OR3 = 17.0  # Freshness OR open (was decision diamond)
X_37 = 20.0   # Fresh (upper)
X_38 = 20.0   # Stale (lower)
MX_OR3 = 23.0; X_OR4 = 23.8  # Freshness OR close

add_func_block(s, X_31, YM, "F.3.1", "Determine\nSite Location")
add_gate(s, X_OR1, YM, "OR"); add_gate(s, X_OR2, YM, "OR")
add_func_block(s, X_33, YU, "F.3.3", "Query Weather API\n(dew_point,\npressure_hpa)")
add_func_block(s, X_34, YL, "F.3.4", "Apply Indoor\nModel")
add_func_block(s, X_35, YM, "F.3.5", "Validate Data\nFreshness")
add_gate(s, X_OR3, YM, "OR"); add_gate(s, X_OR4, YM, "OR")
add_func_block(s, X_37, YU, "F.3.7", "Store\nEnvironmental\nSnapshot")
add_func_block(s, X_38, YL, "F.3.8", "Use Cached +\nWarn User")

# F.3.1 → OR1 open
add_arrow(s, [(blk_r(X_31), YM), (gate_l(X_OR1), YM)])
# OR1 open → junction
add_line(s, [(gate_r(X_OR1), YM), (JX_OR1, YM)])
# Junction → F.3.3 (upper/outdoor)
add_arrow(s, [(JX_OR1, YM), (JX_OR1, YU), (blk_l(X_33), YU)])
# Junction → F.3.4 (lower/indoor)
add_arrow(s, [(JX_OR1, YM), (JX_OR1, YL), (blk_l(X_34), YL)])
# F.3.3 → OR1 merge
add_line(s, [(blk_r(X_33), YU), (MX_OR1, YU), (MX_OR1, YM)])
# F.3.4 → OR1 merge
add_line(s, [(blk_r(X_34), YL), (MX_OR1, YL), (MX_OR1, YM)])
# OR1 merge → OR2 close
add_arrow(s, [(MX_OR1, YM), (gate_l(X_OR2), YM)])
# OR2 close → F.3.5
add_arrow(s, [(gate_r(X_OR2), YM), (blk_l(X_35), YM)])
# F.3.5 → OR3 open (freshness)
add_arrow(s, [(blk_r(X_35), YM), (gate_l(X_OR3), YM)])
# OR3 open → junction for freshness split
JX_OR3 = gate_r(X_OR3) + 0.3
add_line(s, [(gate_r(X_OR3), YM), (JX_OR3, YM)])
# Junction → F.3.7 (fresh, upper)
add_arrow(s, [(JX_OR3, YM), (JX_OR3, YU), (blk_l(X_37), YU)])
# Junction → F.3.8 (stale, lower)
add_arrow(s, [(JX_OR3, YM), (JX_OR3, YL), (blk_l(X_38), YL)])
# F.3.7 → OR4 merge
add_line(s, [(blk_r(X_37), YU), (MX_OR3, YU), (MX_OR3, YM)])
# F.3.8 → OR4 merge
add_line(s, [(blk_r(X_38), YL), (MX_OR3, YL), (MX_OR3, YM)])
# Merge → OR4 close
add_arrow(s, [(MX_OR3, YM), (gate_l(X_OR4), YM)])

# Labels
add_label(s, JX_OR1 + 0.2, YU - 0.8, "Outdoor", width=3.0)
add_label(s, JX_OR1 + 0.2, YL + 0.6, "Indoor", width=3.0)
add_label(s, JX_OR3 + 0.2, YU - 0.8, "Fresh (< 15 min)", width=4.0)
add_label(s, JX_OR3 + 0.2, YL + 0.6, "Stale – use cached", width=4.0)


# ════════════════════════════════════════
#  SLIDE 5 – F.4 Predict Heat Strain (CORE)
#
#  Row 1 (prep):  F.4.1 → ... → F.4.3b → OR(sport?) → [F.4.3c] → OR → F.4.4
#  Row 2 splits:  AND → upper (PHS): F.4.5 → AND(F.4.6 || F.4.7)
#                      → lower (WBGT): F.4.10 → F.4.11 → F.4.12 → F.4.13
#                 → AND close
#  Row 3 (post):  F.4.14 → F.4.15 → OR(ML?) → [F.4.16] → OR → F.4.8 → F.4.9
#  Ref block:     F.4a Ref (Pre-Event Risk Stratification)
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "F.4 Predict Heat Strain – CORE VALUE",
          "Dual-model (PHS+WBGT), sport-specific OR, ML OR, AND gates – MAE ≤ 0.3°C target")

# --- Row 1: Preparation chain (Y=4.5) ---
Y_R1 = 4.5
R1_xs = [2.0, 5.3, 8.6, 11.9, 15.2, 18.5]  # F.4.1 through F.4.3b

add_func_block(s, R1_xs[0], Y_R1, "F.4.1", "Retrieve\nWorker Profile")
add_func_block(s, R1_xs[1], Y_R1, "F.4.2", "Look Up\nMetabolic Rate\n(ISO 8996)")
add_func_block(s, R1_xs[2], Y_R1, "F.4.2a", "Apply\nAcclimatization\nFactor")
add_func_block(s, R1_xs[3], Y_R1, "F.4.3", "Apply Clothing\nFactor\n(ISO 7933)")
add_func_block(s, R1_xs[4], Y_R1, "F.4.3a", "Retrieve\nEnvironmental\nSnapshot")
add_func_block(s, R1_xs[5], Y_R1, "F.4.3b", "Compute\nWBGT Index")

# Row 1 sequential arrows
for i in range(5):
    add_arrow(s, [(blk_r(R1_xs[i]), Y_R1), (blk_l(R1_xs[i+1]), Y_R1)])

# Sport-specific OR gate
X_OR_S_O = 21.0    # Sport OR open
X_OR_S_C = 25.0    # Sport OR close
X_43c = 23.0        # Sport-specific block
Y_SPORT = 7.5       # Below Row 1
X_44 = 27.5          # F.4.4

add_gate(s, X_OR_S_O, Y_R1, "OR")
add_gate(s, X_OR_S_C, Y_R1, "OR")
add_func_block(s, X_43c, Y_SPORT, "F.4.3c", "Apply Sport-\nSpecific\nFactors")
add_func_block(s, X_44, Y_R1, "F.4.4", "Assemble\nAlgorithm\nInputs")

# F.4.3b → OR open
add_arrow(s, [(blk_r(R1_xs[5]), Y_R1), (gate_l(X_OR_S_O), Y_R1)])
# OR open → junction
JX_S = gate_r(X_OR_S_O) + 0.3
add_line(s, [(gate_r(X_OR_S_O), Y_R1), (JX_S, Y_R1)])
# Standard path: straight to OR close merge
MX_S = X_OR_S_C - GD / 2 - 0.3
add_line(s, [(JX_S, Y_R1), (MX_S, Y_R1)])
# Sport path: down to F.4.3c
add_arrow(s, [(JX_S, Y_R1), (JX_S, Y_SPORT), (blk_l(X_43c), Y_SPORT)])
# F.4.3c → merge back up to OR close
add_line(s, [(blk_r(X_43c), Y_SPORT), (MX_S, Y_SPORT), (MX_S, Y_R1)])
# Merge → OR close
add_arrow(s, [(MX_S, Y_R1), (gate_l(X_OR_S_C), Y_R1)])
# OR close → F.4.4
add_arrow(s, [(gate_r(X_OR_S_C), Y_R1), (blk_l(X_44), Y_R1)])

add_label(s, JX_S + 0.2, Y_R1 - 0.8, "Standard", width=3.0)
add_label(s, JX_S + 0.2, Y_SPORT + 0.7, "Sport-specific", width=4.0)

# --- Row 2: AND split (PHS upper, WBGT lower) ---
X_AND_O = 3.0;  Y_AND = 10.5
Y_PHS = 9.0
X_45 = 7.0;  X_AND_I_O = 11.0
X_46 = 14.5;  Y_46 = 8.0
X_47 = 14.5;  Y_47 = 10.0
MX_I = 17.0;  X_AND_I_C = 17.7

Y_WBGT = 13.0
X_410 = 7.0;  X_411 = 11.0;  X_412 = 15.0;  X_413 = 19.0

X_AND_C = 22.0

add_gate(s, X_AND_O, Y_AND, "AND")

# PHS upper branch
add_func_block(s, X_45, Y_PHS, "F.4.5", "Execute PHS\nAlgorithm\n(ISO 7933)", fill=CORE_FILL)
add_gate(s, X_AND_I_O, Y_PHS, "AND")
add_func_block(s, X_46, Y_46, "F.4.6", "Calculate\nTrec")
add_func_block(s, X_47, Y_47, "F.4.7", "Calculate\nDlim")
add_gate(s, X_AND_I_C, Y_PHS, "AND")

# WBGT lower branch
add_func_block(s, X_410, Y_WBGT, "F.4.10", "Compute\nWBGT\nScreening")
add_func_block(s, X_411, Y_WBGT, "F.4.11", "Apply OSHA\nAction Limits")
add_func_block(s, X_412, Y_WBGT, "F.4.12", "Classify\nWBGT Risk\nLevel")
add_func_block(s, X_413, Y_WBGT, "F.4.13", "Generate\nWBGT\nRecommendation")

add_gate(s, X_AND_C, Y_AND, "AND")

# F.4.4 → down to AND open
add_arrow(s, [(X_44, blk_b(Y_R1)),
              (X_44, Y_AND - 1.0),
              (X_AND_O + 1.5, Y_AND - 1.0),
              (X_AND_O + 1.5, Y_AND),
              (gate_r(X_AND_O), Y_AND)])

# AND open → junction, split
JX_A = gate_r(X_AND_O) + 0.3
add_line(s, [(gate_r(X_AND_O), Y_AND), (JX_A, Y_AND)])
add_arrow(s, [(JX_A, Y_AND), (JX_A, Y_PHS), (blk_l(X_45), Y_PHS)])
add_arrow(s, [(JX_A, Y_AND), (JX_A, Y_WBGT), (blk_l(X_410), Y_WBGT)])

# PHS: F.4.5 → inner AND open
add_arrow(s, [(blk_r(X_45), Y_PHS), (gate_l(X_AND_I_O), Y_PHS)])
JX_I = gate_r(X_AND_I_O) + 0.3
add_line(s, [(gate_r(X_AND_I_O), Y_PHS), (JX_I, Y_PHS)])
add_arrow(s, [(JX_I, Y_PHS), (JX_I, Y_46), (blk_l(X_46), Y_46)])
add_arrow(s, [(JX_I, Y_PHS), (JX_I, Y_47), (blk_l(X_47), Y_47)])
# F.4.6 / F.4.7 → inner AND merge → inner AND close
add_line(s, [(blk_r(X_46), Y_46), (MX_I, Y_46), (MX_I, Y_PHS)])
add_line(s, [(blk_r(X_47), Y_47), (MX_I, Y_47), (MX_I, Y_PHS)])
add_arrow(s, [(MX_I, Y_PHS), (gate_l(X_AND_I_C), Y_PHS)])

# Inner AND close → outer AND merge
MX_AC = X_AND_C - 0.7
add_line(s, [(gate_r(X_AND_I_C), Y_PHS), (MX_AC, Y_PHS), (MX_AC, Y_AND)])

# WBGT: sequential
add_arrow(s, [(blk_r(X_410), Y_WBGT), (blk_l(X_411), Y_WBGT)])
add_arrow(s, [(blk_r(X_411), Y_WBGT), (blk_l(X_412), Y_WBGT)])
add_arrow(s, [(blk_r(X_412), Y_WBGT), (blk_l(X_413), Y_WBGT)])
add_line(s, [(blk_r(X_413), Y_WBGT), (MX_AC, Y_WBGT), (MX_AC, Y_AND)])

# Merge → AND close
add_arrow(s, [(MX_AC, Y_AND), (gate_l(X_AND_C), Y_AND)])

add_label(s, JX_A + 0.2, Y_PHS - 1.3, "PHS deep analysis", width=5.0)
add_label(s, JX_A + 0.2, Y_WBGT + 0.8, "WBGT screening", width=5.0)

# --- Row 3: Post-processing (Y=16.0) ---
Y_R3 = 16.0
X_414 = 5.0;  X_415 = 9.0
X_OR_ML_O = 12.5;  X_OR_ML_C = 16.5  # ML OR gates (was decision diamond)
X_416 = 14.5;  Y_416 = 18.5           # ML correction (below Row 3)
X_48 = 19.5;  X_49 = 23.5

add_func_block(s, X_414, Y_R3, "F.4.14", "Reconcile\nPHS + WBGT\nResults", fill=CORE_FILL)
add_func_block(s, X_415, Y_R3, "F.4.15", "Check\nCompliance\nTrigger")
add_gate(s, X_OR_ML_O, Y_R3, "OR")
add_gate(s, X_OR_ML_C, Y_R3, "OR")
add_func_block(s, X_416, Y_416, "F.4.16", "Apply ML\nCorrection")
add_func_block(s, X_48, Y_R3, "F.4.8", "Store\nPrediction\nResult")
add_func_block(s, X_49, Y_R3, "F.4.9", "Log to\nAudit Trail")

# AND close → down to Row 3 → F.4.14
add_arrow(s, [(gate_r(X_AND_C), Y_AND),
              (X_AND_C + 0.8, Y_AND),
              (X_AND_C + 0.8, Y_R3 + 2.0),
              (X_414 - 2.0, Y_R3 + 2.0),
              (X_414 - 2.0, Y_R3),
              (blk_l(X_414), Y_R3)])

# F.4.14 → F.4.15
add_arrow(s, [(blk_r(X_414), Y_R3), (blk_l(X_415), Y_R3)])
# F.4.15 → OR ML open
add_arrow(s, [(blk_r(X_415), Y_R3), (gate_l(X_OR_ML_O), Y_R3)])
# OR ML open → junction
JX_ML = gate_r(X_OR_ML_O) + 0.3
add_line(s, [(gate_r(X_OR_ML_O), Y_R3), (JX_ML, Y_R3)])
# No model: straight to OR ML close merge
MX_ML = X_OR_ML_C - GD / 2 - 0.3
add_line(s, [(JX_ML, Y_R3), (MX_ML, Y_R3)])
# Model available: down to F.4.16
add_arrow(s, [(JX_ML, Y_R3), (JX_ML, Y_416), (blk_l(X_416), Y_416)])
# F.4.16 → merge back up
add_line(s, [(blk_r(X_416), Y_416), (MX_ML, Y_416), (MX_ML, Y_R3)])
# Merge → OR ML close
add_arrow(s, [(MX_ML, Y_R3), (gate_l(X_OR_ML_C), Y_R3)])
# OR ML close → F.4.8
add_arrow(s, [(gate_r(X_OR_ML_C), Y_R3), (blk_l(X_48), Y_R3)])
# F.4.8 → F.4.9
add_arrow(s, [(blk_r(X_48), Y_R3), (blk_l(X_49), Y_R3)])

add_label(s, JX_ML + 0.2, Y_R3 - 0.8, "No model – skip", width=4.0)
add_label(s, JX_ML + 0.2, Y_416 + 0.7, "Model available", width=4.0)

# Reference block: F.4a Pre-Event Risk Stratification
X_REF = 27.5;  Y_REF = Y_R3
add_ref_block(s, X_REF, Y_REF, "F.4a Ref", "Pre-Event Risk\nStratification")
# Precedes (dashed) arrow from F.4.9 → ref block
add_arrow(s, [(blk_r(X_49), Y_R3), (ref_l(X_REF), Y_REF)], dashed=True)


# ════════════════════════════════════════
#  SLIDE 6 – F.5 Classify & Alert
#  Two OR gate pairs (were decision diamonds)
#  + F.5.4a intervention prescription + F.5.7a-g chain
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "F.5 Classify & Alert",
          "17 functions – OR gates + intervention effectiveness chain")

YM = 7.0
YU = 4.0
YL = 10.5
Y_INT = 14.0

X_51 = 2.0
X_OR1_O = 4.5;  X_OR1_C = 27.0   # Risk changed OR gates (wide span)
X_53 = 7.5      # No change (upper)
X_54 = 7.5      # Changed (main)
X_54a = 10.5
X_55 = 13.5
X_OR2_O = 16.0;  X_OR2_C = 24.5   # Ack OR gates
X_57 = 19.0     # Ack yes (upper)
X_58 = 19.0     # Ack no (lower)
X_59 = 22.0
X_510 = 27.0    # Final (inside OR1 close position)

# Intervention chain X positions
X_57a = 3.0;  X_57b = 6.5;  X_57c = 10.0;  X_57d = 13.5
X_57e = 17.0; X_57f = 20.5; X_57g = 24.0

add_func_block(s, X_51, YM, "F.5.1", "Apply Risk\nThresholds")

# OR1: Risk level changed?
add_gate(s, X_OR1_O, YM, "OR")

add_func_block(s, X_53, YU, "F.5.3", "Update\nDashboard Only")
add_func_block(s, X_54, YM, "F.5.4", "Generate Alert\n+ Guidance")
add_func_block(s, X_54a, YM, "F.5.4a", "Generate\nIntervention\nPrescription")
add_func_block(s, X_55, YM, "F.5.5", "Push to\nWorker")

# OR2: Acknowledged?
add_gate(s, X_OR2_O, YM, "OR")

add_func_block(s, X_57, YU, "F.5.7", "Log\nAcknowledgment")
add_func_block(s, X_58, YL, "F.5.8", "Escalate to\nSupervisor")
add_func_block(s, X_59, YL, "F.5.9", "Log\nEscalation")

add_gate(s, X_OR2_C, YM, "OR")

add_func_block(s, X_510, YM, "F.5.10", "Update\nManager\nDashboard")

add_gate(s, X_OR1_C, YM, "OR")  # OR1 close - at same X as F.5.10

# Intervention chain blocks
add_func_block(s, X_57a, Y_INT, "F.5.7a", "Record\nPre-Intervention\nTrec")
add_func_block(s, X_57b, Y_INT, "F.5.7b", "Worker Completes\nIntervention")
add_func_block(s, X_57c, Y_INT, "F.5.7c", "Record\nCompletion GPS")
add_func_block(s, X_57d, Y_INT, "F.5.7d", "Trigger\nRe-Prediction")
add_func_block(s, X_57e, Y_INT, "F.5.7e", "Record\nPost-Intervention\nTrec")
add_func_block(s, X_57f, Y_INT, "F.5.7f", "Compute\nEffectiveness\nScore")
add_func_block(s, X_57g, Y_INT, "F.5.7g", "Feed ML\nTraining Pipeline")

# F.5.1 → OR1 open
add_arrow(s, [(blk_r(X_51), YM), (gate_l(X_OR1_O), YM)])

# OR1 open → junction
JX1 = gate_r(X_OR1_O) + 0.3
add_line(s, [(gate_r(X_OR1_O), YM), (JX1, YM)])
# No change path (upper): → F.5.3
add_arrow(s, [(JX1, YM), (JX1, YU), (blk_l(X_53), YU)])
# Changed path (main): → F.5.4
add_arrow(s, [(JX1, YM), (blk_l(X_54), YM)])

# F.5.4 → F.5.4a → F.5.5 → OR2 open
add_arrow(s, [(blk_r(X_54), YM), (blk_l(X_54a), YM)])
add_arrow(s, [(blk_r(X_54a), YM), (blk_l(X_55), YM)])
add_arrow(s, [(blk_r(X_55), YM), (gate_l(X_OR2_O), YM)])

# OR2 open → junction
JX2 = gate_r(X_OR2_O) + 0.3
add_line(s, [(gate_r(X_OR2_O), YM), (JX2, YM)])
# Ack'd path (upper): → F.5.7
add_arrow(s, [(JX2, YM), (JX2, YU), (blk_l(X_57), YU)])
# Not ack'd path (lower): → F.5.8
add_arrow(s, [(JX2, YM), (JX2, YL), (blk_l(X_58), YL)])

# F.5.8 → F.5.9
add_arrow(s, [(blk_r(X_58), YL), (blk_l(X_59), YL)])

# All paths merge to OR2 close
MX2 = X_OR2_C - GD / 2 - 0.3
# F.5.7 → merge (after intervention chain routes back)
# F.5.9 → merge
add_line(s, [(blk_r(X_59), YL), (MX2, YL), (MX2, YM)])
# F.5.3 → merge to OR1 close (via upper path straight across)
MX1 = X_OR1_C - GD / 2 - 0.3
add_line(s, [(blk_r(X_53), YU), (MX1, YU), (MX1, YM)])

# OR2 close → F.5.10
add_arrow(s, [(MX2, YM), (gate_l(X_OR2_C), YM)])
add_arrow(s, [(gate_r(X_OR2_C), YM), (blk_l(X_510), YM)])

# F.5.10 → OR1 close (they're at same X, so arrow from F.5.10 right to OR1 close right)
# Actually OR1 close wraps the whole changed path including F.5.10
# The no-change path (F.5.3) and the changed path (→F.5.10) merge at OR1 close
add_arrow(s, [(blk_r(X_510), YM), (gate_l(X_OR1_C), YM)])

# F.5.7 → down to intervention chain
add_arrow(s, [(X_57, blk_b(YU)),
              (X_57, Y_INT - 1.5),
              (X_57a - 1.5, Y_INT - 1.5),
              (X_57a - 1.5, Y_INT),
              (blk_l(X_57a), Y_INT)])

# Intervention chain sequential
add_arrow(s, [(blk_r(X_57a), Y_INT), (blk_l(X_57b), Y_INT)])
add_arrow(s, [(blk_r(X_57b), Y_INT), (blk_l(X_57c), Y_INT)])
add_arrow(s, [(blk_r(X_57c), Y_INT), (blk_l(X_57d), Y_INT)])
add_arrow(s, [(blk_r(X_57d), Y_INT), (blk_l(X_57e), Y_INT)])
add_arrow(s, [(blk_r(X_57e), Y_INT), (blk_l(X_57f), Y_INT)])
add_arrow(s, [(blk_r(X_57f), Y_INT), (blk_l(X_57g), Y_INT)])

# F.5.7g → route up to OR2 merge
add_arrow(s, [(blk_r(X_57g), Y_INT),
              (MX2 + 0.3, Y_INT),
              (MX2 + 0.3, YM)])

# Labels
add_label(s, JX1 + 0.2, YU - 0.8, "No change", width=3.0)
add_label(s, JX1 + 1.5, YM - 0.8, "Changed", width=3.0)
add_label(s, JX2 + 0.2, YU - 0.8, "Acknowledged", width=4.0)
add_label(s, JX2 + 0.2, YL + 0.6, "Not ack'd (AMBER: 5 min, RED: 2 min)", width=7.0)
add_label(s, X_57a - 1.5, Y_INT - 2.0, "Intervention Effectiveness Chain", width=8.0, italic=False, fs=FB)


# ════════════════════════════════════════
#  SLIDE 7 – F.6 Monitor & Report
#  AND gate: 5 branches (renamed to functional)
#  + OR gate for report (was decision diamond)
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "F.6 Monitor & Report",
          "12 functions – 5-branch AND gate + OR gate (report generation)")

Y_B1 = 4.5;  Y_B2 = 7.0;  Y_B3 = 9.5;  Y_B4 = 12.0;  Y_B5 = 14.5
YM = 9.5
YU_RPT = 6.0;  YL_RPT = 13.0

X_61 = 2.5
X_AND1 = 5.0;  JX_A = 5.7
X_BRANCH = 9.0
MX_A = 12.5;  X_AND2 = 13.3
X_64 = 16.0
X_64a = 19.5
X_OR_R_O = 22.0  # Report OR open (was decision)
X_66 = 25.0;  X_67 = 28.0  # Report yes
X_68 = 25.0                  # Continue no
MX_R = 27.0;  X_OR_R_C = 27.8  # Report OR close

add_func_block(s, X_61, YM, "F.6.1", "Aggregate\nWorker Status\nby Site")
add_gate(s, X_AND1, YM, "AND")

# 5 AND branches - RENAMED to functional names
add_func_block(s, X_BRANCH, Y_B1, "F.6.2", "Render Real-Time\nRisk Grid")
add_func_block(s, X_BRANCH, Y_B2, "F.6.3", "Compute Site\nRisk Summary")
add_func_block(s, X_BRANCH, Y_B3, "F.6.3a", "Aggregate Team\nSession Status")
add_func_block(s, X_BRANCH, Y_B4, "F.6.3b", "Compute Optimal\nRotations")
add_func_block(s, X_BRANCH, Y_B5, "F.6.3c", "Retrieve Pre-Event\nRisk Scores")

add_gate(s, X_AND2, YM, "AND")
add_func_block(s, X_64, YM, "F.6.4", "Display on\nManager\nDashboard")
add_func_block(s, X_64a, YM, "F.6.4a", "Display\nAcclimatization\nProgress")
add_gate(s, X_OR_R_O, YM, "OR")
add_func_block(s, X_66, YU_RPT, "F.6.6", "Compile\nCompliance\nReport")
add_func_block(s, X_67, YU_RPT, "F.6.7", "Export\nPDF / Excel")
add_func_block(s, X_68, YL_RPT, "F.6.8", "Continue\nMonitoring")
add_gate(s, X_OR_R_C, YM, "OR")

# F.6.1 → AND open
add_arrow(s, [(blk_r(X_61), YM), (gate_l(X_AND1), YM)])
# AND open → junction
add_line(s, [(gate_r(X_AND1), YM), (JX_A, YM)])
# Junction → 5 branches
add_arrow(s, [(JX_A, YM), (JX_A, Y_B1), (blk_l(X_BRANCH), Y_B1)])
add_arrow(s, [(JX_A, YM), (JX_A, Y_B2), (blk_l(X_BRANCH), Y_B2)])
add_arrow(s, [(JX_A, YM), (blk_l(X_BRANCH), Y_B3)])
add_arrow(s, [(JX_A, YM), (JX_A, Y_B4), (blk_l(X_BRANCH), Y_B4)])
add_arrow(s, [(JX_A, YM), (JX_A, Y_B5), (blk_l(X_BRANCH), Y_B5)])

# 5 branches → AND merge
for yb in [Y_B1, Y_B2, Y_B3, Y_B4, Y_B5]:
    add_line(s, [(blk_r(X_BRANCH), yb), (MX_A, yb), (MX_A, YM)])

# AND merge → AND close
add_arrow(s, [(MX_A, YM), (gate_l(X_AND2), YM)])
# AND close → F.6.4
add_arrow(s, [(gate_r(X_AND2), YM), (blk_l(X_64), YM)])
# F.6.4 → F.6.4a
add_arrow(s, [(blk_r(X_64), YM), (blk_l(X_64a), YM)])
# F.6.4a → OR report open
add_arrow(s, [(blk_r(X_64a), YM), (gate_l(X_OR_R_O), YM)])

# OR report open → junction
JX_R = gate_r(X_OR_R_O) + 0.3
add_line(s, [(gate_r(X_OR_R_O), YM), (JX_R, YM)])
# Report path (upper): → F.6.6 → F.6.7
add_arrow(s, [(JX_R, YM), (JX_R, YU_RPT), (blk_l(X_66), YU_RPT)])
add_arrow(s, [(blk_r(X_66), YU_RPT), (blk_l(X_67), YU_RPT)])
# Continue path (lower): → F.6.8
add_arrow(s, [(JX_R, YM), (JX_R, YL_RPT), (blk_l(X_68), YL_RPT)])

# Both merge to OR close
add_line(s, [(blk_r(X_67), YU_RPT), (MX_R, YU_RPT), (MX_R, YM)])
add_line(s, [(blk_r(X_68), YL_RPT), (MX_R, YL_RPT), (MX_R, YM)])
add_arrow(s, [(MX_R, YM), (gate_l(X_OR_R_C), YM)])

# Labels
add_label(s, JX_A + 0.5, Y_B1 - 0.8, "Real-time grid", width=4.0)
add_label(s, JX_A + 0.5, Y_B2 - 0.8, "Summary stats", width=4.0)
add_label(s, JX_A + 0.5, Y_B3 - 0.8, "Team sessions", width=4.0)
add_label(s, JX_A + 0.5, Y_B4 - 0.8, "Rotation recs", width=4.0)
add_label(s, JX_A + 0.5, Y_B5 - 0.8, "Stratification", width=4.0)
add_label(s, JX_R + 0.2, YU_RPT - 0.8, "Report requested", width=4.0)
add_label(s, JX_R + 0.2, YL_RPT + 0.6, "Continue monitoring", width=5.0)


# ════════════════════════════════════════
#  SLIDE 8 – F.7 Maintain Compliance
#  AND gate: F.7.1 (regulatory) || F.7.2 (acclimatization cron)
#  OR gate within F.7.1 path (was decision diamond)
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "F.7 Maintain Compliance",
          "AND gate (parallel ops) + OR gate (regulatory change) – 30-year retention")

YM = 9.5
YU_REG = 5.5;  YL_REG = 9.5  # OR paths within regulatory branch
Y_ACCL = 14.0                  # Acclimatization branch (lower AND)

# AND gate for parallel: regulatory monitoring || acclimatization cron
X_AND_O = 2.5;  X_AND_C = 22.0

# Regulatory branch (upper AND path, centered at Y=7.5)
Y_REG = 7.5
X_71 = 5.5
X_OR_O = 8.0   # Regulatory OR open
X_73 = 11.5;  X_74 = 15.0;  X_75 = 18.5   # Yes path (upper)
X_76 = 11.5;  X_77 = 15.0                   # No path (lower)
MX_REG = 19.5;  X_OR_C = 20.3              # Regulatory OR close

# Acclimatization branch (lower AND path)
X_72 = 5.5

# Post-AND
X_78 = 25.0

add_gate(s, X_AND_O, YM, "AND")

# Regulatory branch
add_func_block(s, X_71, Y_REG, "F.7.1", "Monitor\nRegulatory\nChanges")
add_gate(s, X_OR_O, Y_REG, "OR")

add_func_block(s, X_73, YU_REG, "F.7.3", "Update\nThreshold\nLibrary")
add_func_block(s, X_74, YU_REG, "F.7.4", "Notify\nAffected Orgs")
add_func_block(s, X_75, YU_REG, "F.7.5", "Update Report\nTemplates")
add_func_block(s, X_76, YL_REG, "F.7.6", "Maintain Data\nRetention\n(30-year)")
add_func_block(s, X_77, YL_REG, "F.7.7", "Process\nData Access\nRequests")

add_gate(s, X_OR_C, Y_REG, "OR")

# Acclimatization branch
add_func_block(s, X_72, Y_ACCL, "F.7.2", "Advance\nAcclimatization\nProtocol")

add_gate(s, X_AND_C, YM, "AND")
add_func_block(s, X_78, YM, "F.7.8", "Run Compliance\nSelf-Audit")

# AND open → junction
JX_AND = gate_r(X_AND_O) + 0.3
add_line(s, [(gate_r(X_AND_O), YM), (JX_AND, YM)])
# Junction → regulatory branch (upper)
add_arrow(s, [(JX_AND, YM), (JX_AND, Y_REG), (blk_l(X_71), Y_REG)])
# Junction → acclimatization branch (lower)
add_arrow(s, [(JX_AND, YM), (JX_AND, Y_ACCL), (blk_l(X_72), Y_ACCL)])

# F.7.1 → OR open
add_arrow(s, [(blk_r(X_71), Y_REG), (gate_l(X_OR_O), Y_REG)])
# OR open → junction
JX_OR = gate_r(X_OR_O) + 0.3
add_line(s, [(gate_r(X_OR_O), Y_REG), (JX_OR, Y_REG)])
# New standard (upper): → F.7.3 → F.7.4 → F.7.5
add_arrow(s, [(JX_OR, Y_REG), (JX_OR, YU_REG), (blk_l(X_73), YU_REG)])
add_arrow(s, [(blk_r(X_73), YU_REG), (blk_l(X_74), YU_REG)])
add_arrow(s, [(blk_r(X_74), YU_REG), (blk_l(X_75), YU_REG)])
# Routine (lower): → F.7.6 → F.7.7
add_arrow(s, [(JX_OR, Y_REG), (JX_OR, YL_REG), (blk_l(X_76), YL_REG)])
add_arrow(s, [(blk_r(X_76), YL_REG), (blk_l(X_77), YL_REG)])

# Both OR paths merge
add_line(s, [(blk_r(X_75), YU_REG), (MX_REG, YU_REG), (MX_REG, Y_REG)])
add_line(s, [(blk_r(X_77), YL_REG), (MX_REG, YL_REG), (MX_REG, Y_REG)])
add_arrow(s, [(MX_REG, Y_REG), (gate_l(X_OR_C), Y_REG)])

# OR close → merge to AND close
MX_AND = X_AND_C - GD / 2 - 0.3
add_line(s, [(gate_r(X_OR_C), Y_REG), (MX_AND, Y_REG), (MX_AND, YM)])

# Acclimatization → merge to AND close
add_line(s, [(blk_r(X_72), Y_ACCL), (MX_AND, Y_ACCL), (MX_AND, YM)])

# AND merge → AND close
add_arrow(s, [(MX_AND, YM), (gate_l(X_AND_C), YM)])
# AND close → F.7.8
add_arrow(s, [(gate_r(X_AND_C), YM), (blk_l(X_78), YM)])

# Labels
add_label(s, JX_OR + 0.2, YU_REG - 0.8, "New standard published", width=6.0)
add_label(s, JX_OR + 0.2, YL_REG + 0.6, "Routine maintenance", width=5.0)
add_label(s, JX_AND + 0.2, Y_REG - 1.5, "Regulatory monitoring", width=6.0)
add_label(s, JX_AND + 0.2, Y_ACCL + 0.7, "Daily acclimatization cron", width=6.0)


# ════════════════════════════════════════
#  SLIDE 9 – F.4a Pre-Event Risk Stratification (NEW)
#  6 blocks, linear flow
# ════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_title(s, "Function 4a: Pre-Event Risk Stratification",
          "6 functions – forecast + historical data → risk score (0–100) per worker per event")

YM = 9.5
xs9 = [2.5, 7.0, 11.5, 16.0, 20.5, 25.0]

add_func_block(s, xs9[0], YM, "F.4a.1", "Retrieve\nWeather\nForecast (48hr)")
add_func_block(s, xs9[1], YM, "F.4a.2", "Retrieve\nHistorical\nWorker Data")
add_func_block(s, xs9[2], YM, "F.4a.3", "Map Planned\nActivity\nSchedule")
add_func_block(s, xs9[3], YM, "F.4a.4", "Compute\nRisk Score\n(0–100)")
add_func_block(s, xs9[4], YM, "F.4a.5", "Store\nStratification\nResult")
add_func_block(s, xs9[5], YM, "F.4a.6", "Log to\nAudit Trail")

for i in range(5):
    add_arrow(s, [(blk_r(xs9[i]), YM), (blk_l(xs9[i+1]), YM)])


# ════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════
prs.save(OUTPUT)
print(f"THG FFBD V3 saved to: {OUTPUT}")
print(f"  9 slides covering all functional flows:")
print(f"  Slide 1: Top-Level (F.1–F.7, 3 IT loops)")
print(f"  Slide 2: F.1 Onboard Organization (6 blocks)")
print(f"  Slide 3: F.2 Onboard Workers (7 blocks, +F.2.5a Acclimatization)")
print(f"  Slide 4: F.3 Collect Environmental Data (OR gates, no decision diamonds)")
print(f"  Slide 5: F.4 Predict Heat Strain – CORE (sport OR, ML OR, AND gates, F.4a ref)")
print(f"  Slide 6: F.5 Classify & Alert (OR gates + intervention chain)")
print(f"  Slide 7: F.6 Monitor & Report (5-branch AND + OR gate, renamed blocks)")
print(f"  Slide 8: F.7 Maintain Compliance (AND + OR gates, +F.7.2 acclimatization)")
print(f"  Slide 9: F.4a Pre-Event Risk Stratification (6 blocks, NEW)")
print(f"\nV3 changes from V2:")
print(f"  - Font: body 24pt, header 20pt, gates 22pt, title 28pt left-aligned")
print(f"  - Arrows: 2.5pt weight, rectilinear, filled triangle arrowheads")
print(f"  - Decision diamonds → OR gate pairs (FFBD convention)")
print(f"  - Reference blocks: bracket style (CESYS523)")
print(f"  - F.4: Sport-specific OR gate + F.4a reference block")
print(f"  - F.6: Blocks renamed to functional names")
print(f"  - F.7: AND gate + F.7.2 Advance Acclimatization Protocol")
print(f"  - F.4a: New slide – Pre-Event Risk Stratification")
