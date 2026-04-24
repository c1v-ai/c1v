#!/usr/bin/env python3
"""
Data-driven UCBD PowerPoint generator.

Reads the JSON instances produced by the Module 2 pipeline and emits one
.pptx with a title slide, one slide per UCBD (combined or split based on
size), and a summary slide.

Inputs (in the project folder, relative to --project):
    system_context_summary.json                (optional — drives title slide)
    ucbd/UC*.ucbd.json                         (one per use case)
    requirements_table.json                    (for REQ-ID lookup + summary)
    constants_table.json                       (for summary slide)

Output:
    <project>/<system_name>_UCBD.pptx

Usage:
    pip install python-pptx
    python generate_ucbd_pptx.py --project /path/to/module-2-requirements/
    python generate_ucbd_pptx.py --project . --output my_ucbd.pptx

No project-specific code. Every value is sourced from the JSONs. To change
the output, edit the JSONs and re-run — never edit this script.
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Ensure local imports work
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from _ucbd_helpers import (
    add_combined_slide,
    add_split_slides,
    can_combine,
    FONT,
    BLACK,
    SW,
    SH,
    BLANK_LAYOUT,
)


# ── JSON loaders ────────────────────────────────────────────────

def load_json(path: Path) -> dict:
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def load_context(project_dir: Path) -> dict:
    path = project_dir / "system_context_summary.json"
    if not path.exists():
        return {
            "system_name": "System",
            "project_metadata": {"project_name": "Project"},
        }
    return load_json(path)


def load_ucbds(project_dir: Path) -> list[dict]:
    ucbd_dir = project_dir / "ucbd"
    if not ucbd_dir.exists():
        # Fallback: scan project_dir root for *.ucbd.json
        candidates = sorted(project_dir.glob("*.ucbd.json"))
    else:
        candidates = sorted(ucbd_dir.glob("*.ucbd.json"))
    if not candidates:
        sys.exit(
            f"ERROR: no *.ucbd.json files found in {ucbd_dir} or {project_dir}. "
            "Run Phases 3-5 of the KB first."
        )
    return [load_json(p) for p in candidates]


def load_requirements(project_dir: Path) -> dict:
    path = project_dir / "requirements_table.json"
    if not path.exists():
        return {"requirements_table": []}
    return load_json(path)


def load_constants(project_dir: Path) -> dict:
    path = project_dir / "constants_table.json"
    if not path.exists():
        return {"constants_table": []}
    return load_json(path)


# ── JSON → tuple adapters (consumed by _ucbd_helpers) ──────────

CONSTANT_PATTERN = re.compile(r"\b[A-Z][A-Z0-9_]{2,}\b")
NON_CONSTANT_TOKENS = {
    "ID", "HTTP", "HTTPS", "URL", "API", "JSON", "XML",
    "SQL", "SSO", "SAML", "OAUTH", "JWT", "UTC", "ISO",
}


def extract_constant_refs(text: str, known_constants: set[str]) -> list[str]:
    """Return constants from `known_constants` that appear in `text`."""
    tokens = set(CONSTANT_PATTERN.findall(text)) - NON_CONSTANT_TOKENS
    return sorted(tokens & known_constants)


def build_metadata(ucbd: dict) -> list[tuple[str, str]]:
    """Convert a UCBD JSON instance into the (label, value) tuples that
    `_ucbd_helpers.add_metadata` expects."""
    use_case_name = ucbd.get("use_case_name", "Untitled Use Case")
    columns = ucbd.get("_columns_plan", {})
    primary_actor = columns.get("A_primary_actor", "Actor")

    initial = "; ".join(
        strip_numbering(item) for item in ucbd.get("initial_conditions", [])
    ) or "—"
    ending = "; ".join(
        strip_numbering(item) for item in ucbd.get("ending_conditions", [])
    ) or "—"

    meta = [
        ("Use Case", use_case_name),
        ("Primary Actor", primary_actor),
    ]

    # Derive "Trigger" from first non-empty primary_actor row, if present
    trigger = first_trigger(ucbd.get("actor_steps_table", []))
    if trigger:
        meta.append(("Trigger", trigger))

    meta.append(("Start", initial))
    meta.append(("End", ending))
    return meta


def strip_numbering(text: str) -> str:
    """Remove leading '1. ', '2. ', etc. from a list item."""
    return re.sub(r"^\s*\d+[\.\)]\s*", "", text).strip()


def first_trigger(rows: list[dict]) -> str | None:
    """Return the first non-empty primary_actor cell, truncated if long."""
    for row in rows:
        text = row.get("primary_actor", "").strip()
        if text:
            return text[:120] + ("…" if len(text) > 120 else "")
    return None


def build_steps(ucbd: dict) -> list[tuple[str, str]]:
    """Convert actor_steps_table rows into (step_number, description) tuples.
    Each step shows which column/actor acted, prefixed to the text."""
    rows = ucbd.get("actor_steps_table", [])
    columns = ucbd.get("_columns_plan", {})

    # Derive column display labels
    labels = {
        "primary_actor": columns.get("A_primary_actor", "Primary Actor"),
        "the_system": "System",
        "other_actors": columns.get("C_other_actor_1", "Other"),
        "extra_actor_col": columns.get("D_other_actor_2", "Other"),
    }

    steps = []
    for i, row in enumerate(rows, start=1):
        for key in ("primary_actor", "the_system", "other_actors", "extra_actor_col"):
            text = (row.get(key) or "").strip()
            if text:
                prefix = labels[key]
                steps.append((str(i), f"{prefix}: {text}"))
                break
        else:
            # Empty row — skip
            continue
    return steps


def build_requirements_for_ucbd(
    ucbd_id: str,
    all_requirements: list[dict],
    known_constants: set[str],
) -> list[tuple[str, str, str]]:
    """Filter requirements whose source_ucbd matches ucbd_id, return (id, text, constants)."""
    out = []
    for req in all_requirements:
        if req.get("source_ucbd") != ucbd_id:
            continue
        rid = req.get("index", "")
        text = req.get("requirement", "")
        refs = extract_constant_refs(text, known_constants)
        const_label = ", ".join(refs) if refs else "—"
        out.append((rid, text, const_label))
    return out


def ucbd_id_from_instance(ucbd: dict) -> str:
    """Extract e.g. 'UC01' from metadata.document_id 'C1V-UCBD-UC01'."""
    doc_id = ucbd.get("metadata", {}).get("document_id", "")
    # Assume last segment after "-" is the UC id
    if "-" in doc_id:
        return doc_id.rsplit("-", 1)[-1]
    # Fallback: infer from filename if present in _output_path
    out_path = ucbd.get("_output_path", "")
    m = re.search(r"(UC\d+)", out_path)
    return m.group(1) if m else "UC??"


# ── Non-UCBD slides (title + summary) ──────────────────────────

def add_title_slide(prs, context: dict, ucbd_count: int, req_count: int, const_count: int):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    system_name = context.get("system_name", "System")
    project_name = context.get("project_metadata", {}).get("project_name", "")

    # Title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.0))
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "Use Case Behavioral Diagrams"
    r.font.size = Pt(32)
    r.font.name = FONT
    r.font.color.rgb = BLACK
    r.font.bold = True

    # Subtitle — project name
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(0.6))
    tf2 = tb2.text_frame
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run()
    r2.text = f"{system_name}" + (f" — {project_name}" if project_name and project_name != system_name else "")
    r2.font.size = Pt(18)
    r2.font.name = FONT
    r2.font.color.rgb = BLACK

    # Info line
    tb3 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.6))
    tf3 = tb3.text_frame
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
    r3 = tf3.paragraphs[0].add_run()
    r3.text = f"{ucbd_count} UCBDs  |  {req_count} Requirements  |  {const_count} Constants"
    r3.font.size = Pt(14)
    r3.font.name = FONT
    r3.font.color.rgb = BLACK
    r3.font.bold = True

    # Description from context
    description = context.get("system_description", "")
    if description:
        tb4 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(1.5))
        tf4 = tb4.text_frame
        tf4.word_wrap = True
        tf4.paragraphs[0].alignment = PP_ALIGN.CENTER
        r4 = tf4.paragraphs[0].add_run()
        r4.text = description
        r4.font.size = Pt(12)
        r4.font.name = FONT
        r4.font.color.rgb = BLACK


def add_summary_slide(prs, ucbd_groups: list[tuple[str, int]], constants: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.4))
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "Requirements Summary"
    r.font.size = Pt(20)
    r.font.name = FONT
    r.font.color.rgb = BLACK
    r.font.bold = True

    # Left table: requirements per UCBD
    rows_data = [("Use Case", "Req Count")]
    rows_data.extend(ucbd_groups)
    total_reqs = sum(c for _, c in ucbd_groups)
    rows_data.append(("TOTAL", str(total_reqs)))

    rows, cols = len(rows_data), 2
    tbl = slide.shapes.add_table(rows, cols, Inches(0.7), Inches(0.8), Inches(5.5), Inches(5.5)).table
    tbl.columns[0].width = Inches(4.0)
    tbl.columns[1].width = Inches(1.5)

    for r_idx, row_data in enumerate(rows_data):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                p.font.color.rgb = BLACK
                if r_idx == 0 or r_idx == len(rows_data) - 1:
                    p.font.bold = True

    # Right: constants list
    tb2 = slide.shapes.add_textbox(Inches(7.0), Inches(0.9), Inches(5.5), Inches(5.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    header_run = tf2.paragraphs[0].add_run()
    header_run.text = f"Design Constants ({len(constants)}):"
    header_run.font.size = Pt(11)
    header_run.font.name = FONT
    header_run.font.color.rgb = BLACK
    header_run.font.bold = True

    # Blank line
    tf2.add_paragraph()

    for const in constants:
        name = const.get("constant", "")
        value = const.get("value", "")
        units = const.get("units", "")
        status = const.get("estimate_final", "")
        line = f"{name} = {value} {units}".strip()
        if status:
            line += f"  ({status})"
        p = tf2.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9)
        run.font.name = FONT
        run.font.color.rgb = BLACK


# ── Main orchestration ────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--project", default=".", help="Project folder containing the Module 2 JSONs.")
    ap.add_argument("--output", default=None, help="Output pptx path. Defaults to <project>/<system>_UCBD.pptx.")
    args = ap.parse_args()

    project_dir = Path(args.project).resolve()

    context = load_context(project_dir)
    ucbds = load_ucbds(project_dir)
    req_bundle = load_requirements(project_dir)
    const_bundle = load_constants(project_dir)

    all_requirements = req_bundle.get("requirements_table", [])
    all_constants = const_bundle.get("constants_table", [])
    known_constants = {c.get("constant") for c in all_constants if c.get("constant")}

    # Presentation setup
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    # Title slide
    add_title_slide(prs, context, len(ucbds), len(all_requirements), len(all_constants))

    # One slide per UCBD
    ucbd_groups: list[tuple[str, int]] = []
    for ucbd in ucbds:
        ucbd_id = ucbd_id_from_instance(ucbd)
        metadata = build_metadata(ucbd)
        steps = build_steps(ucbd)
        requirements = build_requirements_for_ucbd(ucbd_id, all_requirements, known_constants)

        use_case_name = ucbd.get("use_case_name", ucbd_id)
        flow_title = f"UCBD: {ucbd_id} — {use_case_name}"

        if can_combine(len(steps), len(requirements)):
            add_combined_slide(
                prs,
                title=f"{flow_title} ({len(requirements)} requirements)",
                metadata=metadata,
                steps=steps,
                requirements=requirements,
            )
        else:
            add_split_slides(
                prs,
                flow_title=flow_title,
                req_title=f"{ucbd_id} Requirements — {use_case_name} ({len(requirements)})",
                metadata=metadata,
                steps=steps,
                requirements=requirements,
            )

        ucbd_groups.append((f"{ucbd_id} — {use_case_name}", len(requirements)))

    # Summary slide
    add_summary_slide(prs, ucbd_groups, all_constants)

    # Save
    if args.output:
        out_path = Path(args.output).resolve()
    else:
        system_slug = re.sub(r"[^A-Za-z0-9]+", "_", context.get("system_name", "System")).strip("_")
        out_path = project_dir / f"{system_slug}_UCBD.pptx"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"  UCBDs: {len(ucbds)}")
    print(f"  Requirements: {len(all_requirements)}")
    print(f"  Constants: {len(all_constants)}")


if __name__ == "__main__":
    main()
