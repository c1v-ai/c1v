#!/usr/bin/env python3
"""
Generate Cornell-standard Context Diagram for Team Heat Guard (THG)
Following CESYS521 formatting rules:
  - Black and white only (no color)
  - Square corners (never rounded)
  - All boxes the SAME size (including system box)
  - Names CAPITALIZED, same font and font size
  - Dashed border = system boundary
  - 18 external elements (8-20 range per CESYS521)
  - Interaction labels as verb phrases (not data attribute lists)
  - Two labels per line: near-element and near-system
  - External-to-external connections where they exist
  - Layout: top(4 actors) + left(5 environmental) + right(5 regulatory) + bottom(4 systems)

Produces two slides:
  Slide 1: Basic  — boxes and lines only
  Slide 2: Detailed — boxes, lines, interaction labels, ext-to-ext connections

Usage:
  pip install python-pptx
  python generate_context_diagram_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ─── CONSTANTS ───────────────────────────────────────────────────

FONT_NAME = "Arial"
BOX_FONT = Pt(9)
LABEL_FONT = Pt(6)
TITLE_FONT = Pt(20)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(130, 130, 130)
LINE_W = Pt(1.2)
DASH_W = Pt(2.0)
EXT_LINE_W = Pt(0.8)

SW = 13.333  # slide width inches
SH = 7.5     # slide height inches
CX = SW / 2
CY = SH / 2

BW = 1.35    # box width  — ALL boxes same size (Cornell requirement)
BH = 0.55    # box height


# ─── DRAWING HELPERS ─────────────────────────────────────────────

def add_box(slide, left, top, text, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(BW), Inches(BH)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = LINE_W
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text.upper()
    run.font.size = BOX_FONT
    run.font.name = FONT_NAME
    run.font.color.rgb = BLACK
    run.font.bold = bold
    return shape


def add_dashed_boundary(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.background()
    shape.line.color.rgb = BLACK
    shape.line.width = DASH_W
    shape.line.dash_style = 4
    return shape


def edge_point(cx, cy, hw, hh, tx, ty):
    """Point on box edge nearest to target (tx, ty).
    cx,cy = box center; hw,hh = half-width, half-height."""
    dx, dy = tx - cx, ty - cy
    if dx == 0 and dy == 0:
        return cx, cy
    hits = []
    if dx != 0:
        t = hw / abs(dx)
        py = cy + t * dy
        if abs(py - cy) <= hh + 0.01:
            hits.append((cx + hw * (1 if dx > 0 else -1), py))
    if dy != 0:
        t = hh / abs(dy)
        px = cx + t * dx
        if abs(px - cx) <= hw + 0.01:
            hits.append((px, cy + hh * (1 if dy > 0 else -1)))
    if not hits:
        return cx, cy
    return min(hits, key=lambda p: (p[0] - tx) ** 2 + (p[1] - ty) ** 2)


def add_line(slide, x1, y1, x2, y2, width=LINE_W, color=BLACK):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = width
    return c


def connect(slide, el, et, sl, st):
    """Straight line from element box edge to system box edge."""
    ecx, ecy = el + BW / 2, et + BH / 2
    scx, scy = sl + BW / 2, st + BH / 2
    ex, ey = edge_point(ecx, ecy, BW / 2, BH / 2, scx, scy)
    sx, sy = edge_point(scx, scy, BW / 2, BH / 2, ecx, ecy)
    add_line(slide, ex, ey, sx, sy)


def add_label(slide, x, y, text, font_size=LABEL_FONT, max_w=1.4, color=BLACK):
    w, h = max_w, 0.48
    tb = slide.shapes.add_textbox(Inches(x - w / 2), Inches(y - h / 2), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.name = FONT_NAME
    run.font.color.rgb = color


# ─── LAYOUT ──────────────────────────────────────────────────────
# 18 elements arranged: 4 top, 5 left, 5 right, 4 bottom

sys_l = CX - BW / 2   # system box left
sys_t = CY - BH / 2   # system box top

# Dashed boundary
BP = 0.55
bnd_l = sys_l - BP
bnd_t = sys_t - BP
bnd_w = BW + 2 * BP
bnd_h = BH + 2 * BP

# --- TOP ROW: 4 human actors ---
TY = 0.3
TX = [1.3, 4.2, 7.2, 10.3]
TOP = [
    ("WORKERS",           TX[0], TY),
    ("SAFETY\nMANAGERS",  TX[1], TY),
    ("SITE\nSUPERVISORS", TX[2], TY),
    ("EMPLOYER\nADMINS",  TX[3], TY),
]

# --- LEFT COLUMN: 5 environmental / physical ---
LX = 0.15
LY = [1.5, 2.65, 3.75, 4.9, 6.05]
LEFT = [
    ("OUTDOOR\nENVIRONMENT", LX, LY[0]),
    ("INDOOR\nENVIRONMENT",  LX, LY[1]),
    ("WORK\nACTIVITIES",     LX, LY[2]),
    ("CLOTHING /\nPPE",      LX, LY[3]),
    ("MOBILE\nDEVICES",      LX, LY[4]),
]

# --- RIGHT COLUMN: 5 regulatory / compliance ---
RX = SW - BW - 0.15
RY = [1.5, 2.65, 3.75, 4.9, 6.05]
RIGHT = [
    ("OSHA\nREGULATIONS",     RX, RY[0]),
    ("STATE HEAT\nSTANDARDS", RX, RY[1]),
    ("NIOSH /\nACGIH",        RX, RY[2]),
    ("AUDIT\nSYSTEMS",        RX, RY[3]),
    ("WORKERS\nCOMP",         RX, RY[4]),
]

# --- BOTTOM ROW: 4 systems / infrastructure ---
BY = 6.6
BX = [1.3, 4.2, 7.2, 10.3]
BOTTOM = [
    ("WEATHER\nAPI",        BX[0], BY),
    ("SSO /\nIDENTITY",     BX[1], BY),
    ("WEB\nBROWSERS",       BX[2], BY),
    ("EMERGENCY\nSERVICES", BX[3], BY),
]

ALL = TOP + LEFT + RIGHT + BOTTOM  # 18 total


# ─── INTERACTION LABELS ──────────────────────────────────────────
# (near_element_label, near_system_label)
# Verb-phrase style per CESYS521 — describes what each entity DOES, not data attributes.
# Empty string = one-directional, no label on that end.

LABELS = {
    # --- Actors (bidirectional) ---
    "WORKERS":           ("enters biometric data into,\nselects activity in",
                          "predicts safe duration for,\nsends risk alerts to"),
    "SAFETY\nMANAGERS":  ("monitors team risk via,\nconfigures thresholds in",
                          "displays risk grid for,\ngenerates reports for"),
    "SITE\nSUPERVISORS": ("acknowledges escalations in",
                          "escalates unacknowledged\nalerts to"),
    "EMPLOYER\nADMINS":  ("provisions accounts in,\nconfigures sites in",
                          "exports compliance\nreports for"),

    # --- Environmental (inbound) ---
    "OUTDOOR\nENVIRONMENT": ("solar radiation, heat,\nhumidity affect",     ""),
    "INDOOR\nENVIRONMENT":  ("enclosure heat, limited\nventilation affect", ""),
    "WORK\nACTIVITIES":     ("metabolic rate from\nphysical labor affects", ""),
    "CLOTHING /\nPPE":      ("insulation, vapor\nresistance affect",        ""),
    "MOBILE\nDEVICES":      ("workers access\nthrough",                     "delivers app and\nalerts through"),

    # --- Regulatory (mostly inbound) ---
    "OSHA\nREGULATIONS":     ("compliance requirements\ngovern",             "demonstrates\ncompliance to"),
    "STATE HEAT\nSTANDARDS": ("state-specific thresholds\nconstrain",        ""),
    "NIOSH /\nACGIH":        ("accuracy benchmarks,\nTLV values constrain", ""),
    "AUDIT\nSYSTEMS":        ("",                                            "sends compliance\nlogs to"),
    "WORKERS\nCOMP":         ("",                                            "sends risk reduction\nevidence to"),

    # --- Systems / Infrastructure ---
    "WEATHER\nAPI":        ("provides environmental\nconditions to",  "requests by GPS\ncoordinates from"),
    "SSO /\nIDENTITY":     ("authenticates enterprise\nusers for",   ""),
    "WEB\nBROWSERS":       ("managers access\ndashboard through",    ""),
    "EMERGENCY\nSERVICES": ("",                                      "dispatches critical\nalerts with GPS to"),
}

# External-to-external connections (not through the system)
EXT_EXT = [
    ("OUTDOOR\nENVIRONMENT", "WORKERS",          "physically exposes\nto heat stress"),
    ("OSHA\nREGULATIONS",    "EMPLOYER\nADMINS", "imposes compliance\nobligations on"),
]


def find_pos(name):
    for n, x, y in ALL:
        if n == name:
            return x, y
    return None, None


# ─── SLIDE BUILDER ───────────────────────────────────────────────

def build_slide(prs, title, detailed=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.01), Inches(12.5), Inches(0.3))
    tf = tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = TITLE_FONT
    run.font.name = FONT_NAME
    run.font.color.rgb = BLACK
    run.font.bold = True

    # Dashed system boundary
    add_dashed_boundary(slide, bnd_l, bnd_t, bnd_w, bnd_h)

    # System box (same size as element boxes)
    add_box(slide, sys_l, sys_t, "The System", bold=True)

    # Element boxes + connector lines
    scx = sys_l + BW / 2
    scy = sys_t + BH / 2

    for name, left, top in ALL:
        add_box(slide, left, top, name)
        connect(slide, left, top, sys_l, sys_t)

        if detailed and name in LABELS:
            from_lbl, to_lbl = LABELS[name]
            ecx = left + BW / 2
            ecy = top + BH / 2

            if from_lbl:
                lx = ecx + (scx - ecx) * 0.22
                ly = ecy + (scy - ecy) * 0.22
                add_label(slide, lx, ly, from_lbl)
            if to_lbl:
                lx = ecx + (scx - ecx) * 0.72
                ly = ecy + (scy - ecy) * 0.72
                add_label(slide, lx, ly, to_lbl)

    # External-to-external connections (detailed only)
    if detailed:
        for from_name, to_name, label in EXT_EXT:
            fx, fy = find_pos(from_name)
            tx, ty = find_pos(to_name)
            if fx is None or tx is None:
                continue
            fcx, fcy = fx + BW / 2, fy + BH / 2
            tcx, tcy = tx + BW / 2, ty + BH / 2
            fex, fey = edge_point(fcx, fcy, BW / 2, BH / 2, tcx, tcy)
            tex, tey = edge_point(tcx, tcy, BW / 2, BH / 2, fcx, fcy)
            add_line(slide, fex, fey, tex, tey, width=EXT_LINE_W, color=GRAY)
            mx = (fex + tex) / 2
            my = (fey + tey) / 2
            add_label(slide, mx, my, label, color=GRAY, max_w=1.2)

    return slide


# ─── MAIN ────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    build_slide(prs, "Context Diagram", detailed=False)
    build_slide(prs, "Context Diagram - Detailed", detailed=True)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "THG_Context_Diagram.pptx")
    prs.save(out_path)

    print(f"Saved: {out_path}")
    print(f"Elements: {len(ALL)}")
    print("Layout: top(4 actors) + left(5 environ) + right(5 regulatory) + bottom(4 systems)")
    print(f"Ext-to-ext: {len(EXT_EXT)}")
