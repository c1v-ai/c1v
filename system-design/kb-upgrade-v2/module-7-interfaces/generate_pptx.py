#!/usr/bin/env python3
"""
Generate c1v Module 6 — Defining Interfaces PowerPoint decks.

Produces THREE pptx files in the same directory as this script:
  1. n2_chart.pptx            — 14×14 N2 matrix (2 slides: grid + cell detail)
  2. data_flow_diagram.pptx   — Top-level DFD (2 slides: overview + legend/detail)
  3. sequence_diagrams.pptx   — 4 sequence diagrams (4 slides, one per SEQ)

Style: c1v brand (Firefly / Porcelain / Tangerine / Danube), 16:9, Calibri.
Source data: n2_chart.json, data_flow_diagram.mmd, sequence_diagrams.mmd,
             interface_matrix.json (for criticality).

Run:
    python3 system-design/module-7-interfaces/generate_pptx.py
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ═══════════════════════════════════════════════════════
# BRAND + DIMENSIONS
# ═══════════════════════════════════════════════════════

# 16:9 — Cornell-wide (30 × 16.875)
SW = 30.0
SH = 16.875

FONT = "Calibri"

# c1v brand colors
FIREFLY   = RGBColor(0x0B, 0x2C, 0x29)  # dark
PORCELAIN = RGBColor(0xFB, 0xFC, 0xFC)  # light bg
TANGERINE = RGBColor(0xF1, 0x8F, 0x01)  # accent / critical
DANUBE    = RGBColor(0x59, 0x98, 0xC5)  # secondary blue / normal

# Utility colors (derived)
BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY_L  = RGBColor(0xE8, 0xE8, 0xE8)
GREY_M  = RGBColor(0xBD, 0xBD, 0xBD)
GREY_D  = RGBColor(0x70, 0x70, 0x70)
# Shades for N2 fill by criticality / type
CRIT_FILL = RGBColor(0xF1, 0x8F, 0x01)   # TANGERINE — critical
HIGH_FILL = RGBColor(0xFA, 0xC8, 0x6E)   # TANGERINE 40%
MED_FILL  = RGBColor(0x5C, 0x9E, 0xC7)   # DANUBE
LOW_FILL  = RGBColor(0xBD, 0xD6, 0xE8)   # DANUBE 40%
DIAG_FILL = FIREFLY

# Category colors for DFD boxes (internal subsystems)
CAT_UI       = RGBColor(0x5C, 0x9E, 0xC7)  # Danube
CAT_PLATFORM = RGBColor(0x0B, 0x2C, 0x29)  # Firefly
CAT_ENGINE   = RGBColor(0xF1, 0x8F, 0x01)  # Tangerine
CAT_DATA     = RGBColor(0xA6, 0xB2, 0xB0)  # neutral
CAT_EXT_SS   = RGBColor(0xE8, 0xC8, 0x8B)  # probe SDK — border-crossing

# Line widths
LW_THIN   = Pt(1.0)
LW_NORMAL = Pt(1.5)
LW_HEAVY  = Pt(2.25)

# Paths
HERE = os.path.dirname(os.path.abspath(__file__))
N2_JSON   = os.path.join(HERE, "n2_chart.json")
DFD_MMD   = os.path.join(HERE, "data_flow_diagram.mmd")
SEQ_MMD   = os.path.join(HERE, "sequence_diagrams.mmd")
IFM_JSON  = os.path.join(HERE, "interface_matrix.json")


# ═══════════════════════════════════════════════════════
# SHARED HELPERS
# ═══════════════════════════════════════════════════════

def new_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_anchor_mid(shape):
    tb = shape.text_frame._txBody
    bodyPr = tb.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.15),
                                  Inches(SW - 0.7), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(24); r.font.name = FONT; r.font.bold = True
    r.font.color.rgb = FIREFLY
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(13); r2.font.name = FONT; r2.font.italic = True
        r2.font.color.rgb = GREY_D


def add_box(slide, x, y, w, h, text, *,
            fill=WHITE, line=FIREFLY, line_w=LW_NORMAL,
            font_size=11, text_color=FIREFLY, bold=False, dashed=False,
            align=PP_ALIGN.CENTER):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = line_w
    if dashed:
        s.line.dash_style = 2
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    # clear the default paragraph
    tf.text = ""
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(font_size); r.font.name = FONT
        r.font.color.rgb = text_color
        r.font.bold = bold
    set_anchor_mid(s)
    return s


def add_textbox(slide, x, y, w, h, text, *,
                font_size=10, color=FIREFLY, bold=False, italic=False,
                align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.text = ""
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(font_size); r.font.name = FONT
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb


def _set_arrowhead(connector, tail=True, head=False):
    spPr = connector._element.find(qn('p:spPr'))
    if spPr is None:
        return
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    if tail:
        tEl = etree.SubElement(ln, qn('a:tailEnd'))
        tEl.set('type', 'triangle'); tEl.set('w', 'med'); tEl.set('len', 'med')
    if head:
        hEl = etree.SubElement(ln, qn('a:headEnd'))
        hEl.set('type', 'triangle'); hEl.set('w', 'med'); hEl.set('len', 'med')


def add_segment(slide, x1, y1, x2, y2, *,
                color=FIREFLY, width=LW_NORMAL, dashed=False,
                arrow_tail=False, arrow_head=False):
    cxn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = color
    cxn.line.width = width
    if dashed:
        cxn.line.dash_style = 2
    if arrow_tail or arrow_head:
        _set_arrowhead(cxn, tail=arrow_tail, head=arrow_head)
    return cxn


def add_arrow(slide, points, *,
              color=FIREFLY, width=LW_NORMAL, dashed=False, label=None,
              label_offset=0.0):
    for i in range(len(points) - 1):
        is_last = (i == len(points) - 2)
        add_segment(slide,
                    points[i][0], points[i][1],
                    points[i + 1][0], points[i + 1][1],
                    color=color, width=width, dashed=dashed,
                    arrow_tail=is_last)
    if label:
        mid_i = len(points) // 2
        if mid_i > 0:
            x1, y1 = points[mid_i - 1]
            x2, y2 = points[mid_i]
            lx = (x1 + x2) / 2
            ly = (y1 + y2) / 2 - 0.15 + label_offset
            add_textbox(slide, lx - 1.25, ly - 0.12, 2.5, 0.26, label,
                        font_size=8, italic=True, color=FIREFLY,
                        align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# DATA LOAD
# ═══════════════════════════════════════════════════════

def load_n2():
    with open(N2_JSON) as f:
        return json.load(f)


def load_ifm():
    with open(IFM_JSON) as f:
        return json.load(f)


def build_criticality_map(ifm):
    """IF-id → criticality string ('critical' | 'high' | 'med' | 'low')."""
    out = {}
    for r in ifm.get("interface_registry", []):
        out[r["id"]] = r.get("criticality", "med")
    for r in ifm.get("external_interface_registry", []):
        out[r["id"]] = r.get("criticality", "med")
    return out


# ═══════════════════════════════════════════════════════
# DECK 1 — N2 CHART
# ═══════════════════════════════════════════════════════

def build_n2_chart_pptx():
    n2  = load_n2()
    ifm = load_ifm()
    crit = build_criticality_map(ifm)

    prs = new_presentation()
    build_n2_grid_slide(prs, n2, crit)
    build_n2_detail_slide(prs, n2, crit)
    out = os.path.join(HERE, "n2_chart.pptx")
    prs.save(out)
    return out, len(prs.slides)


def crit_fill(c):
    return {"critical": CRIT_FILL, "high": HIGH_FILL,
            "med": MED_FILL, "low": LOW_FILL}.get(c, MED_FILL)


def crit_text_color(c):
    # light background for med/low so use dark text; critical/high colored fill contrast
    if c in ("critical", "high"):
        return FIREFLY
    return WHITE


def build_n2_grid_slide(prs, n2, crit):
    slide = add_slide(prs)
    add_title(slide,
        "N2 Chart — c1v Dual-Mode Platform",
        sub="14 subsystems on the diagonal · upper triangle = forward flows · lower triangle = feedback · 32 internal interfaces (IF-01..IF-32)")

    subs = n2["subsystems"]
    n = len(subs)

    # Build interface lookup: (from,to) → interface dict
    ifm_lookup = {}
    for iface in n2["interfaces"]:
        ifm_lookup[(iface["from"], iface["to"])] = iface

    # Grid geometry — leave title band + side column for SS labels + bottom legend
    grid_x0 = 3.6
    grid_y0 = 1.6
    grid_w  = 22.0
    grid_h  = 13.2
    cell_w  = grid_w / n
    cell_h  = grid_h / n

    # Row labels (left) + column labels (top) — SS id + short name
    for i, ss in enumerate(subs):
        # Left label (row header)
        y = grid_y0 + i * cell_h
        add_box(slide,
                grid_x0 - 3.5, y, 3.5, cell_h,
                f"{ss['id']}  {ss['name']}",
                fill=FIREFLY, line=FIREFLY, line_w=LW_THIN,
                text_color=PORCELAIN, bold=True, font_size=9,
                align=PP_ALIGN.LEFT)
        # Top label (column header) — rotated-ish we emulate with short-stack text
        cx = grid_x0 + i * cell_w
        add_box(slide,
                cx, grid_y0 - 1.0, cell_w, 1.0,
                f"{ss['id']}\n{ss['name']}",
                fill=FIREFLY, line=FIREFLY, line_w=LW_THIN,
                text_color=PORCELAIN, bold=True, font_size=7)

    # Grid cells
    for i in range(n):
        for j in range(n):
            x = grid_x0 + j * cell_w
            y = grid_y0 + i * cell_h
            if i == j:
                # Diagonal — subsystem node (solid Firefly, show SS id only)
                add_box(slide, x, y, cell_w, cell_h,
                        subs[i]["id"],
                        fill=DIAG_FILL, line=PORCELAIN, line_w=LW_NORMAL,
                        text_color=TANGERINE, bold=True, font_size=12)
            else:
                src = subs[i]["id"]
                tgt = subs[j]["id"]
                iface = ifm_lookup.get((src, tgt))
                if iface is None:
                    # Empty cell — very light
                    add_box(slide, x, y, cell_w, cell_h, "",
                            fill=PORCELAIN, line=GREY_M, line_w=LW_THIN)
                else:
                    c = crit.get(iface["id"], "med")
                    fill = crit_fill(c)
                    tcol = crit_text_color(c)
                    # Label: IF-id + 1st EC
                    ec_str = ",".join(iface.get("driven_by_ec", [])[:2])
                    label = f"{iface['id']}\n{ec_str}"
                    add_box(slide, x, y, cell_w, cell_h, label,
                            fill=fill, line=FIREFLY,
                            line_w=(LW_HEAVY if c == "critical" else
                                    LW_NORMAL if c == "high" else LW_THIN),
                            text_color=tcol, bold=(c in ("critical","high")),
                            font_size=7)

    # Triangle guide (upper = forward, lower = feedback)
    add_textbox(slide, grid_x0 + grid_w * 0.55, grid_y0 + 0.08,
                grid_w * 0.4, 0.32,
                "▲ UPPER triangle — forward flow (request)",
                font_size=9, italic=True, color=PORCELAIN, bold=True)
    add_textbox(slide, grid_x0 + 0.1, grid_y0 + grid_h - 0.38,
                grid_w * 0.5, 0.32,
                "▼ LOWER triangle — feedback (response/return)",
                font_size=9, italic=True, color=PORCELAIN, bold=True)

    # Legend bottom-left
    lx, ly = 0.4, 15.2
    add_textbox(slide, lx, ly, 6, 0.3, "Criticality (cell fill)",
                font_size=10, bold=True, color=FIREFLY)
    swatches = [("critical", CRIT_FILL), ("high", HIGH_FILL),
                ("med", MED_FILL),       ("low", LOW_FILL)]
    for i, (name, col) in enumerate(swatches):
        sx = lx + i * 1.6
        add_box(slide, sx, ly + 0.32, 0.35, 0.28, "",
                fill=col, line=FIREFLY, line_w=LW_THIN)
        add_textbox(slide, sx + 0.42, ly + 0.32, 1.15, 0.28, name,
                    font_size=9, color=FIREFLY)

    # Right side mini legend (subsystem totals)
    rx, ry = 26.0, 15.2
    add_textbox(slide, rx, ry, 3.6, 0.3, "Grid totals",
                font_size=10, bold=True, color=FIREFLY)
    add_textbox(slide, rx, ry + 0.32, 3.6, 0.9,
                f"subsystems: {n}\ninternal IFs: {len(n2['interfaces'])}\nexternal IFs: {len(n2['external_interfaces'])}",
                font_size=9, color=GREY_D)


def build_n2_detail_slide(prs, n2, crit):
    slide = add_slide(prs)
    add_title(slide,
        "N2 Chart — Interface Cell Detail",
        sub="Every non-zero cell · from → to · data summary · EC budget drivers")

    # Two-column table of all 32 interfaces
    ifaces = n2["interfaces"]
    col_h = 14.6  # usable area under title
    y0 = 1.35
    rows_per_col = (len(ifaces) + 1) // 2  # 16 + 16
    row_h = col_h / rows_per_col

    cols = [
        {"x0": 0.4,  "w": 14.5},
        {"x0": 15.2, "w": 14.5},
    ]

    # Header rows
    for c in cols:
        hx = c["x0"]
        add_box(slide, hx, y0, 1.2, row_h,
                "IF", fill=FIREFLY, line=FIREFLY,
                text_color=PORCELAIN, bold=True, font_size=9)
        add_box(slide, hx + 1.2, y0, 2.0, row_h,
                "from → to", fill=FIREFLY, line=FIREFLY,
                text_color=PORCELAIN, bold=True, font_size=9)
        add_box(slide, hx + 3.2, y0, 9.5, row_h,
                "data / EC drivers", fill=FIREFLY, line=FIREFLY,
                text_color=PORCELAIN, bold=True, font_size=9)
        add_box(slide, hx + 12.7, y0, 1.8, row_h,
                "crit", fill=FIREFLY, line=FIREFLY,
                text_color=PORCELAIN, bold=True, font_size=9)

    for idx, iface in enumerate(ifaces):
        col = cols[0] if idx < rows_per_col else cols[1]
        row_idx = (idx % rows_per_col) + 1
        y = y0 + row_idx * row_h
        c = crit.get(iface["id"], "med")
        fill = crit_fill(c) if c in ("critical", "high") else PORCELAIN
        tcol = FIREFLY

        add_box(slide, col["x0"], y, 1.2, row_h,
                iface["id"],
                fill=PORCELAIN, line=GREY_M, line_w=LW_THIN,
                text_color=FIREFLY, bold=True, font_size=8)
        add_box(slide, col["x0"] + 1.2, y, 2.0, row_h,
                f"{iface['from']} → {iface['to']}",
                fill=PORCELAIN, line=GREY_M, line_w=LW_THIN,
                text_color=FIREFLY, font_size=8)
        ec_drivers = ",".join(iface.get("driven_by_ec", []))
        data = (iface.get("name") or "") + " · " + ec_drivers
        add_box(slide, col["x0"] + 3.2, y, 9.5, row_h,
                data,
                fill=PORCELAIN, line=GREY_M, line_w=LW_THIN,
                text_color=FIREFLY, font_size=7, align=PP_ALIGN.LEFT)
        add_box(slide, col["x0"] + 12.7, y, 1.8, row_h,
                c,
                fill=fill, line=GREY_M, line_w=LW_THIN,
                text_color=tcol, bold=(c in ("critical", "high")),
                font_size=8)


# ═══════════════════════════════════════════════════════
# DECK 2 — DATA FLOW DIAGRAM
# ═══════════════════════════════════════════════════════

def build_dfd_pptx():
    n2  = load_n2()
    ifm = load_ifm()
    crit = build_criticality_map(ifm)

    prs = new_presentation()
    build_dfd_overview_slide(prs, n2, crit)
    build_dfd_legend_slide(prs, n2, ifm, crit)
    out = os.path.join(HERE, "data_flow_diagram.pptx")
    prs.save(out)
    return out, len(prs.slides)


# Subsystem positions (x,y,w,h) arranged in functional bands.
# Layout:
#   Band A (UI edge, left col)            : SS1
#   Band B (platform / session, col 2)    : SS2, SS13
#   Band C (engine core, center)          : SS3, SS5, SS4, SS6, SS7, SS8
#   Band D (probe/metrics, lower)         : SS9, SS10, SS11, SS12
#   Band E (audit, right edge)            : SS14
#
# External actors line the periphery (top band + right band + left band).

# Internal subsystem positions (center_x, center_y, w, h)
INTERNAL_POS = {
    "SS1":  (3.5,  9.0, 2.6, 1.2),
    "SS2":  (7.0,  7.0, 2.6, 1.2),
    "SS13": (7.0, 11.0, 2.6, 1.2),
    "SS3":  (11.3, 7.0, 2.8, 1.4),
    "SS5":  (11.3, 9.2, 2.6, 1.1),
    "SS4":  (11.3, 11.4,2.6, 1.2),
    "SS6":  (15.3, 7.0, 2.6, 1.2),
    "SS7":  (15.3, 9.2, 2.6, 1.1),
    "SS8":  (15.3, 11.4,2.6, 1.2),
    "SS14": (19.0, 11.4,2.6, 1.2),
    "SS9":  (19.0, 9.2, 2.6, 1.2),
    "SS10": (22.4, 9.2, 2.6, 1.2),
    "SS11": (25.6, 9.2, 2.6, 1.1),
    "SS12": (25.6, 7.0, 2.6, 1.2),
}

# External actor positions (center_x, center_y, w, h) — periphery
EXTERNAL_POS = {
    # top band
    "FOUND":  (3.5,  3.4, 3.2, 1.1),
    "IDP":    (7.0,  3.4, 3.0, 1.1),
    "LLM":    (11.3, 3.4, 3.0, 1.1),
    "DOCS":   (15.3, 3.4, 3.0, 1.1),
    "IDE":    (19.0, 3.4, 3.0, 1.1),
    "VCS":    (22.4, 3.4, 3.0, 1.1),
    "CICD":   (25.6, 3.4, 3.0, 1.1),
    # bottom
    "CUST":   (19.0, 14.8, 3.4, 1.1),
    "CLOUD":  (11.3, 14.8, 3.2, 1.1),
    "COMPLY": (22.8, 14.8, 3.4, 1.1),
}

EXTERNAL_LABELS = {
    "FOUND":  "Founders / PMs /\nEngineers / Admins",
    "IDP":    "Identity Providers\n(OIDC / SAML)",
    "LLM":    "LLM Providers\n(Anthropic / OpenAI / local)",
    "DOCS":   "Documentation Sources\n(vendor specs, RFCs)",
    "IDE":    "IDE / CLI Clients\n(Claude Code / Cursor)",
    "VCS":    "Version Control\n(git / HTTPS)",
    "CICD":   "CI / CD Systems",
    "CUST":   "Customer Prod Systems\n+ Observability Tools",
    "CLOUD":  "Cloud Providers\n(AWS / GCP / Azure)",
    "COMPLY": "Compliance Frameworks\n(SOC2/HIPAA/GDPR)",
}


def category_for(ss_id):
    cat_map = {
        "SS1": "ui",
        "SS2": "platform", "SS13": "platform",
        "SS3": "engine", "SS4": "engine", "SS6": "engine",
        "SS8": "engine", "SS10": "engine", "SS12": "engine",
        "SS5": "data", "SS7": "data", "SS11": "data", "SS14": "data",
        "SS9": "ext_ss",
    }
    return cat_map.get(ss_id, "engine")


def fill_for_category(cat):
    return {"ui": CAT_UI, "platform": CAT_PLATFORM,
            "engine": CAT_ENGINE, "data": CAT_DATA,
            "ext_ss": CAT_EXT_SS}.get(cat, WHITE)


def text_for_category(cat):
    return {"ui": WHITE, "platform": WHITE,
            "engine": FIREFLY, "data": FIREFLY,
            "ext_ss": FIREFLY}.get(cat, FIREFLY)


def subsystem_edge(ss_id, side):
    cx, cy, w, h = INTERNAL_POS[ss_id]
    if side == "left":   return (cx - w/2, cy)
    if side == "right":  return (cx + w/2, cy)
    if side == "top":    return (cx, cy - h/2)
    if side == "bottom": return (cx, cy + h/2)
    return (cx, cy)


def external_edge(ext_id, side):
    cx, cy, w, h = EXTERNAL_POS[ext_id]
    if side == "left":   return (cx - w/2, cy)
    if side == "right":  return (cx + w/2, cy)
    if side == "top":    return (cx, cy - h/2)
    if side == "bottom": return (cx, cy + h/2)
    return (cx, cy)


def build_dfd_overview_slide(prs, n2, crit):
    slide = add_slide(prs)
    add_title(slide,
        "Top-Level Data Flow Diagram — c1v Dual-Mode Platform",
        sub="14 subsystems (solid) · 10 external actors (dashed) · every arrow labels data + EC budget · critical flows drawn heavier + Tangerine")

    # c1v boundary — large rounded rect around internal subsystems
    bx, by, bw, bh = 1.7, 5.5, 26.4, 7.6
    boundary = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(bx), Inches(by),
                                      Inches(bw), Inches(bh))
    boundary.fill.solid(); boundary.fill.fore_color.rgb = PORCELAIN
    boundary.line.color.rgb = FIREFLY; boundary.line.width = LW_HEAVY
    try:
        boundary.adjustments[0] = 0.04
    except Exception:
        pass
    add_textbox(slide, bx + 0.2, by + 0.05, 6, 0.35,
                "c1v Dual-Mode Platform (boundary)",
                font_size=10, bold=True, color=FIREFLY)

    # External actors (dashed) — draw first so subsystems render above
    for eid, (cx, cy, w, h) in EXTERNAL_POS.items():
        add_box(slide, cx - w/2, cy - h/2, w, h,
                EXTERNAL_LABELS[eid],
                fill=PORCELAIN, line=FIREFLY, line_w=LW_NORMAL, dashed=True,
                text_color=FIREFLY, bold=False, font_size=9)

    # Internal subsystems
    ss_by_id = {s["id"]: s for s in n2["subsystems"]}
    for ss_id, (cx, cy, w, h) in INTERNAL_POS.items():
        cat = category_for(ss_id)
        ss = ss_by_id[ss_id]
        # Subtitle snippets
        subtitle_map = {
            "SS2":  "\nmode: cloud / local",
            "SS3":  "\nLangGraph",
            "SS5":  "\n≥70% hit target (EC7)",
            "SS6":  "\nMermaid v1",
            "SS7":  "\nTTL 24h · 7d refresh",
            "SS8":  "\n≤5 MB (EC14)",
            "SS9":  "\ncustomer-side",
            "SS11": "\n60-min window",
            "SS13": "\n90d rotation (EC17)",
            "SS14": "\n90d retention (EC18)",
        }
        label = f"{ss_id}  {ss['name']}" + subtitle_map.get(ss_id, "")
        add_box(slide, cx - w/2, cy - h/2, w, h, label,
                fill=fill_for_category(cat),
                line=FIREFLY, line_w=LW_HEAVY,
                text_color=text_for_category(cat),
                bold=True, font_size=9)

    # Interface → arrow drawing helper
    def style_for(iface_id):
        c = crit.get(iface_id, "med")
        if c == "critical":
            return {"color": TANGERINE, "width": LW_HEAVY}
        if c == "high":
            return {"color": FIREFLY,  "width": LW_NORMAL}
        return {"color": GREY_D, "width": LW_THIN}

    # Internal → internal arrows (selected set — all 32 IFs)
    internal_arrows = [
        # (from_ss, from_side, to_ss, to_side, label, mid_points)
        ("SS1",  "right",  "SS2",  "left",  "IF-01 login+mode", []),
        ("SS2",  "left",   "SS1",  "right", "IF-02 SSE ≤50ms",  []),
        ("SS1",  "bottom", "SS6",  "top",   "IF-31 approve", [(3.5,13.0),(15.3,13.0)]),
        ("SS6",  "bottom", "SS1",  "bottom","IF-32 mermaid", [(15.3,13.5),(3.5,13.5)]),

        ("SS2",  "right",  "SS3",  "left",  "IF-03 dispatch (mode LOCKED)", []),
        ("SS3",  "left",   "SS2",  "right", "IF-04 run state",  []),
        ("SS2",  "right",  "SS4",  "left",  "IF-09 routing policy", [(8.5,11.4)]),

        ("SS3",  "bottom", "SS5",  "top",   "IF-05 lookup (EC7)", []),
        ("SS5",  "top",    "SS3",  "bottom","IF-06 hit/miss",     []),
        ("SS5",  "bottom", "SS4",  "top",   "IF-10 prefix reuse", []),
        ("SS3",  "right",  "SS4",  "top",   "IF-07 LLM call ≤8000 tok", [(12.7,7.0),(12.7,11.4)]),
        ("SS4",  "right",  "SS3",  "bottom","IF-08 completion",   [(13.5,11.4),(13.5,7.8)]),

        ("SS3",  "right",  "SS6",  "left",  "IF-11 gen spec (EC12/EC13)", []),
        ("SS6",  "left",   "SS3",  "right", "IF-12 draft+pending", []),

        ("SS6",  "bottom", "SS7",  "top",   "IF-13 citation check HARD floor (EC10)", []),
        ("SS7",  "top",    "SS6",  "bottom","IF-14 citations TTL=24h", []),

        ("SS3",  "right",  "SS8",  "left",  "IF-15 emit CLI (citation gate PASS)", [(14.5,7.0),(14.5,11.4)]),
        ("SS8",  "left",   "SS3",  "right", "IF-16 receipt+sig", [(14.2,11.4),(14.2,7.2)]),

        ("SS9",  "right",  "SS10", "left",  "IF-17 probe ≤6/min ≤4KB (EC1/3/4)", []),
        ("SS10", "left",   "SS9",  "right", "IF-18 backpressure", []),
        ("SS10", "right",  "SS11", "left",  "IF-19 aggregate 60-min (EC2)", []),
        ("SS11", "top",    "SS12", "bottom","IF-20 deviation pull", []),

        ("SS12", "left",   "SS7",  "right", "IF-21 citation lookup", []),
        ("SS7",  "right",  "SS12", "left",  "IF-22 cited refs",      [(18.0,9.5)]),
        ("SS12", "top",    "SS3",  "right", "IF-23 rec → next spec (F.6→F.2)",
                    [(25.6,5.8),(14.0,5.8)]),

        ("SS3",  "bottom", "SS14", "top",   "IF-24 audit orchestration (EC18 90d)",
                    [(11.3,12.6),(19.0,12.6)]),
        ("SS8",  "right",  "SS14", "left",  "IF-25 CLI emit audit", []),
        ("SS10", "bottom", "SS14", "bottom","IF-26 rejection audit",
                    [(22.4,13.4),(19.0,13.4)]),

        ("SS13", "top",    "SS2",  "bottom","IF-27 signed JWT (EC17)", []),
        ("SS2",  "bottom", "SS13", "top",   "IF-28 rotate trigger (EC17)", []),
        ("SS13", "right",  "SS4",  "left",  "IF-29 provider creds",  [(9.5,11.0),(9.5,11.4)]),
        ("SS13", "right",  "SS9",  "left",  "IF-30 READ-ONLY handles (write REJECT)",
                    [(8.5,11.0),(8.5,9.2)]),
    ]

    for entry in internal_arrows:
        ss_a, side_a, ss_b, side_b, label, mid = entry
        if_id = label.split()[0]
        style = style_for(if_id)
        p_start = subsystem_edge(ss_a, side_a)
        p_end   = subsystem_edge(ss_b, side_b)
        points = [p_start] + list(mid) + [p_end]
        add_arrow(slide, points, label=label,
                  color=style["color"], width=style["width"])

    # External → internal arrows
    ext_arrows = [
        # founders → SS1
        ("FOUND", "bottom", "SS1",  "top",    "IF-EXT-01 UI events SSE (EC16)", []),
        # IdP → SS13
        ("IDP",   "bottom", "SS13", "top",    "IF-EXT-02 OIDC/SAML read-only",
                    [(7.0,5.2)]),
        # SS4 → LLM
        ("SS4",   "top",    "LLM",  "bottom", "IF-EXT-03 prompts→completions (EC8/EC16)", []),
        # SS4 → Cloud
        ("SS4",   "bottom", "CLOUD","top",    "IF-EXT-09 service catalogs", []),
        # SS8 → VCS
        ("SS8",   "top",    "VCS",  "bottom", "IF-EXT-04 signed CLI bundle",
                    [(15.3,5.2),(22.4,5.2)]),
        # SS8 → IDE
        ("SS8",   "top",    "IDE",  "bottom", "IF-EXT-10 SKILL.md + CLAUDE.md",
                    [(15.3,5.0),(19.0,5.0)]),
        # SS8 → CICD
        ("SS8",   "top",    "CICD", "bottom", "IF-EXT-08 spec gate / hooks",
                    [(15.6,4.9),(25.6,4.9)]),
        # DOCS → SS7
        ("DOCS",  "bottom", "SS7",  "top",    "IF-EXT-06 RAG retrieval refresh 7d",
                    [(15.3,5.2)]),
        # CUST → SS9
        ("CUST",  "top",    "SS9",  "bottom", "IF-EXT-05 read-only handles ≤2% overhead (EC1/3)",
                    []),
        # SS14 → COMPLY
        ("SS14",  "bottom", "COMPLY","top",   "IF-EXT-07 compliance evidence (EC18)",
                    []),
    ]

    for src, side_a, tgt, side_b, label, mid in ext_arrows:
        if_id = label.split()[0]
        style = style_for(if_id)
        # Resolve endpoints
        p_start = (external_edge(src, side_a) if src in EXTERNAL_POS
                   else subsystem_edge(src, side_a))
        p_end   = (external_edge(tgt, side_b) if tgt in EXTERNAL_POS
                   else subsystem_edge(tgt, side_b))
        points = [p_start] + list(mid) + [p_end]
        add_arrow(slide, points, label=label,
                  color=style["color"], width=style["width"], dashed=True)

    # Compact legend bottom-right
    lx, ly = 0.4, 16.1
    add_textbox(slide, lx, ly, 3, 0.3, "Legend",
                font_size=10, bold=True, color=FIREFLY)
    add_textbox(slide, lx + 0.6, ly + 0.32, 28, 0.3,
                "critical flow = Tangerine heavy   ·   "
                "high = Firefly normal   ·   "
                "med/low = grey thin   ·   "
                "solid = internal   ·   "
                "dashed = external crossing",
                font_size=9, color=GREY_D)


def build_dfd_legend_slide(prs, n2, ifm, crit):
    slide = add_slide(prs)
    add_title(slide,
        "DFD — Category Legend & Critical-Path Summary",
        sub="Subsystem categories · external-actor treatment · critical interfaces prioritized")

    # Left: category swatches
    lx, ly = 0.6, 1.6
    add_textbox(slide, lx, ly, 8, 0.4,
                "Subsystem categories (box fill)",
                font_size=14, bold=True, color=FIREFLY)
    cats = [
        ("UI (user-facing)",         CAT_UI,       WHITE,  ["SS1"]),
        ("Platform (session/auth)",   CAT_PLATFORM, WHITE,  ["SS2","SS13"]),
        ("Engine (compute)",          CAT_ENGINE,  FIREFLY, ["SS3","SS4","SS6","SS8","SS10","SS12"]),
        ("Data (storage/cache)",      CAT_DATA,    FIREFLY, ["SS5","SS7","SS11","SS14"]),
        ("External SS (in customer)", CAT_EXT_SS,  FIREFLY, ["SS9"]),
    ]
    for i, (name, fill, tcol, ids) in enumerate(cats):
        sy = ly + 0.6 + i * 0.9
        add_box(slide, lx, sy, 0.8, 0.7, "",
                fill=fill, line=FIREFLY, line_w=LW_NORMAL)
        add_textbox(slide, lx + 1.0, sy, 9, 0.3, name,
                    font_size=12, bold=True, color=FIREFLY)
        add_textbox(slide, lx + 1.0, sy + 0.33, 9, 0.3,
                    ", ".join(ids), font_size=10, color=GREY_D)

    # External actor style
    ly2 = ly + 6.2
    add_textbox(slide, lx, ly2, 10, 0.4,
                "External actor style (outside c1v boundary)",
                font_size=14, bold=True, color=FIREFLY)
    add_box(slide, lx, ly2 + 0.5, 0.8, 0.7, "",
            fill=PORCELAIN, line=FIREFLY, line_w=LW_NORMAL, dashed=True)
    add_textbox(slide, lx + 1.0, ly2 + 0.5, 9, 0.3,
                "Dashed border, porcelain fill",
                font_size=12, bold=True, color=FIREFLY)
    add_textbox(slide, lx + 1.0, ly2 + 0.85, 9, 0.3,
                "10 external actors: founders, IdP, LLM, VCS, customer, docs, compliance, CI/CD, cloud, IDE.",
                font_size=10, color=GREY_D)

    # Right: critical interfaces list
    rx, ry = 14.5, 1.6
    add_textbox(slide, rx, ry, 14, 0.4,
                "Critical interfaces (heavier + Tangerine arrows)",
                font_size=14, bold=True, color=FIREFLY)

    registry = ifm.get("interface_registry", [])
    criticals = [r for r in registry if r.get("criticality") == "critical"]

    header_row_h = 0.4
    row_h = 0.5
    cols_w = [1.4, 2.4, 6.0, 4.0]  # IF / route / data / EC
    headers = ["IF", "from → to", "data", "EC"]
    cx = rx
    hy = ry + 0.55
    for i, h in enumerate(headers):
        add_box(slide, cx, hy, cols_w[i], header_row_h, h,
                fill=FIREFLY, line=FIREFLY, line_w=LW_THIN,
                text_color=PORCELAIN, bold=True, font_size=10)
        cx += cols_w[i]

    for ri, r in enumerate(criticals):
        y = hy + header_row_h + ri * row_h
        cx = rx
        add_box(slide, cx, y, cols_w[0], row_h, r["id"],
                fill=PORCELAIN, line=GREY_M, line_w=LW_THIN,
                text_color=FIREFLY, bold=True, font_size=9)
        cx += cols_w[0]
        add_box(slide, cx, y, cols_w[1], row_h,
                f"{r.get('source_ss','ext')} → {r.get('target_ss','ext')}",
                fill=PORCELAIN, line=GREY_M, line_w=LW_THIN,
                text_color=FIREFLY, font_size=9)
        cx += cols_w[1]
        add_box(slide, cx, y, cols_w[2], row_h,
                r.get("data_flow", ""),
                fill=PORCELAIN, line=GREY_M, line_w=LW_THIN,
                text_color=FIREFLY, font_size=8,
                align=PP_ALIGN.LEFT)
        cx += cols_w[2]
        ec = ", ".join(r.get("driven_by_ec", []))
        add_box(slide, cx, y, cols_w[3], row_h,
                ec,
                fill=CRIT_FILL, line=FIREFLY, line_w=LW_THIN,
                text_color=FIREFLY, bold=True, font_size=9)


# ═══════════════════════════════════════════════════════
# DECK 3 — SEQUENCE DIAGRAMS
# ═══════════════════════════════════════════════════════

def build_sequence_pptx():
    prs = new_presentation()
    build_seq_1(prs)
    build_seq_2(prs)
    build_seq_3(prs)
    build_seq_4(prs)
    out = os.path.join(HERE, "sequence_diagrams.pptx")
    prs.save(out)
    return out, len(prs.slides)


# Lifeline geometry helpers
def draw_lifelines(slide, participants, *, y_top=2.3, y_bottom=15.5,
                   head_h=0.9, head_w=2.4):
    """
    Draw participant headers + vertical dashed lifelines.
    participants: list[(id, display_text, category)]
        category ∈ 'internal' | 'external'
    Returns dict id → center_x.
    """
    n = len(participants)
    # Distribute evenly across usable width
    x_left = 1.0
    x_right = SW - 1.0
    span = x_right - x_left
    if n > 1:
        step = span / (n - 1)
    else:
        step = 0
    centers = {}
    for i, (pid, text, cat) in enumerate(participants):
        cx = x_left + i * step
        centers[pid] = cx
        # Head box
        fill = FIREFLY if cat == "internal" else PORCELAIN
        text_color = PORCELAIN if cat == "internal" else FIREFLY
        dashed = (cat != "internal")
        add_box(slide, cx - head_w/2, y_top, head_w, head_h, text,
                fill=fill, line=FIREFLY, line_w=LW_NORMAL, dashed=dashed,
                text_color=text_color, bold=True, font_size=10)
        # Lifeline
        add_segment(slide, cx, y_top + head_h, cx, y_bottom,
                    color=GREY_D, width=LW_THIN, dashed=True)
    return centers


def add_sequence_arrow(slide, centers, from_id, to_id, y, label, *,
                       kind="sync",  # 'sync' | 'return' | 'self'
                       heavy=False):
    """Horizontal arrow between two lifeline centers at height y."""
    x_from = centers[from_id]
    x_to   = centers[to_id]
    dashed = (kind == "return")
    color  = TANGERINE if heavy else FIREFLY
    width  = LW_HEAVY if heavy else LW_NORMAL

    if kind == "self" or from_id == to_id:
        # Self-loop: out to the right, down, back
        off = 0.8
        add_segment(slide, x_from, y, x_from + off, y,
                    color=color, width=width, dashed=dashed)
        add_segment(slide, x_from + off, y, x_from + off, y + 0.35,
                    color=color, width=width, dashed=dashed)
        add_segment(slide, x_from + off, y + 0.35, x_from, y + 0.35,
                    color=color, width=width, dashed=dashed,
                    arrow_tail=True)
        # label
        add_textbox(slide, x_from + 0.9, y - 0.02, 7, 0.35, label,
                    font_size=9, color=FIREFLY, italic=True)
    else:
        add_segment(slide, x_from, y, x_to, y,
                    color=color, width=width, dashed=dashed,
                    arrow_tail=True)
        # Label above the arrow
        lx = min(x_from, x_to) + 0.1
        lw = abs(x_to - x_from) - 0.2
        if lw < 2.0:
            lw = 2.0
        add_textbox(slide, lx, y - 0.28, lw, 0.26, label,
                    font_size=9, color=FIREFLY, italic=True,
                    align=PP_ALIGN.CENTER)


def add_sequence_note(slide, centers, over_ids, y, text, *,
                      height=0.5):
    """Yellow note spanning over_ids at height y."""
    xs = [centers[p] for p in over_ids]
    x_from = min(xs) - 1.0
    x_to   = max(xs) + 1.0
    w = x_to - x_from
    add_box(slide, x_from, y, w, height, text,
            fill=RGBColor(0xFF, 0xF3, 0xCD), line=TANGERINE, line_w=LW_NORMAL,
            text_color=FIREFLY, bold=False, font_size=9,
            align=PP_ALIGN.CENTER)


def add_activation(slide, cx, y_start, y_end):
    """Thin activation bar centered on a lifeline."""
    w = 0.18
    add_box(slide, cx - w/2, y_start, w, y_end - y_start, "",
            fill=TANGERINE, line=FIREFLY, line_w=LW_THIN)


def add_alt_frame(slide, centers, over_ids, y_start, y_end, label):
    """Draw a labelled frame rectangle (no fill) to indicate alt/loop/par."""
    xs = [centers[p] for p in over_ids]
    x_from = min(xs) - 1.1
    x_to   = max(xs) + 1.1
    # frame rect
    w = x_to - x_from
    h = y_end - y_start
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(x_from), Inches(y_start),
                               Inches(w), Inches(h))
    s.fill.background()
    s.line.color.rgb = DANUBE; s.line.width = LW_NORMAL
    # label tab
    add_box(slide, x_from, y_start, 1.2, 0.35, label,
            fill=DANUBE, line=DANUBE, line_w=LW_THIN,
            text_color=WHITE, bold=True, font_size=9,
            align=PP_ALIGN.CENTER)


# --- SEQ-1 -----------------------------------------------------------------

def build_seq_1(prs):
    slide = add_slide(prs)
    add_title(slide,
        "SEQ-1 · Probe SDK ↔ Collector — joint tuning",
        sub="Resolves EC1↔EC2 (−1) and EC2↔EC4 (−1). Contract: collector owns EC1+EC2+EC4 jointly, time-based flush fallback prevents batch-starvation.")

    participants = [
        ("SS9",  "SS9 Probe SDK\n(customer-side)", "external"),
        ("SS10", "SS10 Metrics\nCollector",         "internal"),
        ("SS11", "SS11 Metric\nStore",              "internal"),
    ]
    centers = draw_lifelines(slide, participants, y_top=2.3, y_bottom=15.5)

    # Note over SS9,SS10
    add_sequence_note(slide, centers, ["SS9","SS10"], 3.5,
                      "IF-18 initial config push — collector owns EC1 + EC2 + EC4 jointly")
    add_sequence_arrow(slide, centers, "SS10", "SS9", 4.5,
        "cfg { probe_freq=6/min (EC1), agg_window=60 min (EC2), batch_size=50 (EC4), flush_fallback=10 }",
        heavy=True)

    # loop frame
    add_alt_frame(slide, centers, ["SS9","SS10"], 5.3, 11.3,
                  "loop every probe tick (≤6/min · EC1)")
    # enqueue (self)
    add_sequence_arrow(slide, centers, "SS9", "SS9", 6.0,
                       "enqueue event (≤4 KB · EC3)", kind="self")

    # alt frame (batch full OR time flush)
    add_alt_frame(slide, centers, ["SS9","SS10"], 7.3, 11.0,
                  "alt  batch=50 (EC4) / else time flush")
    add_sequence_arrow(slide, centers, "SS9", "SS10", 8.2,
                       "IF-17 probe payload (batch=50)", heavy=True)
    add_sequence_note(slide, centers, ["SS9","SS10"], 8.9,
                      "else: 60-min window elapses without batch fill (roof EC2↔EC4)")
    add_sequence_arrow(slide, centers, "SS9", "SS10", 10.0,
                       "IF-17 TIME-BASED FLUSH (batch ≥ flush_fallback)",
                       heavy=True)

    # aggregate + write (outside loop)
    add_sequence_arrow(slide, centers, "SS10", "SS10", 12.0,
                       "aggregate (60-min window · EC2)", kind="self")
    add_sequence_arrow(slide, centers, "SS10", "SS11", 12.9,
                       "IF-19 aggregated write")

    # alt — overhead violation
    add_alt_frame(slide, centers, ["SS9","SS10"], 13.5, 15.2,
                  "alt  overhead violation (EC1/EC3 breach)")
    add_sequence_arrow(slide, centers, "SS10", "SS9", 14.1,
                       "IF-18 backpressure · REJECT · cfg{probe_freq↓}",
                       heavy=True)
    add_sequence_arrow(slide, centers, "SS10", "SS10", 14.8,
                       "IF-26 audit probe-rejection (I.EXT.9)", kind="self")


# --- SEQ-2 -----------------------------------------------------------------

def build_seq_2(prs):
    slide = add_slide(prs)
    add_title(slide,
        "SEQ-2 · Traceback cache invalidation",
        sub="Resolves EC9↔EC11 (−2, STRONGEST roof negative). Intra-SS7 invalidation: fetcher-driven expiry of TTL cache on vendor-doc refresh.")

    participants = [
        ("SS3",  "SS3 Agent\nOrchestrator",       "internal"),
        ("SS6",  "SS6 Spec\nGenerator",           "internal"),
        ("CA",   "SS7.cache\n(TTL 24h · EC9)",    "internal"),
        ("FE",   "SS7.fetcher\n(refresh 7d · EC11)", "internal"),
        ("DOCS", "Documentation\nSources (I.EXT.5)", "external"),
    ]
    centers = draw_lifelines(slide, participants, y_top=2.3, y_bottom=15.5)

    add_sequence_note(slide, centers, ["CA","FE"], 3.5,
                      "cache and fetcher live inside SS7 — invalidation is intra-subsystem")

    add_sequence_arrow(slide, centers, "SS6", "CA", 4.6,
                       "IF-13 citation_check(claim)", heavy=True)

    # alt HIT / MISS
    add_alt_frame(slide, centers, ["SS6","DOCS"], 5.2, 10.2,
                  "alt  cache HIT (age < 24h · EC9) / else MISS")
    add_sequence_arrow(slide, centers, "CA", "SS6", 6.0,
                       "IF-14 traced citation", kind="return")
    # else branch
    add_sequence_arrow(slide, centers, "CA",   "FE",   6.9,
                       "fetch(source_uri)")
    add_sequence_arrow(slide, centers, "FE",   "DOCS", 7.6,
                       "GET vendor doc")
    add_sequence_arrow(slide, centers, "DOCS", "FE",   8.3,
                       "doc content", kind="return")
    add_sequence_arrow(slide, centers, "FE",   "CA",   9.0,
                       "write(citation, fetched_at=now)")
    add_sequence_arrow(slide, centers, "CA",   "SS6",  9.7,
                       "IF-14 traced citation", kind="return")

    # periodic refresh (every 7d)
    add_sequence_note(slide, centers, ["FE"], 10.8,
                      "every 7d per source (EC11)")
    add_sequence_arrow(slide, centers, "FE", "DOCS", 11.6,
                       "refresh(source_uri)")
    add_sequence_arrow(slide, centers, "DOCS", "FE", 12.3,
                       "new doc content (hash changed?)", kind="return")

    # alt hash changed
    add_alt_frame(slide, centers, ["CA","DOCS"], 12.9, 15.3,
                  "alt  new hash / else same hash")
    add_sequence_arrow(slide, centers, "FE", "CA", 13.7,
                       "INVALIDATE all entries keyed to source_uri",
                       heavy=True)
    add_sequence_note(slide, centers, ["CA","FE"], 14.2,
                      "resolves EC9↔EC11 roof (−2): stale TTL not allowed to outlive refresh")
    add_sequence_arrow(slide, centers, "CA", "SS3", 14.9,
                       "invalidation event (for in-flight specs)", kind="return")


# --- SEQ-3 -----------------------------------------------------------------

def build_seq_3(prs):
    slide = add_slide(prs)
    add_title(slide,
        "SEQ-3 · Session lock → parallel dispatch",
        sub="Resolves EC5↔EC6 (−1). Synchronous IF-09 routing policy write MUST ack before IF-03 dispatch — prevents race between per-session routing and parallel agent fan-out.")

    participants = [
        ("SS1",  "SS1 Founder\nWeb UI",            "internal"),
        ("SS2",  "SS2 Session\nManager",           "internal"),
        ("SS13", "SS13 Credential\nVault",         "internal"),
        ("SS4",  "SS4 LLM\nProvider Layer",        "internal"),
        ("SS3",  "SS3 Agent\nOrchestrator",        "internal"),
    ]
    centers = draw_lifelines(slide, participants, y_top=2.3, y_bottom=15.5)

    add_sequence_arrow(slide, centers, "SS1", "SS2", 3.5,
                       "IF-01 login + mode select (cloud | local)")
    add_sequence_arrow(slide, centers, "SS2", "SS13", 4.2,
                       "auth principal")
    add_sequence_arrow(slide, centers, "SS13", "SS2", 4.9,
                       "IF-27 signed session JWT (EC17)", kind="return")

    add_sequence_note(slide, centers, ["SS2"], 5.6,
                      "lock_mode(cloud | local) — pinned per session (EC5)")

    add_sequence_arrow(slide, centers, "SS2", "SS4", 6.5,
                       "IF-09 routing policy write (scale=2 per-session · EC5)",
                       heavy=True)
    add_sequence_arrow(slide, centers, "SS13", "SS4", 7.2,
                       "IF-29 provider credentials")
    add_sequence_arrow(slide, centers, "SS4", "SS2", 7.9,
                       "ack(routing_set=true)", kind="return")

    add_sequence_note(slide, centers, ["SS2","SS3"], 8.6,
                      "EC5↔EC6 resolution — dispatch cannot occur before routing is confirmed")

    add_sequence_arrow(slide, centers, "SS2", "SS3", 9.7,
                       "IF-03 dispatch(principal, locked_mode, intake_payload)",
                       heavy=True)
    add_activation(slide, centers["SS3"], 9.8, 13.4)

    add_alt_frame(slide, centers, ["SS3","SS4"], 10.1, 12.8,
                  "par  up to 5 parallel agents (EC6)")
    add_sequence_arrow(slide, centers, "SS3", "SS4", 10.9,
                       "IF-07 LLM call · agent A (≤8000 tok · EC8)")
    add_sequence_arrow(slide, centers, "SS3", "SS4", 11.6,
                       "IF-07 LLM call · agent B (≤8000 tok · EC8)")
    add_sequence_arrow(slide, centers, "SS3", "SS4", 12.3,
                       "IF-07 LLM call · agent C (≤8000 tok · EC8)")

    add_sequence_arrow(slide, centers, "SS3", "SS2", 13.1,
                       "IF-04 agent run state", kind="return")

    add_sequence_note(slide, centers, ["SS2","SS3"], 14.0,
                      "if dispatch were issued before IF-09 ack, agents could fan out to the wrong provider — the invariant BLOCKS that")


# --- SEQ-4 -----------------------------------------------------------------

def build_seq_4(prs):
    slide = add_slide(prs)
    add_title(slide,
        "SEQ-4 · Citation-gated CLI emission",
        sub="EC10 HARD FLOOR — 100% citation coverage. Emission MUST NOT fire with any unresolved claim. Contract: SS6→SS7 check → SS3 review → SS8 emit/block.")

    participants = [
        ("SS1",  "SS1 Founder\nWeb UI",       "internal"),
        ("SS3",  "SS3 Agent\nOrchestrator",   "internal"),
        ("SS6",  "SS6 Spec\nGenerator",       "internal"),
        ("SS7",  "SS7 Traceback\nStore",      "internal"),
        ("SS8",  "SS8 CLI\nEmitter",          "internal"),
        ("SS14", "SS14 Audit\nLog",           "internal"),
        ("VCS",  "Version Control\n(I.EXT.3)", "external"),
    ]
    centers = draw_lifelines(slide, participants, y_top=2.3, y_bottom=15.5)

    add_sequence_arrow(slide, centers, "SS3", "SS6", 3.5,
                       "IF-11 generate spec(requirements)")
    add_sequence_arrow(slide, centers, "SS6", "SS6", 4.2,
                       "draft spec (Mermaid v1 · EC12)", kind="self")

    add_sequence_note(slide, centers, ["SS6","SS7"], 5.0,
                      "EC10 HARD FLOOR — every claim must resolve")
    add_sequence_arrow(slide, centers, "SS6", "SS7", 5.9,
                       "IF-13 citation_check(claims[])", heavy=True)
    add_sequence_arrow(slide, centers, "SS7", "SS6", 6.6,
                       "IF-14 citations[] + unresolved[]", kind="return")

    # alt unresolved / all resolved
    add_alt_frame(slide, centers, ["SS1","VCS"], 7.2, 15.3,
                  "alt  unresolved > 0 (BLOCK) / else 100% resolved (EMIT)")

    # BLOCK branch
    add_sequence_arrow(slide, centers, "SS6", "SS3", 8.1,
                       "IF-12 draft · pending_review=true · citation_miss=N",
                       kind="return")
    add_sequence_arrow(slide, centers, "SS3", "SS14", 8.8,
                       "IF-24 audit(citation_miss)")
    add_sequence_arrow(slide, centers, "SS3", "SS1", 9.5,
                       "IF-02 'spec blocked: N unresolved citations'",
                       kind="return")
    add_sequence_note(slide, centers, ["SS3"], 10.0,
                      "EMISSION BLOCKED — no IF-15 fires")

    # EMIT branch
    add_sequence_arrow(slide, centers, "SS6", "SS3", 10.9,
                       "IF-12 draft · pending_review=true · citations_ok",
                       kind="return")
    add_sequence_arrow(slide, centers, "SS3", "SS1", 11.5,
                       "IF-02 render spec for approval")
    add_sequence_arrow(slide, centers, "SS1", "SS6", 12.1,
                       "IF-31 approve")
    add_sequence_arrow(slide, centers, "SS6", "SS1", 12.7,
                       "IF-32 rendered Mermaid (review)", kind="return")

    add_sequence_arrow(slide, centers, "SS3", "SS8", 13.3,
                       "IF-15 emit CLI bundle(spec_uri) — citation gate PASSED",
                       heavy=True)
    add_sequence_arrow(slide, centers, "SS8", "SS8", 13.9,
                       "package + sign (≤5 MB · EC14)", kind="self")
    add_sequence_arrow(slide, centers, "SS8", "VCS", 14.5,
                       "IF-EXT-04 signed commit", heavy=True)
    add_sequence_arrow(slide, centers, "SS8", "SS14", 15.1,
                       "IF-25 audit(bundle_id, hash, signer, vcs_commit)")


# ═══════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════

if __name__ == "__main__":
    outputs = []
    for builder in (build_n2_chart_pptx, build_dfd_pptx, build_sequence_pptx):
        path, slide_count = builder()
        size = os.path.getsize(path)
        outputs.append((path, slide_count, size))
        print(f"Saved: {path}  ({slide_count} slides, {size} bytes)")
    print("\nAll decks generated.")
